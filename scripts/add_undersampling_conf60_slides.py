"""
Add slides for undersampling and conf60 experiments:
Insert before "Rencana Eksperimen Lanjutan" (slide 107)

1. Undersampling motivation + results
2. Undersampling per-class F1 analysis
3. Conf60: motivation + confidence analysis
4. Conf60: complete results (2x2 class config)
5. Conf60: temuan breakthrough
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_GREEN = RGBColor(0x0F, 0x9D, 0x58)
COLOR_LIGHT_GREEN = RGBColor(0xE6, 0xF4, 0xEA)
COLOR_RED = RGBColor(0xD9, 0x3C, 0x3C)
COLOR_ORANGE = RGBColor(0xE6, 0x51, 0x00)
COLOR_TEXT = RGBColor(0x20, 0x20, 0x20)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_BEST = RGBColor(0xC6, 0xEF, 0xCE)


def set_fill(shape, hex_val):
    spPr = shape._element.spPr
    for nf in spPr.findall(qn('a:noFill')): spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr')); c.set('val', hex_val)


def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def para(tf, text, sz=10, bold=False, italic=False, color=None, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = color


def cell(c, text, bold=False, bg=None, fg=None, sz=9, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]; p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    if fg: r.font.color.rgb = fg
    if bg: c.fill.solid(); c.fill.fore_color.rgb = bg


def table(slide, headers, rows, l, t, w, h, cw=None, hsz=9, rsz=8.5, best_cells=None):
    shape = slide.shapes.add_table(len(rows)+1, len(headers), Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None: tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw): tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WHITE, sz=hsz, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LIGHT_BLUE if r % 2 == 0 else COLOR_WHITE
        for c, val in enumerate(row):
            is_best = best_cells and (r, c) in best_cells
            bg_c = COLOR_BEST if is_best else bg
            cell(tbl.cell(r+1, c), str(val), bold=(c==0 or is_best),
                 bg=bg_c, fg=COLOR_TEXT, sz=rsz,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


def insert(prs, index):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    xml_slides = prs.slides._sldIdLst
    el = xml_slides[-1]; xml_slides.remove(el); xml_slides.insert(index, el)
    return slide


# ── LOAD ──
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Insert 5 slides at index 106 (before "Rencana Eksperimen Lanjutan" slide 107)
slides = [insert(prs, 106 + i) for i in range(5)]

# ═══ SLIDE 1: Undersampling Motivation + Result ═══
s = slides[0]
t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Tahap 5A: Undersampling Neutral -- Mengatasi Imbalance', sz=18, bold=True, color=COLOR_TEXT, first=True)

# Motivation
m = tb(s, 0.3, 0.75, 9.4, 0.65)
m.text_frame.word_wrap = True; set_fill(m, 'FFF0CC')
para(m.text_frame, 'Masalah: Imbalance ratio 36.8:1 (neutral vs negative)', sz=10, bold=True, color=COLOR_TEXT, first=True)
para(m.text_frame, 'Model cenderung prediksi semua neutral -> accuracy 95% tapi F1 rendah', sz=9.5, color=COLOR_GRAY)

# Strategy table
table(s, ['Variasi', 'Neutral Train', 'Total Train', 'Rasio N:Neg'], [
    ['Original', '4,192', '5,348', '36.8:1'],
    ['Under-660', '660', '1,816', '5.8:1'],
    ['Under-382', '382', '1,538', '3.4:1'],
    ['Under-114', '114', '1,270', '1:1'],
], 0.3, 1.55, 9.4, 1.55, cw=[2.2, 2.3, 2.3, 2.6], hsz=9.5, rsz=9.5)

# Results summary
table(s, ['Model', 'Original F1', 'Under-660 F1', 'Improvement'], [
    ['Intermediate TL', '0.363', '0.257', '-29%'],
    ['FCNN', '0.263', '0.348', '+32%'],
    ['Late Fusion', '0.296', '0.405', '+37%'],
], 0.3, 3.25, 9.4, 1.2, cw=[2.5, 2.3, 2.3, 2.3], hsz=9.5, rsz=9.5,
   best_cells=[(2, 2), (2, 3)])
print('  Slide 1: Undersampling Result')

# ═══ SLIDE 2: Undersampling Per-Class Analysis ═══
s = slides[1]
t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Tahap 5A: Analisis Per-Kelas Undersampling', sz=18, bold=True, color=COLOR_TEXT, first=True)

# Per-class table (Late Fusion - best)
t_lbl = tb(s, 0.3, 0.72, 9.4, 0.3)
para(t_lbl.text_frame, 'Late Fusion -- Per-Class F1 Score', sz=10, bold=True, color=COLOR_BLUE, first=True)

table(s, ['Dataset', 'Macro F1', 'neutral', 'happy', 'sad', 'negative'], [
    ['Original', '0.296', '0.967', '0.158', '0.061', '0.000'],
    ['Under-660', '0.405', '0.963', '0.075', '0.581', '0.000'],
    ['Under-382', '0.294', '0.931', '0.024', '0.222', '0.000'],
    ['Under-114', '0.160', '0.516', '0.010', '0.115', '0.000'],
], 0.3, 1.05, 9.4, 1.5, cw=[2.0, 1.8, 1.5, 1.5, 1.4, 1.7], hsz=9.5, rsz=9,
   best_cells=[(1, 4)])  # sad F1 under-660

# Temuan boxes
b1 = tb(s, 0.3, 2.7, 4.5, 1.5)
b1.text_frame.word_wrap = True; set_fill(b1, 'E6F4EA')
para(b1.text_frame, 'Temuan Positif: Sad F1 naik drastis', sz=10, bold=True, color=COLOR_GREEN, first=True)
para(b1.text_frame, '', sz=4)
para(b1.text_frame, 'Late Fusion sad: 0.061 -> 0.581 (+852%)', sz=9.5, color=COLOR_TEXT)
para(b1.text_frame, 'FCNN sad: 0.000 -> 0.341', sz=9.5, color=COLOR_TEXT)
para(b1.text_frame, '', sz=3)
para(b1.text_frame, 'Under-660 = sweet spot:', sz=9.5, bold=True, color=COLOR_TEXT)
para(b1.text_frame, '- Cukup seimbang', sz=9, color=COLOR_TEXT)
para(b1.text_frame, '- Data training masih cukup', sz=9, color=COLOR_TEXT)

b2 = tb(s, 4.95, 2.7, 4.75, 1.5)
b2.text_frame.word_wrap = True; set_fill(b2, 'FCE4E4')
para(b2.text_frame, 'Masalah: Negative tetap F1 ~ 0', sz=10, bold=True, color=COLOR_RED, first=True)
para(b2.text_frame, '', sz=4)
para(b2.text_frame, 'Bahkan dengan rasio 1:1 (under-114),', sz=9.5, color=COLOR_TEXT)
para(b2.text_frame, 'negative masih tidak terdeteksi (F1 < 0.1)', sz=9.5, color=COLOR_TEXT)
para(b2.text_frame, '', sz=3)
para(b2.text_frame, 'Penyebab: Test set terlalu kecil', sz=9.5, bold=True, color=COLOR_TEXT)
para(b2.text_frame, '- Happy: 10 | Sad: 29 | Negative: 16', sz=9, color=COLOR_TEXT)
para(b2.text_frame, '- < 30 sampel -> evaluasi tidak reliable', sz=9, color=COLOR_TEXT)
print('  Slide 2: Per-Class Analysis')

# ═══ SLIDE 3: Conf60 Motivation ═══
s = slides[2]
t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Tahap 5B: Confidence Filtering >= 60% -- BREAKTHROUGH', sz=18, bold=True, color=COLOR_TEXT, first=True)

# Analysis
a = tb(s, 0.3, 0.75, 9.4, 0.65)
a.text_frame.word_wrap = True; set_fill(a, 'FFF0CC')
para(a.text_frame, 'Analisis: Confidence Face API per kelas', sz=10, bold=True, color=COLOR_TEXT, first=True)
para(a.text_frame, 'Kelas minoritas punya confidence rendah -> kemungkinan label SALAH (noise)', sz=9.5, color=COLOR_GRAY)

# Confidence table
table(s, ['Emosi', 'Confidence Rata-rata', 'Status'], [
    ['Neutral', '0.959', 'Tinggi (reliable)'],
    ['Happy', '0.878', 'Tinggi'],
    ['Sad', '0.770', 'Moderat'],
    ['Angry', '0.634', 'RENDAH'],
    ['Fearful', '0.671', 'RENDAH'],
    ['Disgusted', '0.566', 'SANGAT RENDAH'],
    ['Surprised', '0.730', 'Moderat'],
], 0.3, 1.55, 5.8, 2.5, cw=[1.8, 2.3, 2.3], hsz=9.5, rsz=9,
   best_cells=[(3, 2), (4, 2), (5, 2)])

# Strategy
st = tb(s, 6.3, 1.55, 3.4, 2.5)
st.text_frame.word_wrap = True; set_fill(st, 'E8F0FE')
para(st.text_frame, 'Strategi Filtering', sz=10, bold=True, color=COLOR_BLUE, first=True)
para(st.text_frame, '', sz=4)
para(st.text_frame, 'Filter sampel confidence < 60%', sz=9.5, color=COLOR_TEXT)
para(st.text_frame, '', sz=3)
para(st.text_frame, 'Dampak:', sz=9.5, bold=True, color=COLOR_TEXT)
para(st.text_frame, '- Total: 7,091 -> 6,795', sz=9, color=COLOR_TEXT)
para(st.text_frame, '- Hanya 4.2% data dihilangkan', sz=9, color=COLOR_TEXT)
para(st.text_frame, '- Kelas minoritas -50-57%', sz=9, color=COLOR_TEXT)
para(st.text_frame, '', sz=3)
para(st.text_frame, 'Prinsip: Data bersih lebih', sz=9, italic=True, color=COLOR_GRAY)
para(st.text_frame, 'penting dari data banyak', sz=9, italic=True, color=COLOR_GRAY)

# Comparison best
c = tb(s, 0.3, 4.2, 9.4, 0.4)
c.text_frame.word_wrap = True; set_fill(c, '0F9D58')
para(c.text_frame, 'Hasil: Best F1 naik dari 0.412 -> 0.521 (+26%) | Intermediate TL 4-class B3',
     sz=10, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER, first=True)
print('  Slide 3: Conf60 Motivation')

# ═══ SLIDE 4: Conf60 Complete Results ═══
s = slides[3]
t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Tahap 5B: Hasil Lengkap Conf60 (Best per Config)', sz=18, bold=True, color=COLOR_TEXT, first=True)

# 7-class comparison
t7 = tb(s, 0.3, 0.75, 4.5, 0.3)
para(t7.text_frame, '7-Class (Best Macro F1)', sz=10, bold=True, color=COLOR_BLUE, first=True)

table(s, ['Model', 'Original', 'Conf60', 'Diff'], [
    ['CNN', '0.137', '0.277', '+102%'],
    ['FCNN', '0.158', '0.244', '+54%'],
    ['Intermediate', '0.137', '0.261', '+91%'],
    ['Late Fusion', '0.175', '0.289', '+65%'],
    ['CNN TL', '0.154', '0.273', '+77%'],
    ['Intermediate TL', '0.180', '0.292', '+62%'],
    ['Late Fusion TL', '0.167', '0.301', '+80%'],
], 0.3, 1.05, 4.5, 2.4, cw=[1.8, 1.2, 1.2, 1.3], hsz=9, rsz=8.5,
   best_cells=[(6, 2)])

# 4-class comparison
t4 = tb(s, 5.1, 0.75, 4.6, 0.3)
para(t4.text_frame, '4-Class (Best Macro F1)', sz=10, bold=True, color=COLOR_BLUE, first=True)

table(s, ['Model', 'Original', 'Conf60', 'Diff'], [
    ['CNN', '0.265', '0.448', '+69%'],
    ['FCNN', '0.361', '0.460', '+27%'],
    ['Intermediate', '0.269', '0.445', '+65%'],
    ['Late Fusion', '0.394', '0.482', '+22%'],
    ['CNN TL', '0.274', '0.507', '+85%'],
    ['Intermediate TL', '0.412', '0.521', '+26%'],
    ['Late Fusion TL', '0.372', '0.513', '+38%'],
], 5.1, 1.05, 4.6, 2.4, cw=[1.8, 1.2, 1.2, 1.3], hsz=9, rsz=8.5,
   best_cells=[(5, 2), (5, 3)])

# Bottom insight
b = tb(s, 0.3, 3.65, 9.4, 0.7)
b.text_frame.word_wrap = True; set_fill(b, 'E6F4EA')
para(b.text_frame, 'Semua model mengalami peningkatan signifikan dengan conf60', sz=10, bold=True, color=COLOR_GREEN, first=True)
para(b.text_frame, 'Best overall: Intermediate TL 4-class B3 = 0.521 (vs 0.412 sebelumnya)', sz=9.5, color=COLOR_TEXT)
para(b.text_frame, 'Peningkatan rata-rata: 7-class +76%, 4-class +47%', sz=9.5, color=COLOR_TEXT)
print('  Slide 4: Conf60 Complete Results')

# ═══ SLIDE 5: Conf60 Temuan ═══
s = slides[4]
t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Analisis Hasil -- Temuan Confidence Filtering', sz=18, bold=True, color=COLOR_TEXT, first=True)

# Temuan 24
t24 = tb(s, 0.3, 0.78, 9.3, 0.25)
para(t24.text_frame, 'Temuan 24: Breakthrough -- F1 0.412 -> 0.521 (+26%)', sz=9, bold=True, color=COLOR_BLUE, first=True)

table(s, ['Tahap', 'Best Model', 'Macro F1'], [
    ['Tahap 3 (TL)', 'Intermediate TL B1', '0.412'],
    ['Tahap 4 (Front-only)', 'Intermediate TL B1', '0.412'],
    ['Tahap 5 (Conf60)', 'Intermediate TL B3', '0.521'],
], 0.3, 1.08, 9.3, 1.05, cw=[2.5, 3.5, 3.3], hsz=9, rsz=8.5,
   best_cells=[(2, 2)])

# Temuan 25
t25 = tb(s, 0.3, 2.25, 4.5, 0.25)
para(t25.text_frame, 'Temuan 25: Hanya 4.2% data dihilangkan', sz=9, bold=True, color=COLOR_BLUE, first=True)

m25 = tb(s, 0.3, 2.55, 4.5, 1.2)
m25.text_frame.word_wrap = True; set_fill(m25, 'E6F4EA')
para(m25.text_frame, 'Trade-off: Sedikit data tapi bersih', sz=9.5, bold=True, color=COLOR_GREEN, first=True)
para(m25.text_frame, '', sz=3)
para(m25.text_frame, '- 7,091 -> 6,795 sampel (4.2% filtered)', sz=9, color=COLOR_TEXT)
para(m25.text_frame, '- Kelas minoritas: 50-57% filtered', sz=9, color=COLOR_TEXT)
para(m25.text_frame, '- Tapi F1 naik signifikan', sz=9, color=COLOR_TEXT)
para(m25.text_frame, '', sz=3)
para(m25.text_frame, 'Bersih menang vs banyak', sz=9, italic=True, color=COLOR_GRAY)

# Temuan 26
t26 = tb(s, 5.0, 2.25, 4.7, 0.25)
para(t26.text_frame, 'Temuan 26: Label noise adalah masalah utama', sz=9, bold=True, color=COLOR_BLUE, first=True)

m26 = tb(s, 5.0, 2.55, 4.7, 1.2)
m26.text_frame.word_wrap = True; set_fill(m26, 'FCE4E4')
para(m26.text_frame, 'Penyebab performa rendah sebelumnya:', sz=9.5, bold=True, color=COLOR_RED, first=True)
para(m26.text_frame, '', sz=3)
para(m26.text_frame, '- BUKAN karakteristik data natural', sz=9, color=COLOR_TEXT)
para(m26.text_frame, '- BUKAN arsitektur yang kurang', sz=9, color=COLOR_TEXT)
para(m26.text_frame, '- TAPI label noise Face API', sz=9, bold=True, color=COLOR_TEXT)
para(m26.text_frame, '', sz=3)
para(m26.text_frame, 'Memperkuat alasan validasi ahli', sz=9, italic=True, color=COLOR_GRAY)

# Conclusion
c = tb(s, 0.3, 3.85, 9.4, 0.4)
c.text_frame.word_wrap = True; set_fill(c, '1A73E8')
para(c.text_frame, 'Kesimpulan: Kualitas label > kuantitas data. Confidence filtering terbukti efektif meningkatkan performa model.',
     sz=9.5, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER, first=True)
print('  Slide 5: Conf60 Temuan')

# ── SAVE ──
prs.save(PPTX_PATH)
print(f'\nSaved! Total: {len(prs.slides)} slides')

for i in range(104, len(prs.slides)):
    slide = prs.slides[i]
    title = next((s.text_frame.text.strip()[:58].replace(chr(10),' ').encode('ascii','replace').decode()
                  for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), '')
    print(f'  Slide {i+1}: {title}')
