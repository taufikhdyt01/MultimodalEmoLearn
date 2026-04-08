"""
Rebuild slide 67 (index 66) - Pendekatan Validasi SSL
Clear all shapes, rebuild with proven approach (same as slides 74/79).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE        = RGBColor(0x1A, 0x73, 0xE8)
COLOR_ORANGE      = RGBColor(0xFF, 0x6D, 0x00)
COLOR_GREEN       = RGBColor(0x0F, 0x9D, 0x58)
COLOR_LIGHT_BLUE  = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_LIGHT_OG    = RGBColor(0xFF, 0xF0, 0xCC)
COLOR_LIGHT_GREEN = RGBColor(0xE6, 0xF4, 0xEA)
COLOR_TEXT        = RGBColor(0x20, 0x20, 0x20)
COLOR_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY        = RGBColor(0x66, 0x66, 0x66)
COLOR_HEADER_BG   = RGBColor(0x1A, 0x73, 0xE8)


def set_fill(tb, hex_val):
    """Add solidFill to textbox spPr — must remove noFill first!"""
    spPr = tb._element.spPr
    # Remove noFill that python-pptx adds by default (this blocks solidFill)
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', hex_val)


def add_tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))


def add_para(tf, text, size=10, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, first=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return para


def style_cell(cell, text, bold=False, bg=None, fg=None,
               size=9, align=PP_ALIGN.LEFT):
    cell.text = text
    para = cell.text_frame.paragraphs[0]
    para.alignment = align
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if fg:
        run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def add_table(slide, headers, rows, left, top, width, height,
              col_widths=None, h_size=9, r_size=8.5):
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None:
            tblPr.remove(sid)
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(Inches(width) * cw / total)
    for c, h in enumerate(headers):
        style_cell(tbl.cell(0, c), h, bold=True,
                   bg=COLOR_HEADER_BG, fg=COLOR_WHITE,
                   size=h_size, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LIGHT_BLUE if r % 2 == 0 else COLOR_WHITE
        for c, val in enumerate(row):
            style_cell(tbl.cell(r+1, c), str(val),
                       bold=(c == 0), bg=bg, fg=COLOR_TEXT,
                       size=r_size,
                       align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


# ── LOAD ──────────────────────────────────────────────
prs = Presentation(PPTX_PATH)
slide = prs.slides[66]  # slide 67

# Clear all existing shapes
sp_tree = slide.shapes._spTree
for child in list(sp_tree):
    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
    if tag not in ('nvGrpSpPr', 'grpSpPr'):
        sp_tree.remove(child)
print(f'Cleared. Remaining: {len(list(slide.shapes))} shapes')

# ── TITLE ─────────────────────────────────────────────
tb = add_tb(slide, 0.3, 0.12, 9.4, 0.6)
tf = tb.text_frame
add_para(tf, 'Pendekatan Validasi: Expert (146 Sampel) + Self-Supervised Learning',
         size=20, bold=True, color=COLOR_TEXT, align=PP_ALIGN.LEFT, first=True)

# ── FLOW BOXES (3 kotak + panah sebagai teks) ─────────
# Box 1 - Total Dataset (biru)
tb1 = add_tb(slide, 0.25, 0.85, 2.1, 0.85)
tf1 = tb1.text_frame; tf1.word_wrap = True
add_para(tf1, 'Total Dataset', size=10, bold=True, color=COLOR_WHITE,
         align=PP_ALIGN.CENTER, first=True)
add_para(tf1, '9,894 sampel', size=9.5, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
set_fill(tb1, '1A73E8')

# Panah 1
tb_a1 = add_tb(slide, 2.4, 1.0, 0.4, 0.5)
tf_a1 = tb_a1.text_frame
add_para(tf_a1, '>', size=16, bold=True, color=COLOR_TEXT,
         align=PP_ALIGN.CENTER, first=True)

# Box 2 - Validasi Ahli (oranye)
tb2 = add_tb(slide, 2.85, 0.85, 2.6, 0.85)
tf2 = tb2.text_frame; tf2.word_wrap = True
add_para(tf2, 'Validasi Ahli', size=10, bold=True, color=COLOR_WHITE,
         align=PP_ALIGN.CENTER, first=True)
add_para(tf2, '146 sampel (1.5%)', size=9.5, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_para(tf2, '3 validator psikologi', size=9.5, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
set_fill(tb2, 'E65100')

# Panah 2
tb_a2 = add_tb(slide, 5.5, 1.0, 0.4, 0.5)
tf_a2 = tb_a2.text_frame
add_para(tf_a2, '>', size=16, bold=True, color=COLOR_TEXT,
         align=PP_ALIGN.CENTER, first=True)

# Box 3 - SSL (hijau)
tb3 = add_tb(slide, 5.95, 0.85, 2.6, 0.85)
tf3 = tb3.text_frame; tf3.word_wrap = True
add_para(tf3, 'SSL Verification', size=10, bold=True, color=COLOR_WHITE,
         align=PP_ALIGN.CENTER, first=True)
add_para(tf3, '~9,748 sampel (98.5%)', size=9.5, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_para(tf3, 'ResNet18 embedding', size=9.5, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
set_fill(tb3, '0F9D58')

# Panah 3
tb_a3 = add_tb(slide, 8.6, 1.0, 0.35, 0.5)
tf_a3 = tb_a3.text_frame
add_para(tf_a3, '>', size=16, bold=True, color=COLOR_TEXT,
         align=PP_ALIGN.CENTER, first=True)

# Box 4 - Hasil (biru muda)
tb4 = add_tb(slide, 8.98, 0.88, 0.9, 0.75)
tf4 = tb4.text_frame; tf4.word_wrap = True
add_para(tf4, 'Dataset Valid', size=8.5, bold=True, color=COLOR_TEXT,
         align=PP_ALIGN.CENTER, first=True)
set_fill(tb4, 'E8F0FE')

# Label referensi bawah flow
tb_ref0 = add_tb(slide, 0.25, 1.78, 9.6, 0.28)
tf_ref0 = tb_ref0.text_frame
add_para(tf_ref0,
         'Mengacu MER2024 (Lian et al., 2024, ACM MM): standar 1% labeled + 99% SSL untuk emotion recognition',
         size=8, italic=True, color=COLOR_GRAY, align=PP_ALIGN.CENTER, first=True)

# ── TABEL PERBANDINGAN ────────────────────────────────
tb_lbl = add_tb(slide, 0.25, 2.15, 5.7, 0.32)
tf_lbl = tb_lbl.text_frame
add_para(tf_lbl, 'Perbandingan dengan MER2024 (Benchmark Resmi)',
         size=10, bold=True, color=COLOR_TEXT, first=True)

add_table(slide,
    headers=['Aspek', 'MER2024', 'Penelitian Ini'],
    rows=[
        ['Total sampel', '115,595', '9,894'],
        ['Labeled (divalidasi ahli)', '1,169 (1%)', '146 (1.5%)'],
        ['Unlabeled (SSL)', '114,426 (99%)', '~9,748 (98.5%)'],
        ['Jumlah anotator', '5 orang', '3 orang'],
        ['Threshold agreement', '4/5 = 80%', '2/3 = 67%'],
        ['Metrik kesepakatan', 'Agreement rate', "Fleiss' Kappa >= 0.61"],
    ],
    left=0.25, top=2.52, width=5.7, height=2.6,
    col_widths=[2.4, 1.65, 1.65], h_size=9.5, r_size=9,
)

# ── KOTAK DETAIL VALIDATOR ────────────────────────────
tb_vlbl = add_tb(slide, 6.15, 2.15, 3.6, 0.32)
tf_vlbl = tb_vlbl.text_frame
add_para(tf_vlbl, 'Detail Validasi',
         size=10, bold=True, color=COLOR_TEXT, first=True)

tb_vbox = add_tb(slide, 6.15, 2.52, 3.6, 2.6)
tf_vbox = tb_vbox.text_frame; tf_vbox.word_wrap = True
lines = [
    ('3 validator psikologi', True),
    ('(mahasiswa S2/S3 psikologi)', False),
    ('', False),
    ('Masing-masing validasi 146 sampel', False),
    ('secara independen', False),
    ('', False),
    ("Hitung Fleiss' Kappa (3+ rater)", False),
    ('Target: k >= 0.61 (substantial)', False),
    ('', False),
    ('Sampling: stratified per kelas', False),
    ('Tool: web app Streamlit', False),
]
for i, (text, bold) in enumerate(lines):
    add_para(tf_vbox, text, size=9.5, bold=bold, color=COLOR_TEXT,
             first=(i == 0))
set_fill(tb_vbox, 'F1F3F4')

# ── REFERENSI ─────────────────────────────────────────
tb_ref = add_tb(slide, 0.25, 5.2, 9.6, 0.35)
tf_ref = tb_ref.text_frame; tf_ref.word_wrap = True
add_para(tf_ref,
         'Lian, Z., et al. (2024). MER 2024: Semi-Supervised Learning, Noise Robustness, and '
         'Open-Vocabulary Multimodal Emotion Recognition. ACM MM 2024. arXiv:2404.17113',
         size=7.5, italic=True, color=COLOR_GRAY, first=True)

# ── SAVE ──────────────────────────────────────────────
prs.save(PPTX_PATH)
print('Saved!')
print(f'Slide 67 shapes: {len(list(prs.slides[66].shapes))}')
