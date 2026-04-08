"""
Insert a new slide after Slide 66 (Validasi Ahli Psikologi):
  "Pendekatan Validasi: SSL + Expert (~100 Sampel)"

Shows:
  - Flow diagram: 9,894 total → 100 expert → SSL for rest
  - Comparison table: penelitian ini vs MER2024
  - Validator info: 3 orang, Fleiss' Kappa >= 0.61
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE      = RGBColor(0x1A, 0x73, 0xE8)
COLOR_ORANGE    = RGBColor(0xFF, 0x6D, 0x00)
COLOR_GREEN     = RGBColor(0x0F, 0x9D, 0x58)
COLOR_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_LIGHT_ORANGE = RGBColor(0xFF, 0xF0, 0xCC)
COLOR_LIGHT_GREEN  = RGBColor(0xE6, 0xF4, 0xEA)
COLOR_TEXT      = RGBColor(0x20, 0x20, 0x20)
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_HEADER_BG = RGBColor(0x1A, 0x73, 0xE8)


def set_box_fill(shape, hex_color):
    spPr = shape._element.spPr
    # Remove existing fills
    for old in spPr.findall(qn('a:solidFill')):
        spPr.remove(old)
    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', hex_color)


def add_textbox(slide, text, left, top, width, height,
                font_size=10, bold=False, color=None, bg_hex=None,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if bg_hex:
        set_box_fill(tb, bg_hex)
    return tb


def add_multiline_textbox(slide, lines, left, top, width, height,
                           font_size=10, bg_hex=None, color=None,
                           bold_first=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if bg_hex:
        set_box_fill(tb, bg_hex)
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = (bold_first and i == 0)
        if color:
            run.font.color.rgb = color
    return tb


def style_cell(cell, text, bold=False, bg=None, fg=None, size=9, align=PP_ALIGN.LEFT):
    cell.text = text
    para = cell.text_frame.paragraphs[0]
    para.alignment = align
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if fg:
        run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def add_table(slide, headers, rows, left, top, width, height,
              col_widths=None, h_size=9, r_size=8.5):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        styleId = tblPr.find(qn('a:tableStyleId'))
        if styleId is not None:
            tblPr.remove(styleId)
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / total)
    for c, h in enumerate(headers):
        style_cell(tbl.cell(0, c), h, bold=True,
                   bg=COLOR_HEADER_BG, fg=COLOR_WHITE,
                   size=h_size, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LIGHT_BLUE if r % 2 == 0 else COLOR_WHITE
        for c, val in enumerate(row):
            bold_cell = (c == 0)
            style_cell(tbl.cell(r + 1, c), str(val),
                       bold=bold_cell, bg=bg, fg=COLOR_TEXT,
                       size=r_size,
                       align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)


def insert_slide_after(prs, index):
    layout = prs.slide_layouts[10]
    slide = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst
    new_elem = xml_slides[-1]
    xml_slides.remove(new_elem)
    xml_slides.insert(index, new_elem)
    return slide


# ── LOAD ────────────────────────────────────────────────
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Insert after slide 66 (index 66)
slide = insert_slide_after(prs, 66)

# ── TITLE ───────────────────────────────────────────────
add_textbox(slide,
    'Pendekatan Validasi: Expert (~100 Sampel) + Self-Supervised Learning',
    Inches(0.33), Inches(0.12), Inches(9.3), Inches(0.58),
    font_size=18, bold=True, color=COLOR_TEXT)

# ── FLOW DIAGRAM (3 boxes + arrows) ─────────────────────
# Box 1: Total dataset
add_multiline_textbox(slide,
    ['Total Dataset', '9,894 sampel'],
    Inches(0.25), Inches(0.85), Inches(2.2), Inches(0.75),
    font_size=10, bg_hex='1A73E8', color=COLOR_WHITE,
    bold_first=True, align=PP_ALIGN.CENTER)

# Arrow 1
add_textbox(slide, '\u2192', Inches(2.5), Inches(0.95), Inches(0.4), Inches(0.5),
            font_size=16, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER)

# Box 2: Expert validation
add_multiline_textbox(slide,
    ['Validasi Ahli', '~100 sampel (1%)', 'stratified per kelas', '3 validator psikologi'],
    Inches(2.95), Inches(0.85), Inches(2.5), Inches(0.9),
    font_size=9.5, bg_hex='FF6D00', color=COLOR_WHITE,
    bold_first=True, align=PP_ALIGN.CENTER)

# Arrow 2
add_textbox(slide, '\u2192', Inches(5.5), Inches(0.95), Inches(0.4), Inches(0.5),
            font_size=16, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER)

# Box 3: SSL verification
add_multiline_textbox(slide,
    ['SSL Verification', '~9,794 sampel (99%)', 'embedding ResNet18', 'centroid consistency check'],
    Inches(5.95), Inches(0.85), Inches(2.5), Inches(0.9),
    font_size=9.5, bg_hex='0F9D58', color=COLOR_WHITE,
    bold_first=True, align=PP_ALIGN.CENTER)

# Arrow 3
add_textbox(slide, '\u2192', Inches(8.5), Inches(0.95), Inches(0.4), Inches(0.5),
            font_size=16, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER)

# Box 4: Result
add_multiline_textbox(slide,
    ['Dataset', 'Tervalidasi'],
    Inches(8.9), Inches(0.9), Inches(0.95), Inches(0.7),
    font_size=9, bg_hex='4285F4', color=COLOR_WHITE,
    bold_first=False, align=PP_ALIGN.CENTER)

# ── LABEL BELOW BOXES ───────────────────────────────────
add_textbox(slide,
    'Mengacu pada MER2024 (Lian et al., 2024, ACM MM): 1% labeled + 99% SSL',
    Inches(0.25), Inches(1.82), Inches(9.6), Inches(0.3),
    font_size=8.5, italic=True, color=RGBColor(0x55, 0x55, 0x55),
    align=PP_ALIGN.CENTER)

# ── COMPARISON TABLE ────────────────────────────────────
add_textbox(slide,
    'Perbandingan dengan MER2024 (Benchmark Resmi ACM MM + IJCAI 2024)',
    Inches(0.25), Inches(2.2), Inches(5.7), Inches(0.35),
    font_size=10, bold=True, color=COLOR_TEXT)

add_table(slide,
    headers=['Aspek', 'MER2024', 'Penelitian Ini'],
    rows=[
        ['Total sampel', '115,595', '9,894'],
        ['Labeled (divalidasi)', '1,169 (1%)', '~100 (1%)'],
        ['Unlabeled (SSL)', '99%', '99%'],
        ['Jumlah anotator', '5 orang', '3 orang'],
        ['Threshold agreement', '4/5 = 80%', '2/3 = 67%'],
        ['Metrik', 'Agreement rate', "Fleiss' Kappa \u2265 0.61"],
    ],
    left=Inches(0.25), top=Inches(2.6), width=Inches(5.7), height=Inches(2.55),
    col_widths=[2.2, 1.75, 1.75], h_size=9.5, r_size=9,
)

# ── VALIDATOR INFO BOX ──────────────────────────────────
add_textbox(slide,
    'Detail Validator',
    Inches(6.2), Inches(2.2), Inches(3.6), Inches(0.35),
    font_size=10, bold=True, color=COLOR_TEXT)

add_multiline_textbox(slide,
    [
        '3 validator berlatar psikologi',
        '(mahasiswa S2/S3 psikologi)',
        '',
        'Validasi 100 sampel secara independen',
        '-> hitung Fleiss\' Kappa',
        '-> target \u03ba \u2265 0.61 (substantial)',
        '',
        'Sampel dipilih: stratified random',
        'sampling per kelas emosi',
        '',
        'Tool: web app Streamlit',
        '(klik gambar \u2192 pilih label)',
    ],
    Inches(6.2), Inches(2.6), Inches(3.6), Inches(2.55),
    font_size=9.5, bg_hex='F8F9FA', color=COLOR_TEXT,
    bold_first=False)

# ── REFERENCE ───────────────────────────────────────────
add_textbox(slide,
    'Referensi: Lian, Z., et al. (2024). MER 2024: Semi-Supervised Learning, Noise Robustness, and Open-Vocabulary '
    'Multimodal Emotion Recognition. ACM MM 2024. arXiv:2404.17113',
    Inches(0.25), Inches(5.2), Inches(9.6), Inches(0.35),
    font_size=7.5, italic=True, color=RGBColor(0x66, 0x66, 0x66))

# ── SAVE ────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f'Saved! Total slides: {len(prs.slides)}')

prs2 = Presentation(PPTX_PATH)
for i in range(64, 70):
    slide = prs2.slides[i]
    title = next((s.text_frame.text.strip()[:60].replace(chr(10), ' ').encode('ascii', 'replace').decode()
                  for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), '')
    print(f'  Slide {i+1}: {title}')
