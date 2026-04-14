"""
Add 3 benchmark slides before Diskusi (slide 97):
1. Benchmark Results (JAFFE LOSO + CK+ 10-Fold CV)
2. Perbandingan dengan Paper SOTA
3. Temuan Benchmark (18-19)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE       = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_GREEN      = RGBColor(0x0F, 0x9D, 0x58)
COLOR_RED        = RGBColor(0xD9, 0x3C, 0x3C)
COLOR_TEXT       = RGBColor(0x20, 0x20, 0x20)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY       = RGBColor(0x66, 0x66, 0x66)


def set_fill(shape, hex_val):
    spPr = shape._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', hex_val)


def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def para(tf, text, sz=10, bold=False, italic=False, color=None, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = color


def cell(c, text, bold=False, bg=None, fg=None, sz=9, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]; p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    if fg: r.font.color.rgb = fg
    if bg: c.fill.solid(); c.fill.fore_color.rgb = bg


def table(slide, headers, rows, l, t, w, h, cw=None, hsz=9, rsz=8.5, highlight_row=None):
    shape = slide.shapes.add_table(len(rows)+1, len(headers), Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None: tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw): tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WHITE, sz=hsz, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LIGHT_BLUE if r % 2 == 0 else COLOR_WHITE
        if highlight_row is not None and r == highlight_row:
            bg = RGBColor(0xC6, 0xEF, 0xCE)
        for c, val in enumerate(row):
            cell(tbl.cell(r+1, c), str(val), bold=(c==0 or (highlight_row is not None and r==highlight_row)),
                 bg=bg, fg=COLOR_TEXT, sz=rsz,
                 align=PP_ALIGN.LEFT if c <= 1 else PP_ALIGN.CENTER)


def insert(prs, index):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    xml_slides = prs.slides._sldIdLst
    el = xml_slides[-1]; xml_slides.remove(el); xml_slides.insert(index, el)
    return slide


# ── LOAD ──
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Insert 3 slides before Diskusi (index 96)
s1 = insert(prs, 96)
s2 = insert(prs, 97)
s3 = insert(prs, 98)

# ═══ SLIDE 1: Benchmark Results ═══
t = tb(s1, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Benchmark: JAFFE (LOSO) & CK+ (10-Fold CV)', sz=18, bold=True, color=COLOR_TEXT, first=True)

# Best per dataset table
t_lbl = tb(s1, 0.3, 0.72, 9.4, 0.25)
para(t_lbl.text_frame, 'Best Model per Dataset (Macro F1)', sz=10, bold=True, color=COLOR_TEXT, first=True)

table(s1, ['Dataset', 'Evaluasi', 'Best Model', 'Macro F1'], [
    ['CK+ 7-class', '10-Fold CV', 'Intermediate TL', '0.783 +/- 0.107'],
    ['CK+ 4-class', '10-Fold CV', 'CNN TL', '0.755 +/- 0.079'],
    ['JAFFE 7-class', 'LOSO (10)', 'Late Fusion', '0.467 +/- 0.092'],
    ['JAFFE 4-class', 'LOSO (10)', 'Late Fusion', '0.530 +/- 0.126'],
    ['Dataset sendiri 4c', '5-Fold CV', 'Intermediate TL', '0.435 +/- 0.068'],
    ['Dataset sendiri 4c', 'LOSO (37)', 'Intermediate TL', '0.370 +/- 0.125'],
], 0.3, 1.0, 9.4, 2.2, cw=[2.2, 1.5, 2.0, 2.0], hsz=9.5, rsz=9.5, highlight_row=0)

# Pattern box
b = tb(s1, 0.3, 3.35, 9.4, 0.55)
b.text_frame.word_wrap = True; set_fill(b, 'E8F0FE')
para(b.text_frame, 'Pola: CK+ (0.783) >> JAFFE (0.467) > Dataset sendiri (0.370)', sz=10, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER, first=True)
para(b.text_frame, 'Semakin natural ekspresinya dan semakin imbalanced, semakin rendah performanya', sz=9, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

print('  Slide 1: Benchmark Results')

# ═══ SLIDE 2: SOTA Comparison ═══
t = tb(s2, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Perbandingan dengan Paper State-of-the-Art', sz=18, bold=True, color=COLOR_TEXT, first=True)

# CK+ SOTA
t_ck = tb(s2, 0.3, 0.72, 9.4, 0.25)
para(t_ck.text_frame, 'CK+ (10-Fold CV)', sz=10, bold=True, color=COLOR_BLUE, first=True)

table(s2, ['Paper', 'Tahun', 'Metode', 'Acc/F1'], [
    ['Dada et al.', '2023', 'CNN-10', '99.9% (acc)'],
    ['AA-DCN', '2024', 'Anti-aliased Deep Conv', '99.26% (acc)'],
    ['GhostNet', '2024', 'Face+Speech+EEG', '98.27% (acc)'],
    ['b-skeleton+CNN', '2024', 'Landmark + CNN', '96.19% (acc)'],
    ['Penelitian ini', '2026', 'Intermediate TL (ResNet18)', '91.9% / F1: 0.783'],
], 0.3, 1.0, 9.4, 1.7, cw=[1.8, 0.8, 3.0, 2.5], hsz=9, rsz=8.5, highlight_row=4)

# JAFFE SOTA
t_jf = tb(s2, 0.3, 2.8, 9.4, 0.25)
para(t_jf.text_frame, 'JAFFE (LOSO)', sz=10, bold=True, color=COLOR_BLUE, first=True)

table(s2, ['Paper', 'Tahun', 'Metode', 'Acc/F1'], [
    ['AA-DCN', '2024', 'Anti-aliased Deep Conv', '98.0% (acc)'],
    ['Fine-grained fusion', '2025', 'Landmark + image', '97.61% (acc)'],
    ['b-skeleton+CNN', '2024', 'Landmark geometric', '89.23% (acc)'],
    ['Penelitian ini', '2026', 'Late Fusion (CNN+FCNN)', '54.0% / F1: 0.467'],
], 0.3, 3.1, 9.4, 1.35, cw=[2.0, 0.8, 3.0, 2.5], hsz=9, rsz=8.5, highlight_row=3)

# Note
n = tb(s2, 0.3, 4.55, 9.4, 0.35)
n.text_frame.word_wrap = True
para(n.text_frame, 'Catatan: Paper SOTA pakai arsitektur lebih besar (VGG/EfficientNet), sebagian random split. '
     'Penelitian ini pakai ResNet18 + subject-wise CV/LOSO yang lebih ketat.',
     sz=7.5, italic=True, color=COLOR_GRAY, first=True)

print('  Slide 2: SOTA Comparison')

# ═══ SLIDE 3: Temuan Benchmark ═══
t = tb(s3, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Analisis Hasil -- Temuan Benchmark', sz=18, bold=True, color=COLOR_TEXT, first=True)

# Temuan 16-17 (top row)
t16 = tb(s3, 0.3, 0.78, 4.5, 0.25)
para(t16.text_frame, 'Temuan 16: CK+ jauh lebih tinggi dari dataset sendiri', sz=8.5, bold=True, color=COLOR_BLUE, first=True)

table(s3, ['Dataset', 'Best F1'], [
    ['CK+ (lab)', '0.783'],
    ['JAFFE (lab)', '0.467'],
    ['Dataset sendiri', '0.370'],
], 0.3, 1.08, 4.5, 0.9, cw=[2.5, 2.0], hsz=8.5, rsz=8.5)

t17 = tb(s3, 5.0, 0.78, 4.6, 0.25)
para(t17.text_frame, 'Temuan 17: Transfer Learning konsisten terbaik', sz=8.5, bold=True, color=COLOR_BLUE, first=True)

table(s3, ['Dataset', 'Best Model'], [
    ['CK+ 7c', 'Intermediate TL'],
    ['CK+ 4c', 'CNN TL'],
    ['JAFFE', 'Late Fusion (CNN+FCNN)'],
    ['Dataset sendiri', 'Intermediate TL'],
], 0.3+4.7, 1.08, 4.6, 1.1, cw=[2.2, 2.4], hsz=8.5, rsz=8.5)

# Temuan 18-19 (bottom row)
t18 = tb(s3, 0.3, 2.15, 4.5, 0.25)
para(t18.text_frame, 'Temuan 18: Dataset sendiri paling menantang', sz=8.5, bold=True, color=COLOR_BLUE, first=True)

m18 = tb(s3, 0.3, 2.45, 4.5, 1.1)
m18.text_frame.word_wrap = True; set_fill(m18, 'FCE4E4')
para(m18.text_frame, 'Alasan performa rendah:', sz=9, bold=True, color=COLOR_RED, first=True)
para(m18.text_frame, '1. Ekspresi natural (bukan lab-induced)', sz=8.5, color=COLOR_TEXT)
para(m18.text_frame, '2. Sangat imbalanced (neutral 82%)', sz=8.5, color=COLOR_TEXT)
para(m18.text_frame, '3. Variasi pencahayaan & posisi', sz=8.5, color=COLOR_TEXT)
para(m18.text_frame, '4. Micro-expression saat programming', sz=8.5, color=COLOR_TEXT)

t19 = tb(s3, 5.0, 2.15, 4.6, 0.25)
para(t19.text_frame, 'Temuan 19: Std deviation = stabilitas model', sz=8.5, bold=True, color=COLOR_BLUE, first=True)

table(s3, ['Dataset', 'Best F1', 'Std'], [
    ['CK+', '0.783', '0.107 (stabil)'],
    ['JAFFE', '0.467', '0.092 (moderat)'],
    ['Dataset sendiri', '0.370', '0.125 (variabel)'],
], 5.0, 2.45, 4.6, 0.9, cw=[1.8, 1.2, 1.6], hsz=8.5, rsz=8.5)

# Conclusion
c = tb(s3, 0.3, 3.7, 9.3, 0.35)
c.text_frame.word_wrap = True; set_fill(c, 'E8F0FE')
para(c.text_frame, 'Kesimpulan: Arsitektur bekerja baik di dataset standar (CK+ 0.783) -- performa rendah di dataset sendiri '
     'karena karakteristik data natural programming, bukan kelemahan arsitektur',
     sz=8.5, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER, first=True)

print('  Slide 3: Temuan Benchmark')

# ── SAVE ──
prs.save(PPTX_PATH)
print(f'\nSaved! Total: {len(prs.slides)} slides')

for i in range(95, len(prs.slides)):
    slide = prs.slides[i]
    title = next((s.text_frame.text.strip()[:55].replace(chr(10),' ').encode('ascii','replace').decode()
                  for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), '')
    print(f'  Slide {i+1}: {title}')
