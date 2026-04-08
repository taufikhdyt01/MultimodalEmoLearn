"""
Insert 2 new dedicated table slides:
  - After slide 74 (Rancangan Tahap 2): class merging table
  - After new slide 80 (Rancangan Tahap 3, after offset): TL comparison table

Also removes the existing (broken) tables from slides 74 and 79.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_HEADER_BG   = RGBColor(0x1A, 0x73, 0xE8)
COLOR_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ROW_ALT     = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_ROW_NORMAL  = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT        = RGBColor(0x20, 0x20, 0x20)
COLOR_ORANGE_LIGHT = RGBColor(0xFF, 0xF0, 0xCC)


def insert_blank_slide(prs, index):
    layout = prs.slide_layouts[10]  # blank-ish layout
    slide = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst
    new_elem = xml_slides[-1]
    xml_slides.remove(new_elem)
    xml_slides.insert(index, new_elem)
    return slide


def add_title(slide, text, font_size=20):
    tb = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(9.3), Inches(0.6))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT


def style_cell(cell, text, bold=False, bg=None, fg=None, size=9, align=PP_ALIGN.LEFT):
    cell.text = text
    para = cell.text_frame.paragraphs[0]
    para.alignment = align
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if fg:
        run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def add_table(slide, headers, rows, left, top, width, height,
              col_widths=None, h_size=10, r_size=9):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    tbl = shape.table
    # Remove tableStyleId to avoid PowerPoint rendering conflicts
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        styleId = tblPr.find(qn('a:tableStyleId'))
        if styleId is not None:
            tblPr.remove(styleId)
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / total)
    for c, h in enumerate(headers):
        style_cell(tbl.cell(0, c), h, bold=True,
                   bg=COLOR_HEADER_BG, fg=COLOR_HEADER_TEXT,
                   size=h_size, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_ROW_ALT if r % 2 == 0 else COLOR_ROW_NORMAL
        for c, val in enumerate(row):
            style_cell(tbl.cell(r + 1, c), str(val),
                       bg=bg, fg=COLOR_TEXT, size=r_size,
                       align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)


def strip_table_from_slide(slide, slide_num):
    """Remove the graphicFrame (table) from the slide's spTree."""
    sp_tree = slide.shapes._spTree
    removed = 0
    for child in list(sp_tree):
        tag = child.tag.split('}')[-1]
        if tag == 'graphicFrame':
            sp_tree.remove(child)
            removed += 1
    print(f'  Slide {slide_num}: removed {removed} graphicFrame(s)')


# ── LOAD ────────────────────────────────────────────────
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Remove existing (broken) tables from slides 74 and 79 (indices 73, 78)
print('Removing existing tables from slides 74 and 79...')
strip_table_from_slide(prs.slides[73], 74)
strip_table_from_slide(prs.slides[78], 79)

# ── INSERT SLIDE AFTER 74 (index 74): Class Merging Table ──
s_class = insert_blank_slide(prs, 74)
add_title(s_class, 'Penggabungan Kelas: 7-Class \u2192 4-Class')
add_table(s_class,
    headers=['Kelas Baru', 'Kelas Asal yang Digabung', 'Total Sample', 'Alasan'],
    rows=[
        ['neutral',  'neutral', '8,356', 'Tetap \u2014 kelas dominan'],
        ['happy',    'happy',   '783',   'Tetap \u2014 cukup data'],
        ['sad',      'sad',     '576',   'Tetap \u2014 cukup data'],
        ['negative', 'angry (63) + fearful (13) + disgusted (24) + surprised (79)',
         '179', 'Digabung \u2014 masing-masing < 80 sample'],
    ],
    left=Inches(0.4), top=Inches(1.0), width=Inches(9.2), height=Inches(2.2),
    col_widths=[1.4, 4.0, 1.4, 2.4], h_size=11, r_size=11,
)
tb = s_class.shapes.add_textbox(Inches(0.4), Inches(3.35), Inches(9.2), Inches(0.3))
tf = tb.text_frame
tf.text = '4 model \u00d7 3 skenario = 12 eksperimen tambahan  |  Total kumulatif: 24 eksperimen'
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = COLOR_TEXT

print('Inserted class merging table slide at position 75')

# After insertion, Rancangan Tahap 3 is now at index 79 (was 78), so insert after it at index 80
s_tl = insert_blank_slide(prs, 80)
add_title(s_tl, 'Perubahan Arsitektur: From Scratch \u2192 Transfer Learning')
add_table(s_tl,
    headers=['Komponen', 'Tahap 1 & 2 (From Scratch)', 'Tahap 3 (Transfer Learning)', 'Catatan'],
    rows=[
        ['CNN backbone', 'EmotionCNN (dari nol)', 'ResNet18 pretrained ImageNet',
         'Fine-tune seluruh layer'],
        ['FCNN', 'Fully-connected (dari nol)', 'Sama \u2014 tidak berubah',
         'Landmark = numerik, tidak perlu pretrained'],
        ['Late Fusion', 'CNN scratch + FCNN', 'CNN TL + FCNN', ''],
        ['Intermediate Fusion', 'CNN scratch + FCNN (feature-level)',
         'ResNet18 + FCNN (feature-level)', ''],
        ['Learning Rate', '0.0001', '0.00005', 'Lebih kecil untuk fine-tuning'],
    ],
    left=Inches(0.4), top=Inches(1.0), width=Inches(9.2), height=Inches(2.5),
    col_widths=[2.0, 2.6, 2.6, 2.0], h_size=11, r_size=10,
)
tb = s_tl.shapes.add_textbox(Inches(0.4), Inches(3.65), Inches(9.2), Inches(0.4))
tf = tb.text_frame
tf.text = ('4 model \u00d7 3 skenario \u00d7 2 konfigurasi kelas (7-class + 4-class) '
           '= 24 eksperimen tambahan\nTotal keseluruhan: 12 + 12 + 24 = 48 eksperimen')
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = COLOR_TEXT

print('Inserted TL comparison table slide at position 81')

# ── SAVE ────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f'\nSaved! Total slides: {len(prs.slides)}')

# Verify
prs2 = Presentation(PPTX_PATH)
print('\nSlides 72-88:')
for i in range(71, min(len(prs2.slides), 88)):
    slide = prs2.slides[i]
    title = next((s.text_frame.text.strip()[:60].replace(chr(10), ' ')
                  for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), '')
    print(f'  Slide {i+1}: {title}')