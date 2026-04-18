"""Insert 1 intro slide (instruksi dosen + sumber dataset) BEFORE slide 105."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LB = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_WH = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x20, 0x20, 0x20)
COLOR_GREEN = RGBColor(0x0F, 0x9D, 0x58)
COLOR_OK = RGBColor(0xE6, 0xF4, 0xEA)


def set_fill(shape, hex_val):
    spPr = shape._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', hex_val)


def tb(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def para(tf, text, sz=10, bold=False, italic=False, color=None,
         align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    if color:
        r.font.color.rgb = color


def cell(c, text, bold=False, bg=None, fg=None, sz=9, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    if fg:
        r.font.color.rgb = fg
    if bg:
        c.fill.solid()
        c.fill.fore_color.rgb = bg


def table(slide, headers, rows, l, t, w, h, cw=None,
          hsz=9, rsz=8.5):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                    Inches(l), Inches(t),
                                    Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None:
            tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw):
            tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WH,
             sz=hsz, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LB if r % 2 == 0 else COLOR_WH
        for c, val in enumerate(row):
            cell(tbl.cell(r + 1, c), str(val), bold=(c == 0),
                 bg=bg, fg=COLOR_TEXT, sz=rsz,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


prs = Presentation(PPTX)
initial = len(prs.slides)
print(f'Initial slides: {initial}')

# Add slide at end
blank_layout = prs.slide_layouts[6]
s = prs.slides.add_slide(blank_layout)

# Title
t = tb(s, 0.3, 0.1, 9.4, 0.45)
para(t.text_frame,
     'SLIDE 24A: Instruksi Dosen & Sumber Dataset Benchmark',
     sz=14, bold=True, color=COLOR_TEXT, first=True)

# Instruksi box
ins = tb(s, 0.3, 0.6, 9.4, 1.8)
ins.text_frame.word_wrap = True
set_fill(ins, 'E8F0FE')
para(ins.text_frame, 'Instruksi dari Dosen Pembimbing',
     sz=11, bold=True, color=COLOR_BLUE, first=True)
para(ins.text_frame, '', sz=3)
para(ins.text_frame,
     '1. Skema 1 — Self Train-Test: Tiap dataset (CK+, JAFFE, RAF-DB, KDEF, Primer) '
     'dilatih dan diuji dengan data masing-masing.',
     sz=9.5, color=COLOR_TEXT)
para(ins.text_frame, '', sz=2)
para(ins.text_frame,
     '2. Print semua nilai evaluasi: Macro F1, Micro F1, Weighted F1 (semua dicatat).',
     sz=9.5, color=COLOR_TEXT)
para(ins.text_frame, '', sz=2)
para(ins.text_frame,
     '3. Skema 2 — Cross-Dataset: Dataset sekunder (CK+, JAFFE, RAF-DB, KDEF) '
     'digunakan untuk train, data uji menggunakan data test Primer.',
     sz=9.5, color=COLOR_TEXT)
para(ins.text_frame, '', sz=2)
para(ins.text_frame,
     '4. Model: gunakan semua model yang dimiliki (CNN, FCNN, Intermediate, '
     'Late Fusion, CNN TL, Intermediate TL).',
     sz=9.5, color=COLOR_TEXT)

# Sumber dataset table
para(tb(s, 0.3, 2.5, 9.4, 0.28).text_frame,
     'Sumber Dataset', sz=11, bold=True, color=COLOR_BLUE, first=True)

table(s,
      ['Dataset', 'Sumber', 'Jumlah', 'Karakteristik'],
      [
          ['CK+',
           'Kaggle/GitHub (Cohn-Kanade+)',
           '636 imgs (+18 contempt)',
           'Lab posed, 118 subjek, sequence peak'],
          ['JAFFE',
           'Original (Lyons et al., 1998)',
           '213 imgs',
           'Lab posed, 10 wanita Jepang, 7 emosi'],
          ['RAF-DB',
           'Kaggle shuvoalok/raf-db-dataset (basic emotion public release)',
           '14,449 imgs (11,565 train / 2,884 test)',
           'In-the-wild dari web, official split'],
          ['KDEF',
           'Official kdef.se (KDEF_and_AKDEF.zip)',
           '4,900 imgs → 3,307 usable (face-detected)',
           'Lab posed, 70 subjek, 5 sudut'],
          ['Primer',
           'Akuisisi sendiri (sesi programming)',
           '7,091 imgs (5,287 train / 579 val / 929 test)',
           'Natural expression, 37 subjek, front-only'],
      ],
      0.3, 2.8, 9.4, 2.1,
      cw=[1.0, 3.2, 2.3, 2.9], hsz=9, rsz=8)

# Note
note = tb(s, 0.3, 5.0, 9.4, 0.5)
note.text_frame.word_wrap = True
set_fill(note, 'E6F4EA')
para(note.text_frame,
     'Catatan: RAF-DB pakai basic-emotion public release (15,339 imgs official) — BUKAN subset. '
     'KDEF ~32% di-drop karena MediaPipe gagal deteksi wajah di sudut profil penuh.',
     sz=8.5, italic=True, color=COLOR_TEXT, first=True)

# ─────────── Move new slide to position 104 (becomes slide 105) ───────────
# Target: insert after slide index 103 (0-based) = slide 104 "Analisis Temuan Benchmark"
INSERT_AFTER_IDX = 103

xml_slides = prs.slides._sldIdLst
slides_list = list(xml_slides)
new_slide = slides_list[-1]
xml_slides.remove(new_slide)

slides_after = list(xml_slides)
target_elem = slides_after[INSERT_AFTER_IDX]
target_elem.addnext(new_slide)

prs.save(PPTX)
print(f'Final slides: {len(prs.slides)}')
print(f'New intro slide inserted at position: {INSERT_AFTER_IDX + 2}')
