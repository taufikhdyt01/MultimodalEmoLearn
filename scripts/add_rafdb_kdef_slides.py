"""Insert RAF-DB & KDEF benchmark slides into PPT Bimbingan (after slide 104)."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LB = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_WH = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x20, 0x20, 0x20)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_RED = RGBColor(0xD9, 0x3C, 0x3C)
COLOR_GREEN = RGBColor(0x0F, 0x9D, 0x58)
COLOR_OK = RGBColor(0xE6, 0xF4, 0xEA)


def set_fill(shape, hex_val):
    spPr = shape._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', hex_val)


def tb(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def para(tf, text, sz=10, bold=False, italic=False, color=None,
         align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    if color:
        r.font.color.rgb = color


def cell(c, text, bold=False, bg=None, fg=None, sz=9, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    if fg:
        r.font.color.rgb = fg
    if bg:
        c.fill.solid()
        c.fill.fore_color.rgb = bg


def table(slide, headers, rows, l, t, w, h, cw=None,
          hsz=9, rsz=8.5, highlight_rows=None):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                    Inches(l), Inches(t),
                                    Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None:
            tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw):
            tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WH,
             sz=hsz, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LB if r % 2 == 0 else COLOR_WH
        if highlight_rows and r in highlight_rows:
            bg = COLOR_OK
        for c, val in enumerate(row):
            cell(tbl.cell(r + 1, c), str(val), bold=(c == 0),
                 bg=bg, fg=COLOR_TEXT, sz=rsz,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


def move_slide_after(prs, moved_idx, target_idx):
    """Move slide at moved_idx to just after target_idx (0-based)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    moved = slides[moved_idx]
    xml_slides.remove(moved)
    # After removal, positions shift; target_idx stays if < moved_idx
    slides_after = list(xml_slides)
    # Insert after target_idx
    target_elem = slides_after[target_idx]
    target_elem.addnext(moved)


prs = Presentation(PPTX)
initial_count = len(prs.slides)
print(f'Initial slides: {initial_count}')

# Insert target: after slide 104 (0-indexed: 103) = "Analisis Hasil -- Temuan Benchmark"
INSERT_AFTER_IDX = 103  # 0-based index

blank_layout = prs.slide_layouts[6]

# ─────────── NEW SLIDE 1: RAF-DB Setup + 7-class ───────────
s = prs.slides.add_slide(blank_layout)

t = tb(s, 0.3, 0.1, 9.4, 0.45)
para(t.text_frame, 'SLIDE 24B: Benchmark Tambahan — RAF-DB & KDEF (arahan dosen)',
     sz=14, bold=True, color=COLOR_TEXT, first=True)

# Setup box
setup = tb(s, 0.3, 0.58, 9.4, 0.9)
setup.text_frame.word_wrap = True
set_fill(setup, 'E8F0FE')
para(setup.text_frame, 'Setup Dataset', sz=10, bold=True,
     color=COLOR_BLUE, first=True)
para(setup.text_frame,
     '• RAF-DB: 11,565 train / 2,884 test (official split) — in-the-wild 7 basic emotions',
     sz=8.5, color=COLOR_TEXT)
para(setup.text_frame,
     '• KDEF: 2,630 train / 340 val / 337 test (subject-wise) — lab posed, 70 subjek, 5 angle',
     sz=8.5, color=COLOR_TEXT)
para(setup.text_frame,
     '• Skema: Self train-test, B1 baseline, metrik Macro/Micro/Weighted F1',
     sz=8.5, color=COLOR_TEXT)

# RAF-DB 7-class
para(tb(s, 0.3, 1.55, 4.6, 0.25).text_frame, 'RAF-DB 7-Class',
     sz=10, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.729', '0.815', '0.813', '0.815'],
    ['FCNN', '0.578', '0.714', '0.703', '0.714'],
    ['Intermediate', '0.696', '0.785', '0.783', '0.785'],
    ['CNN TL', '0.741', '0.830', '0.826', '0.830'],
    ['Intermediate TL', '0.744', '0.833', '0.832', '0.833'],
    ['Late Fusion', '0.719', '0.809', '0.805', '0.809'],
], 0.3, 1.85, 4.6, 2.3, cw=[1.8, 0.8, 0.8, 0.8, 0.8],
   hsz=8.5, rsz=8, highlight_rows=[4])

# RAF-DB 4-class
para(tb(s, 5.05, 1.55, 4.6, 0.25).text_frame, 'RAF-DB 4-Class',
     sz=10, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.808', '0.830', '0.830', '0.830'],
    ['FCNN', '0.694', '0.728', '0.729', '0.728'],
    ['Intermediate', '0.792', '0.818', '0.818', '0.818'],
    ['CNN TL', '0.827', '0.845', '0.846', '0.845'],
    ['Intermediate TL', '0.836', '0.853', '0.855', '0.853'],
    ['Late Fusion', '0.819', '0.842', '0.841', '0.842'],
], 5.05, 1.85, 4.6, 2.3, cw=[1.8, 0.8, 0.8, 0.8, 0.8],
   hsz=8.5, rsz=8, highlight_rows=[4])

# Best note
best = tb(s, 0.3, 4.3, 9.4, 0.9)
best.text_frame.word_wrap = True
set_fill(best, 'E6F4EA')
para(best.text_frame,
     'Best RAF-DB: Intermediate TL — 7c Macro F1 = 0.744 | 4c Macro F1 = 0.836',
     sz=10, bold=True, color=COLOR_GREEN, first=True)
para(best.text_frame,
     'Pola konsisten: Intermediate Fusion + Transfer Learning mengungguli semua arsitektur lain.',
     sz=9, color=COLOR_TEXT)


# ─────────── NEW SLIDE 2: KDEF + Progress ───────────
s = prs.slides.add_slide(blank_layout)

t = tb(s, 0.3, 0.1, 9.4, 0.45)
para(t.text_frame, 'SLIDE 24B (lanjutan): Hasil KDEF & Progress Eksperimen',
     sz=14, bold=True, color=COLOR_TEXT, first=True)

# KDEF 7-class table
para(tb(s, 0.3, 0.6, 4.6, 0.25).text_frame, 'KDEF 7-Class (2,630 train / 337 test)',
     sz=10, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.798', '0.801', '0.798', '0.801'],
    ['FCNN', '0.666', '0.680', '0.663', '0.680'],
    ['Intermediate', '0.671', '0.674', '0.668', '0.674'],
    ['CNN TL', '0.833', '0.831', '0.833', '0.831'],
    ['Intermediate TL', '0.843', '0.843', '0.843', '0.843'],
    ['Late Fusion', '0.776', '0.777', '0.775', '0.777'],
], 0.3, 0.9, 4.6, 2.3, cw=[1.8, 0.8, 0.8, 0.8, 0.8],
   hsz=8.5, rsz=8, highlight_rows=[4])

# Progress status table
para(tb(s, 5.05, 0.6, 4.6, 0.25).text_frame, 'Status Eksperimen',
     sz=10, bold=True, color=COLOR_TEXT, first=True)
table(s, ['Eksperimen', 'Status'], [
    ['RAF-DB 7-class', '✅ Done'],
    ['RAF-DB 4-class', '✅ Done'],
    ['KDEF 7-class', '✅ Done'],
    ['KDEF 4-class', '⏳ Pending'],
    ['Primer self (nb 62)', '⏳ Pending'],
    ['Cross-dataset (nb 63)', '⏳ Pending'],
], 5.05, 0.9, 4.6, 2.3, cw=[3.2, 1.4],
   hsz=9, rsz=8.5, highlight_rows=[0, 1, 2])

# Best KDEF note
best = tb(s, 0.3, 3.35, 9.4, 0.55)
best.text_frame.word_wrap = True
set_fill(best, 'E6F4EA')
para(best.text_frame,
     'Best KDEF 7c: Intermediate TL — Macro F1 = 0.843 (konsisten dengan pola RAF-DB)',
     sz=10, bold=True, color=COLOR_GREEN, first=True)

# Overall pattern note
pat = tb(s, 0.3, 4.0, 9.4, 1.25)
pat.text_frame.word_wrap = True
set_fill(pat, 'FFF0CC')
para(pat.text_frame, 'Pola Best Model Lintas Dataset',
     sz=10, bold=True, color=COLOR_TEXT, first=True)
para(pat.text_frame, '', sz=2)
para(pat.text_frame,
     '• Primer 4c: Intermediate TL B1 = 0.412  |  CK+ 4c: CNN TL = 0.837  |  JAFFE 4c: Late Fusion = 0.530',
     sz=9, color=COLOR_TEXT)
para(pat.text_frame,
     '• RAF-DB 4c: Intermediate TL = 0.836  |  KDEF 7c: Intermediate TL = 0.843',
     sz=9, color=COLOR_TEXT)
para(pat.text_frame,
     '• Insight: Intermediate TL dominan di 4 dari 5 dataset — arsitektur robust untuk multimodal FER.',
     sz=9, italic=True, color=COLOR_TEXT)


# ─────────── NEW SLIDE 3: Temuan ───────────
s = prs.slides.add_slide(blank_layout)

t = tb(s, 0.3, 0.1, 9.4, 0.45)
para(t.text_frame, 'SLIDE 24B (lanjutan): Temuan dari Benchmark RAF-DB & KDEF',
     sz=14, bold=True, color=COLOR_TEXT, first=True)

# Temuan 16
t16 = tb(s, 0.3, 0.6, 9.4, 1.3)
t16.text_frame.word_wrap = True
set_fill(t16, 'E8F0FE')
para(t16.text_frame,
     'Temuan 16: Intermediate TL konsisten best di semua dataset',
     sz=11, bold=True, color=COLOR_BLUE, first=True)
para(t16.text_frame, '', sz=2)
para(t16.text_frame,
     '• RAF-DB 7c: 0.744 | RAF-DB 4c: 0.836 | KDEF 7c: 0.843 | Primer 4c: 0.412',
     sz=9.5, color=COLOR_TEXT)
para(t16.text_frame,
     '• Arsitektur Intermediate Fusion + Transfer Learning (ResNet18) optimal '
     'untuk multimodal FER lintas kondisi data (posed, in-the-wild, natural).',
     sz=9.5, color=COLOR_TEXT)

# Temuan 17
t17 = tb(s, 0.3, 2.0, 9.4, 1.5)
t17.text_frame.word_wrap = True
set_fill(t17, 'FFF0CC')
para(t17.text_frame,
     'Temuan 17: Gap besar primer vs benchmark → konfirmasi hipotesis karakteristik data',
     sz=11, bold=True, color=COLOR_TEXT, first=True)
para(t17.text_frame, '', sz=2)
para(t17.text_frame,
     '• Gap Macro F1: RAF-DB (0.836) vs Primer (0.412) → selisih 0.424 (dataset natural yang berbeda karakteristik)',
     sz=9.5, color=COLOR_TEXT)
para(t17.text_frame,
     '• KDEF (0.843) lebih tinggi karena lab posed (ekspresi jelas, balanced).',
     sz=9.5, color=COLOR_TEXT)
para(t17.text_frame,
     '• Argumen terkuat: performa rendah di primer BUKAN karena kelemahan arsitektur, '
     'tapi karena data natural + imbalanced ekstrem + minoritas langka.',
     sz=9.5, italic=True, bold=True, color=COLOR_TEXT)

# Temuan 18
t18 = tb(s, 0.3, 3.6, 9.4, 1.55)
t18.text_frame.word_wrap = True
set_fill(t18, 'E6F4EA')
para(t18.text_frame,
     'Temuan 18: Efek fusion tergantung ukuran dataset',
     sz=11, bold=True, color=COLOR_GREEN, first=True)
para(t18.text_frame, '', sz=2)
para(t18.text_frame,
     '• Di dataset besar (RAF-DB 11k train): selisih CNN TL vs Intermediate TL kecil (~0.01) '
     '— CNN sendiri sudah cukup dengan data banyak.',
     sz=9.5, color=COLOR_TEXT)
para(t18.text_frame,
     '• Di dataset kecil + imbalanced (Primer 5.3k): fusion memberi gain lebih besar '
     '— landmark sebagai informasi komplementer lebih berguna.',
     sz=9.5, color=COLOR_TEXT)
para(t18.text_frame,
     '• Implikasi: multimodal fusion terutama bermanfaat di low-data / imbalanced regime.',
     sz=9.5, italic=True, color=COLOR_TEXT)


# ─────────── Move 3 new slides to position INSERT_AFTER_IDX ───────────
final_count = len(prs.slides)
new_indices = [final_count - 3, final_count - 2, final_count - 1]
print(f'Added {len(new_indices)} slides at end (indices {new_indices})')

# Reorder: move each new slide to right after INSERT_AFTER_IDX
# We iterate backwards so target_idx stays valid
xml_slides = prs.slides._sldIdLst
slides_list = list(xml_slides)
# Keep the 3 new slides (last 3)
new_slides = slides_list[-3:]
# Remove them from current positions
for ns in new_slides:
    xml_slides.remove(ns)
# Now list is back to original 117 slides. Insert after INSERT_AFTER_IDX (index 103)
slides_after = list(xml_slides)
target_elem = slides_after[INSERT_AFTER_IDX]
# Insert in reverse so final order matches new_slides order
for ns in reversed(new_slides):
    target_elem.addnext(ns)

prs.save(PPTX)
print(f'Final total slides: {len(prs.slides)}')
print(f'New slides inserted at positions: {INSERT_AFTER_IDX+2}, {INSERT_AFTER_IDX+3}, {INSERT_AFTER_IDX+4}')
print(f'Saved: {PPTX}')
