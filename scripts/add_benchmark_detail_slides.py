"""
Add detailed benchmark slides before the summary slides (97-99).
Insert at index 96 (before current slide 97 Benchmark Results).

Slides:
1. Dataset Benchmark (info table)
2. Hasil Single Split 7-Class (6 models × 3 datasets)
3. Hasil Single Split 4-Class
4. Hasil LOSO/CV 7-Class
5. Hasil LOSO/CV 4-Class
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE       = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_TEXT       = RGBColor(0x20, 0x20, 0x20)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY       = RGBColor(0x66, 0x66, 0x66)
COLOR_GREEN_LIGHT = RGBColor(0xC6, 0xEF, 0xCE)


def set_fill(shape, hex_val):
    spPr = shape._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', hex_val)


def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def para(tf, text, sz=10, bold=False, italic=False, color=None, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = color


def cell(c, text, bold=False, bg=None, fg=None, sz=9, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]; p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    if fg: r.font.color.rgb = fg
    if bg: c.fill.solid(); c.fill.fore_color.rgb = bg


def table(slide, headers, rows, l, t, w, h, cw=None, hsz=9, rsz=9, best_cells=None):
    """best_cells: list of (row, col) tuples to highlight green"""
    shape = slide.shapes.add_table(len(rows)+1, len(headers), Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None: tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw): tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WHITE, sz=hsz, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LIGHT_BLUE if r % 2 == 0 else COLOR_WHITE
        for c, val in enumerate(row):
            is_best = best_cells and (r, c) in best_cells
            bg_c = COLOR_GREEN_LIGHT if is_best else bg
            cell(tbl.cell(r+1, c), str(val), bold=(c==0 or is_best),
                 bg=bg_c, fg=COLOR_TEXT, sz=rsz,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


def insert(prs, index):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    xml_slides = prs.slides._sldIdLst
    el = xml_slides[-1]; xml_slides.remove(el); xml_slides.insert(index, el)
    return slide


# ── LOAD ──
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Insert 5 slides at index 96 (before Benchmark Results summary)
slides = []
for i in range(5):
    slides.append(insert(prs, 96 + i))

# ═══ SLIDE 1: Dataset Benchmark ═══
s = slides[0]
t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Benchmark -- Dataset JAFFE & CK+', sz=20, bold=True, color=COLOR_TEXT, first=True)

# Motivation box
m = tb(s, 0.3, 0.8, 9.4, 0.8)
m.text_frame.word_wrap = True; set_fill(m, 'FFF0CC')
para(m.text_frame, 'Tujuan: Menguji pipeline dan arsitektur yang sama pada dataset standar untuk', sz=9.5, bold=True, color=COLOR_TEXT, first=True)
para(m.text_frame, 'menunjukkan bahwa performa rendah di dataset sendiri disebabkan oleh', sz=9.5, color=COLOR_TEXT)
para(m.text_frame, 'karakteristik data (natural expression, imbalanced), bukan kelemahan arsitektur.', sz=9.5, color=COLOR_TEXT)

table(s, ['Dataset', 'Sampel', 'Emosi', 'Subjek', 'Karakteristik'], [
    ['JAFFE', '213', '7', '10', 'Lab, wanita Jepang, seimbang'],
    ['CK+', '636', '7+contempt', '118', 'Lab, ekspresi peak, semi-balanced'],
    ['Dataset sendiri', '7,091', '7', '37', 'Natural (programming), imbalanced'],
], 0.3, 1.75, 9.4, 1.2, cw=[2.0, 1.0, 0.8, 1.0, 4.6], hsz=10, rsz=10)

# Evaluation info
e = tb(s, 0.3, 3.1, 9.4, 0.8)
e.text_frame.word_wrap = True; set_fill(e, 'E8F0FE')
para(e.text_frame, 'Evaluasi yang dilakukan:', sz=10, bold=True, color=COLOR_TEXT, first=True)
para(e.text_frame, '1. Single Split (80/10/10 subject-wise) -- semua model, B1 Baseline', sz=9.5, color=COLOR_TEXT)
para(e.text_frame, '2. JAFFE: LOSO (10 fold = 10 subjek) -- semua model', sz=9.5, color=COLOR_TEXT)
para(e.text_frame, '3. CK+: 10-Fold CV (subject-wise) -- semua model', sz=9.5, color=COLOR_TEXT)
print('  Slide 1: Dataset Benchmark')

# ═══ SLIDE 2: Single Split 7-Class ═══
s = slides[1]
t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Benchmark 7-Class -- Single Split (Macro F1)', sz=20, bold=True, color=COLOR_TEXT, first=True)

table(s, ['Model', 'JAFFE', 'CK+', 'Dataset Sendiri'], [
    ['CNN', '0.304', '0.461', '0.137'],
    ['FCNN', '0.209', '0.395', '0.158'],
    ['Late Fusion', '0.545', '0.498', '0.175'],
    ['Intermediate', '0.037', '0.316', '0.137'],
    ['CNN TL', '0.464', '0.913', '0.154'],
    ['Intermediate TL', '0.447', '0.833', '0.180'],
], 0.3, 0.8, 9.4, 2.4, cw=[2.5, 2.3, 2.3, 2.3], hsz=10, rsz=10,
   best_cells=[(2, 1), (4, 2)])  # Late Fusion best JAFFE, CNN TL best CK+

b = tb(s, 0.3, 3.3, 9.4, 0.3)
para(b.text_frame, 'Best: JAFFE = Late Fusion (0.545) | CK+ = CNN TL (0.913) | Dataset sendiri = Intermediate TL (0.180)',
     sz=9, bold=True, color=COLOR_TEXT, first=True)
print('  Slide 2: Single Split 7-Class')

# ═══ SLIDE 3: Single Split 4-Class ═══
s = slides[2]
t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Benchmark 4-Class -- Single Split (Macro F1)', sz=20, bold=True, color=COLOR_TEXT, first=True)

table(s, ['Model', 'JAFFE', 'CK+', 'Dataset Sendiri'], [
    ['CNN', '0.177', '0.645', '0.238'],
    ['FCNN', '0.438', '0.592', '0.361'],
    ['Late Fusion', '0.396', '0.592', '0.394'],
    ['Intermediate', '0.177', '0.567', '0.243'],
    ['CNN TL', '0.330', '0.675', '0.274'],
    ['Intermediate TL', '0.375', '0.837', '0.412'],
], 0.3, 0.8, 9.4, 2.4, cw=[2.5, 2.3, 2.3, 2.3], hsz=10, rsz=10,
   best_cells=[(1, 1), (5, 2), (5, 3)])  # FCNN best JAFFE, IntTL best CK+ & sendiri

b = tb(s, 0.3, 3.3, 9.4, 0.3)
para(b.text_frame, 'Best: JAFFE = FCNN (0.438) | CK+ = Intermediate TL (0.837) | Dataset sendiri = Intermediate TL (0.412)',
     sz=9, bold=True, color=COLOR_TEXT, first=True)
print('  Slide 3: Single Split 4-Class')

# ═══ SLIDE 4: LOSO/CV 7-Class ═══
s = slides[3]
t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Benchmark 7-Class -- JAFFE LOSO & CK+ 10-Fold CV', sz=18, bold=True, color=COLOR_TEXT, first=True)

table(s, ['Model', 'JAFFE LOSO', 'CK+ 10-Fold CV', 'Dataset Sendiri'], [
    ['CNN', '0.249 +/- 0.111', '0.404 +/- 0.049', '-'],
    ['FCNN', '0.304 +/- 0.157', '0.478 +/- 0.022', '-'],
    ['Late Fusion', '0.467 +/- 0.092', '0.544 +/- 0.060', '-'],
    ['Intermediate', '0.129 +/- 0.070', '0.226 +/- 0.082', '-'],
    ['CNN TL', '0.426 +/- 0.143', '0.734 +/- 0.082', '-'],
    ['Intermediate TL', '0.293 +/- 0.156', '0.783 +/- 0.107', '0.370 +/- 0.125'],
], 0.3, 0.8, 9.4, 2.4, cw=[2.2, 2.4, 2.4, 2.4], hsz=9.5, rsz=9,
   best_cells=[(2, 1), (5, 2)])

b = tb(s, 0.3, 3.3, 9.4, 0.3)
para(b.text_frame, 'Best: JAFFE = Late Fusion (0.467) | CK+ = Intermediate TL (0.783)',
     sz=9, bold=True, color=COLOR_TEXT, first=True)
print('  Slide 4: LOSO/CV 7-Class')

# ═══ SLIDE 5: LOSO/CV 4-Class ═══
s = slides[4]
t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Benchmark 4-Class -- JAFFE LOSO & CK+ 10-Fold CV', sz=18, bold=True, color=COLOR_TEXT, first=True)

table(s, ['Model', 'JAFFE LOSO', 'CK+ 10-Fold CV', 'Dataset Sendiri'], [
    ['CNN', '0.338 +/- 0.161', '0.584 +/- 0.163', '-'],
    ['FCNN', '0.431 +/- 0.194', '0.598 +/- 0.036', '0.399 +/- 0.062*'],
    ['Late Fusion', '0.530 +/- 0.126', '0.621 +/- 0.031', '0.401 +/- 0.055*'],
    ['Intermediate', '0.202 +/- 0.066', '0.458 +/- 0.172', '-'],
    ['CNN TL', '0.510 +/- 0.155', '0.755 +/- 0.079', '-'],
    ['Intermediate TL', '0.450 +/- 0.214', '0.715 +/- 0.054', '0.435 +/- 0.068*'],
], 0.3, 0.8, 9.4, 2.4, cw=[2.2, 2.4, 2.4, 2.4], hsz=9.5, rsz=9,
   best_cells=[(2, 1), (4, 2), (5, 3)])

b = tb(s, 0.3, 3.3, 9.4, 0.35)
b.text_frame.word_wrap = True
para(b.text_frame, 'Best: JAFFE = Late Fusion (0.530) | CK+ = CNN TL (0.755) | Dataset sendiri = Intermediate TL (0.435)',
     sz=9, bold=True, color=COLOR_TEXT, first=True)
para(b.text_frame, '*5-Fold CV (dataset sendiri) | LOSO Intermediate TL: 0.370 +/- 0.125',
     sz=8, italic=True, color=COLOR_GRAY)
print('  Slide 5: LOSO/CV 4-Class')

# ── SAVE ──
prs.save(PPTX_PATH)
print(f'\nSaved! Total: {len(prs.slides)} slides')

for i in range(95, len(prs.slides)):
    slide = prs.slides[i]
    title = next((s.text_frame.text.strip()[:58].replace(chr(10),' ').encode('ascii','replace').decode()
                  for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), '')
    print(f'  Slide {i+1}: {title}')
