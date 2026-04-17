"""Rebuild dataset problem analysis and solution slides."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'


def set_fill(shape, hex_val):
    spPr = shape._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', hex_val)


def tb(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def para(tf, text, sz=10, bold=False, italic=False, color=None, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    if color:
        r.font.color.rgb = color


def cell(c, text, bold=False, bg=None, fg=None, sz=9, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    if fg:
        r.font.color.rgb = fg
    if bg:
        c.fill.solid()
        c.fill.fore_color.rgb = bg


def table(slide, headers, rows, l, t, w, h, cw=None, hsz=9, rsz=8.5, warning_rows=None):
    COLOR_BLUE = RGBColor(0x1A, 0x73, 0xE8)
    COLOR_LB = RGBColor(0xE8, 0xF0, 0xFE)
    COLOR_WH = RGBColor(0xFF, 0xFF, 0xFF)
    COLOR_TEXT = RGBColor(0x20, 0x20, 0x20)
    COLOR_WARN = RGBColor(0xFC, 0xE4, 0xE4)

    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                    Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None:
            tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw):
            tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WH,
             sz=hsz, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LB if r % 2 == 0 else COLOR_WH
        if warning_rows and r in warning_rows:
            bg = COLOR_WARN
        for c, val in enumerate(row):
            cell(tbl.cell(r + 1, c), str(val), bold=(c == 0),
                 bg=bg, fg=COLOR_TEXT, sz=rsz,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


COLOR_TEXT = RGBColor(0x20, 0x20, 0x20)
COLOR_RED = RGBColor(0xD9, 0x3C, 0x3C)
COLOR_GREEN = RGBColor(0x0F, 0x9D, 0x58)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)


def clear(slide):
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            sp_tree.remove(child)


prs = Presentation(PPTX)

# Find the 2 slides
slide_problem = None
slide_solution = None
for i in range(100, len(prs.slides)):
    for s in prs.slides[i].shapes:
        if s.has_text_frame:
            t = s.text_frame.text
            if 'Analisis Mendalam Masalah' in t:
                slide_problem = prs.slides[i]
            elif 'Solusi yang Dapat' in t:
                slide_solution = prs.slides[i]

print(f'Problem: {slide_problem is not None}, Solution: {slide_solution is not None}')

# ─── SLIDE 1: Problem Analysis (complete) ───
clear(slide_problem)
s = slide_problem

t = tb(s, 0.3, 0.1, 9.4, 0.45)
para(t.text_frame, 'Analisis Mendalam Masalah Dataset',
     sz=17, bold=True, color=COLOR_TEXT, first=True)

# Masalah 1: Per-user diversity
p1 = tb(s, 0.3, 0.6, 4.5, 0.25)
para(p1.text_frame, 'Masalah 1: Diversity Per-User (12/37 = 32%)',
     sz=9.5, bold=True, color=COLOR_RED, first=True)

table(s, ['Kategori', 'User ID'], [
    ['Extreme 100% neutral', '115, 209'],
    ['Extreme 95%+ neutral', '112, 200, 201, 203, 206, 207'],
    ['Terlalu sedikit non-neutral', '101, 113, 118, 197'],
], 0.3, 0.87, 4.5, 1.2, cw=[2.8, 3.2], hsz=9, rsz=8, warning_rows=[0, 1])

# Masalah 2: Minority class
p2 = tb(s, 4.95, 0.6, 4.75, 0.25)
para(p2.text_frame, 'Masalah 2: Kelas Minoritas Langka',
     sz=9.5, bold=True, color=COLOR_RED, first=True)

table(s, ['Emosi', 'Total', 'Conf>=60%', 'Conf>=90%'], [
    ['Neutral', '6,054', '5,906', '5,242'],
    ['Happy', '699', '667', '514'],
    ['Sad', '451', '373', '212'],
    ['Angry', '56', '35', '8'],
    ['Fearful', '10', '5', '3'],
    ['Disgusted', '21', '17', '5'],
    ['Surprised', '55', '39', '19'],
], 4.95, 0.87, 4.75, 2.0, cw=[1.8, 1.2, 1.4, 1.4],
      hsz=8.5, rsz=8, warning_rows=[3, 4, 5, 6])

# Masalah 3: Test set
p3 = tb(s, 0.3, 2.2, 4.5, 0.25)
para(p3.text_frame, 'Masalah 3: Test Set Reliability',
     sz=9.5, bold=True, color=COLOR_RED, first=True)

table(s, ['Test Size', 'Stabilitas F1', 'Contoh'], [
    ['< 10', 'Tidak reliable', '-'],
    ['10-30', 'Hati-hati (variance tinggi)', 'Happy=10, Negative=16'],
    ['30-100', 'Cukup reliable', 'Sad=29'],
    ['> 100', 'Reliable', 'Neutral=950'],
], 0.3, 2.47, 4.5, 1.5, cw=[1.2, 2.5, 2.5],
      hsz=8.5, rsz=8, warning_rows=[0, 1])

# Implikasi box
imp = tb(s, 4.95, 2.92, 4.75, 1.1)
imp.text_frame.word_wrap = True
set_fill(imp, 'FCE4E4')
para(imp.text_frame, 'Implikasi Kritis:', sz=9.5, bold=True, color=COLOR_RED, first=True)
para(imp.text_frame, '', sz=3)
para(imp.text_frame, '- Happy 10 sampel: 1 salah = 10% F1', sz=8.5, color=COLOR_TEXT)
para(imp.text_frame, '- Negative 16 sampel: sangat variable', sz=8.5, color=COLOR_TEXT)
para(imp.text_frame, '- Neutral 950 sampel: 1 salah = 0.1%', sz=8.5, color=COLOR_TEXT)
para(imp.text_frame, '', sz=3)
para(imp.text_frame, 'Variance tinggi = evaluasi tidak stabil',
     sz=8.5, italic=True, bold=True, color=COLOR_RED)

# Bottom conclusion
b = tb(s, 0.3, 4.15, 9.4, 0.3)
b.text_frame.word_wrap = True
set_fill(b, 'FFF0CC')
para(b.text_frame, 'Kesimpulan: F1 0.567 mungkin sudah batas maksimal dengan dataset ini. Perlu konsultasi solusi lanjutan.',
     sz=9, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER, first=True)

print('Slide 1 updated')

# ─── SLIDE 2: Solution Options ───
clear(slide_solution)
s = slide_solution

t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Solusi yang Dapat Dikonsultasikan',
     sz=18, bold=True, color=COLOR_TEXT, first=True)

intro = tb(s, 0.3, 0.75, 9.4, 0.35)
intro.text_frame.word_wrap = True
set_fill(intro, 'E8F0FE')
para(intro.text_frame,
     '3 opsi yang dipertimbangkan + mohon arahan jika ada solusi lain dari Bapak/Ibu',
     sz=9.5, italic=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER, first=True)

options = [
    ('A', 'Drop Problem Users',
     'Hilangkan 4 user extreme (112, 115, 200, 209)',
     'Pro: Data lebih bersih | Kontra: Test set berkurang', 'FFF0CC'),
    ('B', 'Merge ke 3-Class',
     'neutral / happy / negative (gabung sad+angry+fearful+disgusted+surprised)',
     'Pro: Test set lebih reliable (negative=55) | Kontra: Kehilangan granularitas emosi', 'E6F4EA'),
    ('C', 'Binary Classification',
     'neutral vs non-neutral saja (1 vs all)',
     'Pro: Paling mudah dicapai | Kontra: Scope sempit dari proposal awal', 'FCE4E4'),
]

y = 1.2
for letter, title, desc, pros_cons, bg in options:
    box = tb(s, 0.3, y, 9.4, 0.85)
    box.text_frame.word_wrap = True
    set_fill(box, bg)
    para(box.text_frame, 'Opsi ' + letter + ': ' + title,
         sz=10, bold=True, color=COLOR_TEXT, first=True)
    para(box.text_frame, desc, sz=9, color=COLOR_TEXT)
    para(box.text_frame, pros_cons, sz=8.5, italic=True, color=COLOR_GRAY)
    y += 0.9

# Question for supervisor
q = tb(s, 0.3, 3.95, 9.4, 0.5)
q.text_frame.word_wrap = True
set_fill(q, '1A73E8')
para(q.text_frame, 'Pertanyaan untuk Pembimbing:',
     sz=10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
     align=PP_ALIGN.CENTER, first=True)
para(q.text_frame,
     'Apakah Bapak/Ibu punya solusi lain? Atau saran opsi mana yang paling tepat?',
     sz=9, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

print('Slide 2 updated')

prs.save(PPTX)
print('Saved!')
