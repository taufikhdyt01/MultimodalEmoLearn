"""
Add 2 'Analisis Hasil' slides after the front-only results:
- Slide A: Temuan 10-12 (Front-Only vs Front+Side)
- Slide B: Temuan 13-15 (Robustness & Data Leakage)

Format: 2x2 layout with label + mini table (same as slide 74).
Insert after slide 91 (Perbandingan Front-Only vs Front+Side).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE       = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_TEXT       = RGBColor(0x20, 0x20, 0x20)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)


def set_fill(shape, hex_val):
    spPr = shape._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', hex_val)


def add_tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def add_para(tf, text, sz=10, bold=False, color=None, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    if color: r.font.color.rgb = color


def cell(c, text, bold=False, bg=None, fg=None, sz=8.5, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]; p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    if fg: r.font.color.rgb = fg
    if bg: c.fill.solid(); c.fill.fore_color.rgb = bg


def mini_table(slide, headers, rows, l, t, w, h, cw=None):
    shape = slide.shapes.add_table(len(rows)+1, len(headers), Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None: tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw): tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WHITE, sz=8, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LIGHT_BLUE if r % 2 == 0 else COLOR_WHITE
        for c, val in enumerate(row):
            cell(tbl.cell(r+1, c), str(val), bold=(c==0), bg=bg, fg=COLOR_TEXT, sz=8.5,
                 align=PP_ALIGN.LEFT if c==0 else PP_ALIGN.CENTER)


def insert(prs, index):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    xml_slides = prs.slides._sldIdLst
    el = xml_slides[-1]; xml_slides.remove(el); xml_slides.insert(index, el)
    return slide


# ── LOAD ──
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Insert 2 slides after slide 91 (Perbandingan) = index 91
# Slide 92 = Temuan 10-12, Slide 93 = Temuan 13-15
# Current 92 (Evaluasi Robustness) and 93 (Data Leakage) will shift to 94, 95

s1 = insert(prs, 91)
s2 = insert(prs, 92)

# ═══ SLIDE 92: Temuan 10-12 (Front-Only vs Front+Side) ═══
# Title
t = add_tb(s1, 0.3, 0.15, 9.3, 0.6)
add_para(t.text_frame, 'Analisis Hasil \u2014 Temuan Front-Only vs Front+Side',
         sz=20, bold=True, color=COLOR_TEXT, first=True)

# Temuan 10 (top-left)
t10 = add_tb(s1, 0.3, 0.85, 4.5, 0.28)
add_para(t10.text_frame, 'Temuan 10: Side view tidak meningkatkan best model',
         sz=9, bold=True, color=COLOR_BLUE, first=True)

mini_table(s1, ['', 'Best Macro F1'], [
    ['Front-only', '0.412'],
    ['Front+side', '0.407'],
    ['Selisih', '+0.005'],
], 0.3, 1.18, 4.5, 1.0, cw=[2.5, 2.0])

# Temuan 11 (top-right)
t11 = add_tb(s1, 5.0, 0.85, 4.6, 0.28)
add_para(t11.text_frame, 'Temuan 11: FCNN paling terdampak penghapusan side view',
         sz=9, bold=True, color=COLOR_BLUE, first=True)

mini_table(s1, ['Model', 'Front-Only', 'Front+Side', 'Diff'], [
    ['FCNN 7c', '0.158', '0.234', '-0.076'],
    ['FCNN 4c', '0.361', '0.394', '-0.033'],
    ['CNN 7c', '0.137', '0.134', '+0.003'],
], 5.0, 1.18, 4.6, 1.0, cw=[1.5, 1.2, 1.2, 1.0])

# Temuan 12 (bottom-left)
t12 = add_tb(s1, 0.3, 2.35, 4.5, 0.28)
add_para(t12.text_frame, 'Temuan 12: CNN membaik di front-only',
         sz=9, bold=True, color=COLOR_BLUE, first=True)

mini_table(s1, ['Model', 'Front-Only', 'Front+Side', 'Diff'], [
    ['CNN 7c B1', '0.137', '0.133', '+0.005'],
    ['CNN 7c B3', '0.136', '0.119', '+0.017'],
    ['CNN 4c B3', '0.265', '0.238', '+0.027'],
], 0.3, 2.68, 4.5, 1.0, cw=[1.5, 1.2, 1.2, 1.0])

# Temuan 13 (bottom-right)
t13 = add_tb(s1, 5.0, 2.35, 4.6, 0.28)
add_para(t13.text_frame, 'Temuan 13: Konsistensi data > kuantitas data',
         sz=9, bold=True, color=COLOR_BLUE, first=True)

mini_table(s1, ['', 'Front+Side', 'Front-Only'], [
    ['Total sampel', '9,894', '7,091 (-28%)'],
    ['Best Macro F1', '0.407', '0.412 (+1.2%)'],
    ['Best model', 'CNN TL B2', 'Intermediate TL B1'],
], 5.0, 2.68, 4.6, 1.0, cw=[1.8, 1.5, 1.5])

# Conclusion
c = add_tb(s1, 0.3, 3.85, 9.3, 0.35)
c.text_frame.word_wrap = True; set_fill(c, 'E8F0FE')
add_para(c.text_frame, 'Kesimpulan: Menghilangkan side view mengurangi 28% data tapi best model justru sedikit membaik '
         '\u2014 konsistensi sudut kamera lebih penting dari jumlah data',
         sz=9, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER, first=True)

print('  Inserted: Temuan 10-13')

# ═══ SLIDE 93: Temuan 14-15 (Robustness & Leakage) ═══
t = add_tb(s2, 0.3, 0.15, 9.3, 0.6)
add_para(t.text_frame, 'Analisis Hasil \u2014 Temuan Evaluasi Robustness',
         sz=20, bold=True, color=COLOR_TEXT, first=True)

# Temuan 14 (top, full width)
t14 = add_tb(s2, 0.3, 0.85, 9.3, 0.28)
add_para(t14.text_frame, 'Temuan 14: LOSO menunjukkan performa sebenarnya lebih rendah dari single split',
         sz=9, bold=True, color=COLOR_BLUE, first=True)

mini_table(s2, ['Strategi', 'Macro F1', 'Std', 'Keterangan'], [
    ['Single Split', '0.412', '-', 'Fix 5 user test (bisa bias)'],
    ['LOSO (34/37 fold)', '0.370', '0.125', 'Rotasi semua user (gold standard)'],
    ['Selisih', '-0.042', '', 'Single split over-estimate ~10%'],
], 0.3, 1.18, 9.3, 1.05, cw=[2.2, 1.5, 1.0, 4.6])

# Temuan 15 (bottom, full width)
t15 = add_tb(s2, 0.3, 2.4, 9.3, 0.28)
add_para(t15.text_frame, 'Temuan 15: Random Split jauh lebih tinggi \u2014 bukti data leakage',
         sz=9, bold=True, color=COLOR_BLUE, first=True)

mini_table(s2, ['Model', 'User-Based (Single)', 'Random Split', 'Selisih', 'Kenaikan'], [
    ['Intermediate TL', '0.412', '0.586 +/- 0.032', '+0.174', '+42%'],
    ['Late Fusion', '0.394', '0.580 +/- 0.032', '+0.186', '+47%'],
    ['FCNN', '0.361', '0.471 +/- 0.026', '+0.110', '+30%'],
], 0.3, 2.73, 9.3, 1.05, cw=[2.0, 2.0, 2.2, 1.5, 1.6])

# Conclusion
c = add_tb(s2, 0.3, 3.95, 9.3, 0.35)
c.text_frame.word_wrap = True; set_fill(c, 'E8F0FE')
add_para(c.text_frame, 'Kesimpulan: Random Split (0.586) >> Single Split (0.412) > LOSO (0.370) '
         '\u2014 semakin ketat evaluasi, semakin jujur hasilnya. User-based split wajib digunakan.',
         sz=9, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER, first=True)

print('  Inserted: Temuan 14-15')

# ── SAVE ──
prs.save(PPTX_PATH)
print(f'\nSaved! Total: {len(prs.slides)} slides')

for i in range(89, len(prs.slides)):
    slide = prs.slides[i]
    title = next((s.text_frame.text.strip()[:58].replace(chr(10),' ').encode('ascii','replace').decode()
                  for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), '')
    print(f'  Slide {i+1}: {title}')
