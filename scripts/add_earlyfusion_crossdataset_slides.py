"""Insert Early Fusion (SLIDE 30), Cross-Dataset (SLIDE 31), and full
Skema 1/2 table slides (SLIDE 32) into PPT Bimbingan before 'Rencana Eksperimen
Lanjutan' (currently slide 119 = 0-idx 118)."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LB = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_WH = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x20, 0x20, 0x20)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_RED = RGBColor(0xD9, 0x3C, 0x3C)
COLOR_GREEN = RGBColor(0x0F, 0x9D, 0x58)
COLOR_OK = RGBColor(0xE6, 0xF4, 0xEA)
COLOR_WARN = RGBColor(0xFE, 0xF7, 0xE0)


def set_fill(shape, hex_val):
    spPr = shape._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', hex_val)


def tb(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def para(tf, text, sz=10, bold=False, italic=False, color=None,
         align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    if color:
        r.font.color.rgb = color


def cell(c, text, bold=False, bg=None, fg=None, sz=9, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    if fg:
        r.font.color.rgb = fg
    if bg:
        c.fill.solid()
        c.fill.fore_color.rgb = bg


def table(slide, headers, rows, l, t, w, h, cw=None,
          hsz=9, rsz=8, highlight_rows=None):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                    Inches(l), Inches(t),
                                    Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None:
            tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw):
            tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WH,
             sz=hsz, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LB if r % 2 == 0 else COLOR_WH
        if highlight_rows and r in highlight_rows:
            bg = COLOR_OK
        for c, val in enumerate(row):
            cell(tbl.cell(r + 1, c), str(val), bold=(c == 0),
                 bg=bg, fg=COLOR_TEXT, sz=rsz,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


def move_slide_after(prs, moved_idx, target_idx):
    """Move slide at moved_idx to just after target_idx (0-based)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    moved = slides[moved_idx]
    xml_slides.remove(moved)
    slides_after = list(xml_slides)
    target_elem = slides_after[target_idx]
    target_elem.addnext(moved)


def add_title(s, text):
    t = tb(s, 0.3, 0.1, 9.4, 0.45)
    para(t.text_frame, text, sz=14, bold=True,
         color=COLOR_TEXT, first=True)


prs = Presentation(PPTX)
initial = len(prs.slides)
print(f'Initial slides: {initial}')

# Insert target: after slide 118 (0-idx 117 = "Solusi yang Dapat Dikonsultasikan"),
# before slide 119 "Rencana Eksperimen Lanjutan"
INSERT_AFTER_IDX = 117  # 0-based

blank = prs.slide_layouts[6]


# ═══════════ SLIDE 30a: Early Fusion Motivasi + Arsitektur + Hasil 7c ═══════════
s = prs.slides.add_slide(blank)
add_title(s, 'SLIDE 30: Tahap 6 — Early Fusion (Arahan Dosen)')

# Motivasi box
mot = tb(s, 0.3, 0.58, 9.4, 1.0)
mot.text_frame.word_wrap = True
set_fill(mot, 'E8F0FE')
para(mot.text_frame, 'Motivasi & Arsitektur', sz=10, bold=True,
     color=COLOR_BLUE, first=True)
para(mot.text_frame,
     '• Dosen meminta Early Fusion: landmark "ditempel" ke gambar',
     sz=8.5, color=COLOR_TEXT)
para(mot.text_frame,
     '• Pendekatan: HAE-Net (Wu et al., MMM 2020) — landmark → Gaussian heatmap 224×224 → concat sebagai channel ke-4 citra RGB',
     sz=8.5, color=COLOR_TEXT)
para(mot.text_frame,
     '• Fusi terjadi di level input (0%), bukan feature level (Intermediate) atau decision level (Late)',
     sz=8.5, color=COLOR_TEXT)

# Arsitektur detail
arch = tb(s, 0.3, 1.65, 9.4, 0.95)
arch.text_frame.word_wrap = True
set_fill(arch, 'FEF7E0')
para(arch.text_frame, 'Detail Implementasi', sz=10, bold=True,
     color=RGBColor(0xB0, 0x68, 0x00), first=True)
para(arch.text_frame,
     '• Input: Image 224×224×3 + Heatmap 224×224×1 → 224×224×4',
     sz=8.5, color=COLOR_TEXT)
para(arch.text_frame,
     '• Heatmap: 68 landmark → Gaussian blob σ=3px → element-wise max aggregation → single channel [0,1]',
     sz=8.5, color=COLOR_TEXT)
para(arch.text_frame,
     '• TL variant: ResNet18 Conv2d pertama 3→4 channel, weight RGB di-copy, channel ke-4 di-init dari mean RGB',
     sz=8.5, color=COLOR_TEXT)

# 7-class table
para(tb(s, 0.3, 2.75, 4.6, 0.25).text_frame, 'Hasil Early Fusion conf60 — 7-Class',
     sz=10, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Config', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['EF B1', '0.246', '0.794', '0.786', '0.794'],
    ['EF B2', '0.205', '0.520', '0.552', '0.520'],
    ['EF B3', '0.264', '0.680', '0.726', '0.680'],
    ['EF TL B1', '0.253', '0.713', '0.722', '0.713'],
    ['EF TL B2', '0.247', '0.636', '0.663', '0.636'],
    ['EF TL B3', '0.333', '0.753', '0.773', '0.753'],
], 0.3, 3.05, 4.6, 2.35, cw=[1.6, 0.8, 0.8, 0.8, 0.8],
   hsz=8.5, rsz=8, highlight_rows=[5])

# 4-class table
para(tb(s, 5.05, 2.75, 4.6, 0.25).text_frame, 'Hasil Early Fusion conf60 — 4-Class',
     sz=10, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Config', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['EF B1', '0.457', '0.822', '0.816', '0.822'],
    ['EF B2', '0.427', '0.690', '0.728', '0.690'],
    ['EF B3', '0.427', '0.728', '0.752', '0.728'],
    ['EF TL B1', '0.471', '0.770', '0.770', '0.770'],
    ['EF TL B2', '0.424', '0.642', '0.668', '0.642'],
    ['EF TL B3', '0.433', '0.678', '0.709', '0.678'],
], 5.05, 3.05, 4.6, 2.35, cw=[1.6, 0.8, 0.8, 0.8, 0.8],
   hsz=8.5, rsz=8, highlight_rows=[3])

# Best bar
best = tb(s, 0.3, 5.5, 9.4, 0.7)
best.text_frame.word_wrap = True
set_fill(best, 'E6F4EA')
para(best.text_frame,
     'Best Early Fusion: 7c = EF TL B3 (0.333) | 4c = EF TL B1 (0.471)',
     sz=10, bold=True, color=COLOR_GREEN, first=True)
para(best.text_frame,
     '12 config (6 per kelas) × nb 64 — semua metrik Macro/Micro/Weighted F1 + Accuracy.',
     sz=9, color=COLOR_TEXT)


# ═══════════ SLIDE 30b: Perbandingan Fusion + Temuan Early Fusion ═══════════
s = prs.slides.add_slide(blank)
add_title(s, 'SLIDE 30 (lanjutan): Perbandingan Fusion Strategy & Temuan')

# Comparison table
para(tb(s, 0.3, 0.6, 9.4, 0.25).text_frame,
     'Perbandingan 5 Fusion Strategy di Primer conf60 (Best per Arsitektur)',
     sz=10, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Fusion Strategy', 'Best 7c Macro F1', 'Best 4c Macro F1', 'Fusion Point'], [
    ['Single CNN (TL)', '0.273 (B1)', '0.507 (B3)', '— (single modal)'],
    ['Single FCNN', '0.244 (B2)', '0.459 (B2)', '— (single modal)'],
    ['Early Fusion (HAE-Net)', '0.333 (TL B3)', '0.471 (TL B1)', 'Input (0%)'],
    ['Intermediate Fusion', '0.292 (TL B3)', '0.521 (TL B3)', 'Feature (~50%)'],
    ['Late Fusion', '0.301 (TL B1)', '0.567 (TL B3) ★', 'Decision (~95%)'],
], 0.3, 0.9, 9.4, 2.1, cw=[2.5, 1.6, 1.6, 1.8],
   hsz=9, rsz=8.5, highlight_rows=[4])

# Findings
para(tb(s, 0.3, 3.15, 9.4, 0.3).text_frame, 'Temuan (32-35)',
     sz=11, bold=True, color=COLOR_TEXT, first=True)

f32 = tb(s, 0.3, 3.5, 9.4, 0.72)
f32.text_frame.word_wrap = True
set_fill(f32, 'E6F4EA')
para(f32.text_frame, 'Temuan 32: Early Fusion TL B3 tembus best di 7-class (0.333)',
     sz=9, bold=True, color=COLOR_GREEN, first=True)
para(f32.text_frame,
     'Melampaui Intermediate TL B3 (0.292) dan Late Fusion TL B1 (0.301). '
     'Kombinasi heatmap channel + class weights + augmentation + TL bekerja sinergis untuk 7-class.',
     sz=8.5, color=COLOR_TEXT)

f33 = tb(s, 0.3, 4.3, 9.4, 0.65)
f33.text_frame.word_wrap = True
set_fill(f33, 'FEF7E0')
para(f33.text_frame, 'Temuan 33: Early Fusion 4-class underperforms (0.471 vs 0.567 Late Fusion)',
     sz=9, bold=True, color=RGBColor(0xB0, 0x68, 0x00), first=True)
para(f33.text_frame,
     'Heatmap sparse (mostly zeros) kurang informatif ketika class granularity dikurangi; '
     'fusion di level feature/decision lebih optimal untuk kasus ini.',
     sz=8.5, color=COLOR_TEXT)

f34 = tb(s, 0.3, 5.02, 9.4, 0.65)
f34.text_frame.word_wrap = True
set_fill(f34, 'FCE8E6')
para(f34.text_frame, 'Temuan 34: B2 class weights merugikan Early Fusion konsisten',
     sz=9, bold=True, color=COLOR_RED, first=True)
para(f34.text_frame,
     'B2 drop accuracy ~30% vs B1 di semua Early Fusion config. Re-weighting loss merusak learning '
     'saat input 4-channel concat. Intermediate/Late Fusion tetap stabil dengan B2.',
     sz=8.5, color=COLOR_TEXT)

f35 = tb(s, 0.3, 5.74, 9.4, 0.65)
f35.text_frame.word_wrap = True
set_fill(f35, 'E8F0FE')
para(f35.text_frame, 'Temuan 35: Ranking fusion strategy berbeda per granularitas kelas',
     sz=9, bold=True, color=COLOR_BLUE, first=True)
para(f35.text_frame,
     '7-class: Early TL B3 > Late TL B1 > Intermediate TL B3.  '
     '4-class: Late TL B3 > Intermediate TL B3 > Early TL B1.  '
     'Tidak ada single strategy dominan — pemilihan tergantung konteks task.',
     sz=8.5, color=COLOR_TEXT)


# ═══════════ SLIDE 31: Cross-Dataset Best per Source ═══════════
s = prs.slides.add_slide(blank)
add_title(s, 'SLIDE 31: Cross-Dataset Evaluation Lengkap (Skema 2)')

intro = tb(s, 0.3, 0.6, 9.4, 0.55)
intro.text_frame.word_wrap = True
set_fill(intro, 'E8F0FE')
para(intro.text_frame, 'Train di public FER benchmark → Test di Primer conf60 (929 imgs)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
para(intro.text_frame,
     'Melengkapi Skema 2: CK+/JAFFE/RAF-DB/KDEF × 7c/4c × 6 arsitektur. '
     'Sebelumnya hanya CK+ → Primer 7c yang ditabulasi lengkap.',
     sz=8.5, color=COLOR_TEXT)

# Best per source
para(tb(s, 0.3, 1.25, 9.4, 0.25).text_frame,
     'Best Model per Source Dataset → Primer conf60',
     sz=10, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Source', 'Best Model', 'Acc', 'Macro F1', 'Micro F1', 'W-F1'], [
    ['CK+ 7c', 'FCNN', '0.773', '0.194', '0.773', '0.739'],
    ['CK+ 4c', 'CNN', '0.642', '0.396 ★', '0.642', '0.703'],
    ['JAFFE 7c', 'Late Fusion', '0.081', '0.040', '0.081', '0.054'],
    ['JAFFE 4c', 'Intermediate TL', '0.160', '0.093', '0.160', '0.073'],
    ['RAF-DB 7c', 'FCNN', '0.640', '0.183', '0.640', '0.672'],
    ['RAF-DB 4c', 'Intermediate', '0.551', '0.269', '0.551', '0.594'],
    ['KDEF 7c', 'CNN TL', '0.053', '0.038', '0.053', '0.036'],
    ['KDEF 4c', 'CNN', '0.067', '0.079', '0.067', '0.078'],
], 0.3, 1.55, 9.4, 2.8, cw=[1.4, 2.0, 1.1, 1.3, 1.1, 1.1],
   hsz=9, rsz=8.5, highlight_rows=[1, 5])

# Temuan
para(tb(s, 0.3, 4.5, 9.4, 0.3).text_frame, 'Temuan (36-40)',
     sz=11, bold=True, color=COLOR_TEXT, first=True)

bx1 = tb(s, 0.3, 4.82, 4.6, 1.6)
bx1.text_frame.word_wrap = True
set_fill(bx1, 'FCE8E6')
para(bx1.text_frame, 'T36: Semua source gagal generalize',
     sz=9, bold=True, color=COLOR_RED, first=True)
para(bx1.text_frame,
     'Best cross hanya 0.396 (CK+ 4c → Primer dengan CNN B1). '
     'Vs Primer self-training best 0.567 — gap -0.17.',
     sz=8, color=COLOR_TEXT)
para(bx1.text_frame, '\nT39: JAFFE/KDEF → Primer catastrophic',
     sz=9, bold=True, color=COLOR_RED)
para(bx1.text_frame,
     'JAFFE < 10% acc. KDEF 5-7% acc. Lab-posed dataset tidak transferable sama sekali ke natural.',
     sz=8, color=COLOR_TEXT)

bx2 = tb(s, 5.05, 4.82, 4.6, 1.6)
bx2.text_frame.word_wrap = True
set_fill(bx2, 'E6F4EA')
para(bx2.text_frame, 'T38: Landmark > Visual untuk domain shift',
     sz=9, bold=True, color=COLOR_GREEN, first=True)
para(bx2.text_frame,
     'FCNN (landmark-only) best di CK+ 7c dan RAF-DB 7c → Primer. '
     'Landmark geometry lebih tahan domain shift.',
     sz=8, color=COLOR_TEXT)
para(bx2.text_frame, '\nT40: Justifikasi tesis & paper kuat',
     sz=9, bold=True, color=COLOR_GREEN)
para(bx2.text_frame,
     'Transfer learning dari public FER tidak viable untuk natural programming. '
     'Solusi: Primer self-training + conf60 + fusion + aug.',
     sz=8, color=COLOR_TEXT)


# ═══════════ SLIDE 32a: Skema 1 — CK+ & JAFFE (Lengkap dengan Late Fusion TL) ═══════════
s = prs.slides.add_slide(blank)
add_title(s, 'SLIDE 32: Skema 1 Lengkap — CK+ & JAFFE (+ Late Fusion TL)')

intro = tb(s, 0.3, 0.58, 9.4, 0.45)
intro.text_frame.word_wrap = True
set_fill(intro, 'FEF7E0')
para(intro.text_frame,
     'Instruksi dosen: semua 7 model (termasuk Late Fusion TL) × semua metrik Macro/Micro/W-F1/Acc.',
     sz=9, bold=True, color=RGBColor(0xB0, 0x68, 0x00), first=True)

# CK+ 7c
para(tb(s, 0.3, 1.08, 4.6, 0.22).text_frame, 'CK+ 7-Class',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.461', '0.729', '0.659', '0.729'],
    ['FCNN', '0.395', '0.678', '0.614', '0.678'],
    ['Intermediate', '0.316', '0.695', '0.585', '0.695'],
    ['CNN TL', '0.913', '0.949', '0.946', '0.949'],
    ['Intermediate TL', '0.833', '0.881', '0.886', '0.881'],
    ['Late Fusion', '0.498', '0.780', '0.694', '0.780'],
    ['Late Fusion TL', '0.835', '0.881', '0.890', '0.881'],
], 0.3, 1.32, 4.6, 2.5, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[3])

# CK+ 4c
para(tb(s, 5.05, 1.08, 4.6, 0.22).text_frame, 'CK+ 4-Class',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.645', '0.790', '0.776', '0.790'],
    ['FCNN', '0.592', '0.758', '0.740', '0.758'],
    ['Intermediate', '0.567', '0.758', '0.740', '0.758'],
    ['CNN TL', '0.675', '0.903', '0.890', '0.903'],
    ['Intermediate TL', '0.837', '0.903', '0.902', '0.903'],
    ['Late Fusion', '0.592', '0.758', '0.740', '0.758'],
    ['Late Fusion TL', '0.604', '0.806', '0.803', '0.806'],
], 5.05, 1.32, 4.6, 2.5, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[4])

# JAFFE 7c
para(tb(s, 0.3, 3.92, 4.6, 0.22).text_frame, 'JAFFE 7-Class',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.304', '0.450', '0.319', '0.450'],
    ['FCNN', '0.209', '0.250', '0.169', '0.250'],
    ['Intermediate', '0.037', '0.150', '0.039', '0.150'],
    ['CNN TL', '0.464', '0.500', '0.437', '0.500'],
    ['Intermediate TL', '0.447', '0.450', '0.420', '0.450'],
    ['Late Fusion', '0.545', '0.600', '0.522', '0.600'],
    ['Late Fusion TL', '0.146', '0.200', '0.120', '0.200'],
], 0.3, 4.16, 4.6, 2.5, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[5])

# JAFFE 4c
para(tb(s, 5.05, 3.92, 4.6, 0.22).text_frame, 'JAFFE 4-Class',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.177', '0.550', '0.390', '0.550'],
    ['FCNN', '0.438', '0.550', '0.530', '0.550'],
    ['Intermediate', '0.177', '0.550', '0.390', '0.550'],
    ['CNN TL', '0.329', '0.500', '0.476', '0.500'],
    ['Intermediate TL', '0.375', '0.650', '0.558', '0.650'],
    ['Late Fusion', '0.396', '0.650', '0.552', '0.650'],
    ['Late Fusion TL', '0.492', '0.650', '0.615', '0.650'],
], 5.05, 4.16, 4.6, 2.5, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[6])


# ═══════════ SLIDE 32b: Skema 1 — RAF-DB & KDEF (+ Late Fusion TL) ═══════════
s = prs.slides.add_slide(blank)
add_title(s, 'SLIDE 32 (lanjutan): Skema 1 Lengkap — RAF-DB & KDEF')

# RAF-DB 7c
para(tb(s, 0.3, 0.6, 4.6, 0.22).text_frame, 'RAF-DB 7-Class',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.729', '0.815', '0.813', '0.815'],
    ['FCNN', '0.578', '0.714', '0.703', '0.714'],
    ['Intermediate', '0.696', '0.785', '0.783', '0.785'],
    ['CNN TL', '0.741', '0.830', '0.826', '0.830'],
    ['Intermediate TL', '0.744', '0.833', '0.832', '0.833'],
    ['Late Fusion', '0.719', '0.809', '0.805', '0.809'],
    ['Late Fusion TL', '0.735', '0.829', '0.823', '0.829'],
], 0.3, 0.84, 4.6, 2.5, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[4])

# RAF-DB 4c
para(tb(s, 5.05, 0.6, 4.6, 0.22).text_frame, 'RAF-DB 4-Class',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.808', '0.830', '0.830', '0.830'],
    ['FCNN', '0.694', '0.728', '0.729', '0.728'],
    ['Intermediate', '0.792', '0.818', '0.818', '0.818'],
    ['CNN TL', '0.827', '0.845', '0.846', '0.845'],
    ['Intermediate TL', '0.836', '0.853', '0.855', '0.853'],
    ['Late Fusion', '0.819', '0.842', '0.841', '0.842'],
    ['Late Fusion TL', '0.832', '0.849', '0.850', '0.849'],
], 5.05, 0.84, 4.6, 2.5, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[4])

# KDEF 7c
para(tb(s, 0.3, 3.44, 4.6, 0.22).text_frame, 'KDEF 7-Class',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.798', '0.801', '0.798', '0.801'],
    ['FCNN', '0.666', '0.680', '0.663', '0.680'],
    ['Intermediate', '0.671', '0.674', '0.668', '0.674'],
    ['CNN TL', '0.833', '0.831', '0.833', '0.831'],
    ['Intermediate TL', '0.843', '0.843', '0.843', '0.843'],
    ['Late Fusion', '0.776', '0.777', '0.775', '0.777'],
    ['Late Fusion TL', '0.836', '0.834', '0.836', '0.834'],
], 0.3, 3.68, 4.6, 2.5, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[4])

# KDEF 4c
para(tb(s, 5.05, 3.44, 4.6, 0.22).text_frame, 'KDEF 4-Class',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.841', '0.872', '0.872', '0.872'],
    ['FCNN', '0.678', '0.792', '0.766', '0.792'],
    ['Intermediate', '0.776', '0.831', '0.828', '0.831'],
    ['CNN TL', '0.918', '0.929', '0.928', '0.929'],
    ['Intermediate TL', '0.923', '0.929', '0.929', '0.929'],
    ['Late Fusion', '0.859', '0.890', '0.888', '0.890'],
    ['Late Fusion TL', '0.920', '0.932', '0.930', '0.932'],
], 5.05, 3.68, 4.6, 2.5, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[4])

note = tb(s, 0.3, 6.25, 9.4, 0.4)
note.text_frame.word_wrap = True
set_fill(note, 'E6F4EA')
para(note.text_frame,
     'Pola konsisten: Intermediate TL unggul di RAF-DB/KDEF 4c. Late Fusion TL kompetitif (selisih <0.01).',
     sz=9, bold=True, color=COLOR_GREEN, first=True)


# ═══════════ SLIDE 32c: Skema 1 — Primer & Summary ═══════════
s = prs.slides.add_slide(blank)
add_title(s, 'SLIDE 32 (lanjutan): Skema 1 Lengkap — Primer (+ Summary)')

# Primer 7c
para(tb(s, 0.3, 0.6, 4.6, 0.22).text_frame,
     'Primer 7-Class (conf60, B1 baseline)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.270', '0.787', '0.795', '0.787'],
    ['FCNN', '0.261', '0.766', '0.792', '0.766'],
    ['Intermediate', '0.260', '0.799', '0.793', '0.799'],
    ['CNN TL', '0.281', '0.819', '0.814', '0.819'],
    ['Intermediate TL', '0.292', '0.808', '0.819', '0.808'],
    ['Late Fusion', '0.244', '0.778', '0.777', '0.778'],
    ['Late Fusion TL', '0.285', '0.827', '0.828', '0.827'],
], 0.3, 0.84, 4.6, 2.5, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[4])

# Primer 4c
para(tb(s, 5.05, 0.6, 4.6, 0.22).text_frame,
     'Primer 4-Class (conf60, B1 baseline)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.476', '0.803', '0.802', '0.803'],
    ['FCNN', '0.459', '0.763', '0.774', '0.763'],
    ['Intermediate', '0.444', '0.727', '0.753', '0.727'],
    ['CNN TL', '0.378', '0.721', '0.682', '0.721'],
    ['Intermediate TL', '0.482', '0.780', '0.790', '0.780'],
    ['Late Fusion', '0.460', '0.763', '0.779', '0.763'],
    ['Late Fusion TL', '0.472', '0.777', '0.779', '0.777'],
], 5.05, 0.84, 4.6, 2.5, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[4])

note = tb(s, 0.3, 3.45, 9.4, 0.5)
note.text_frame.word_wrap = True
set_fill(note, 'FEF7E0')
para(note.text_frame,
     'Catatan Primer: tabel di atas B1 (baseline) saja untuk apple-to-apple vs benchmark publik.',
     sz=9, bold=True, color=RGBColor(0xB0, 0x68, 0x00), first=True)
para(note.text_frame,
     'Best Primer keseluruhan (dengan B3 + TL + augmentation) = Late Fusion TL 4c B3 = 0.567.',
     sz=9, color=COLOR_TEXT)

# Summary box
para(tb(s, 0.3, 4.05, 9.4, 0.3).text_frame, 'Summary Temuan Skema 1',
     sz=11, bold=True, color=COLOR_TEXT, first=True)

summ = tb(s, 0.3, 4.4, 9.4, 2.5)
summ.text_frame.word_wrap = True
set_fill(summ, 'E8F0FE')
para(summ.text_frame,
     '• Intermediate TL dominan di RAF-DB 7c/4c, KDEF 7c/4c, Primer 7c/4c, CK+ 4c',
     sz=9, color=COLOR_TEXT, first=True)
para(summ.text_frame,
     '• CK+ 7c → CNN TL terbaik (0.913) | JAFFE → Late Fusion atau Late Fusion TL (dataset kecil)',
     sz=9, color=COLOR_TEXT)
para(summ.text_frame,
     '• Late Fusion TL (baru, nb 65): kompetitif di RAF-DB/KDEF/Primer, tapi kolaps di CK+ 7c (0.835 vs 0.913) dan JAFFE 7c (0.146)',
     sz=9, color=COLOR_TEXT)
para(summ.text_frame,
     '• Primer paling challenging: best Primer 4c = 0.482 (Intermediate TL, B1)',
     sz=9, color=COLOR_TEXT)
para(summ.text_frame,
     '• vs benchmark: RAF-DB 4c = 0.836 | KDEF 4c = 0.923 | CK+ 4c = 0.837',
     sz=9, color=COLOR_TEXT)
para(summ.text_frame,
     '• Gap ~0.35-0.45 menegaskan: arsitektur bekerja, Primer memang sulit karena natural + imbalance',
     sz=9, bold=True, color=COLOR_GREEN)


# ═══════════ SLIDE 32d: Skema 2 — CK+ & JAFFE → Primer ═══════════
s = prs.slides.add_slide(blank)
add_title(s, 'SLIDE 32 (lanjutan): Skema 2 Cross-Dataset — CK+ & JAFFE → Primer')

intro = tb(s, 0.3, 0.6, 9.4, 0.4)
intro.text_frame.word_wrap = True
set_fill(intro, 'E8F0FE')
para(intro.text_frame,
     '6 model × 3 metrik per combo. Late Fusion TL belum (butuh checkpoint FCNN+CNN_TL matching).',
     sz=8.5, color=COLOR_TEXT, first=True)

# CK+ → Primer 7c
para(tb(s, 0.3, 1.05, 4.6, 0.22).text_frame, 'CK+ → Primer (7-Class)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.127', '0.719', '0.635', '0.719'],
    ['FCNN', '0.194', '0.773', '0.739', '0.773'],
    ['Intermediate', '0.153', '0.536', '0.593', '0.536'],
    ['CNN TL', '0.163', '0.670', '0.701', '0.670'],
    ['Intermediate TL', '0.103', '0.152', '0.238', '0.152'],
    ['Late Fusion', '0.160', '0.529', '0.580', '0.529'],
], 0.3, 1.3, 4.6, 2.2, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[1])

# CK+ → Primer 4c
para(tb(s, 5.05, 1.05, 4.6, 0.22).text_frame, 'CK+ → Primer (4-Class)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.396', '0.642', '0.703', '0.642'],
    ['FCNN', '0.072', '0.102', '0.093', '0.102'],
    ['Intermediate', '0.186', '0.465', '0.532', '0.465'],
    ['CNN TL', '0.258', '0.386', '0.527', '0.386'],
    ['Intermediate TL', '0.202', '0.271', '0.381', '0.271'],
    ['Late Fusion', '0.245', '0.489', '0.556', '0.489'],
], 5.05, 1.3, 4.6, 2.2, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[0])

# JAFFE → Primer 7c
para(tb(s, 0.3, 3.55, 4.6, 0.22).text_frame, 'JAFFE → Primer (7-Class)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.001', '0.002', '0.000', '0.002'],
    ['FCNN', '0.023', '0.013', '0.018', '0.013'],
    ['Intermediate', '0.017', '0.014', '0.020', '0.014'],
    ['CNN TL', '0.015', '0.054', '0.006', '0.054'],
    ['Intermediate TL', '0.023', '0.028', '0.045', '0.028'],
    ['Late Fusion', '0.040', '0.081', '0.054', '0.081'],
], 0.3, 3.8, 4.6, 2.2, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[5])

# JAFFE → Primer 4c
para(tb(s, 5.05, 3.55, 4.6, 0.22).text_frame, 'JAFFE → Primer (4-Class)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.004', '0.009', '0.000', '0.009'],
    ['FCNN', '0.004', '0.009', '0.000', '0.009'],
    ['Intermediate', '0.005', '0.010', '0.002', '0.010'],
    ['CNN TL', '0.004', '0.009', '0.000', '0.009'],
    ['Intermediate TL', '0.093', '0.160', '0.073', '0.160'],
    ['Late Fusion', '0.007', '0.010', '0.002', '0.010'],
], 5.05, 3.8, 4.6, 2.2, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[4])

note = tb(s, 0.3, 6.1, 9.4, 0.6)
note.text_frame.word_wrap = True
set_fill(note, 'FCE8E6')
para(note.text_frame,
     'CK+ 4c → Primer dengan CNN B1 = 0.396 (satu-satunya combo > 0.30). '
     'JAFFE total collapse di semua combo (<0.10 Macro F1).',
     sz=9, bold=True, color=COLOR_RED, first=True)


# ═══════════ SLIDE 32e: Skema 2 — RAF-DB & KDEF → Primer ═══════════
s = prs.slides.add_slide(blank)
add_title(s, 'SLIDE 32 (lanjutan): Skema 2 Cross-Dataset — RAF-DB & KDEF → Primer')

# RAF-DB → Primer 7c
para(tb(s, 0.3, 0.6, 4.6, 0.22).text_frame, 'RAF-DB → Primer (7-Class)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.076', '0.099', '0.132', '0.099'],
    ['FCNN', '0.183', '0.640', '0.672', '0.640'],
    ['Intermediate', '0.109', '0.237', '0.287', '0.237'],
    ['CNN TL', '0.175', '0.545', '0.611', '0.545'],
    ['Intermediate TL', '0.180', '0.479', '0.562', '0.479'],
    ['Late Fusion', '0.091', '0.166', '0.212', '0.166'],
], 0.3, 0.84, 4.6, 2.2, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[1])

# RAF-DB → Primer 4c
para(tb(s, 5.05, 0.6, 4.6, 0.22).text_frame, 'RAF-DB → Primer (4-Class)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.056', '0.026', '0.026', '0.026'],
    ['FCNN', '0.264', '0.487', '0.553', '0.487'],
    ['Intermediate', '0.269', '0.551', '0.594', '0.551'],
    ['CNN TL', '0.206', '0.256', '0.357', '0.256'],
    ['Intermediate TL', '0.179', '0.306', '0.363', '0.306'],
    ['Late Fusion', '0.170', '0.137', '0.147', '0.137'],
], 5.05, 0.84, 4.6, 2.2, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[2])

# KDEF → Primer 7c
para(tb(s, 0.3, 3.1, 4.6, 0.22).text_frame, 'KDEF → Primer (7-Class)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.024', '0.040', '0.024', '0.040'],
    ['FCNN', '0.007', '0.008', '0.008', '0.008'],
    ['Intermediate', '0.036', '0.020', '0.026', '0.020'],
    ['CNN TL', '0.038', '0.053', '0.036', '0.053'],
    ['Intermediate TL', '0.034', '0.056', '0.023', '0.056'],
    ['Late Fusion', '0.015', '0.008', '0.005', '0.008'],
], 0.3, 3.34, 4.6, 2.2, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[3])

# KDEF → Primer 4c
para(tb(s, 5.05, 3.1, 4.6, 0.22).text_frame, 'KDEF → Primer (4-Class)',
     sz=9.5, bold=True, color=COLOR_BLUE, first=True)
table(s, ['Model', 'Macro', 'Micro', 'W-F1', 'Acc'], [
    ['CNN', '0.079', '0.067', '0.078', '0.067'],
    ['FCNN', '0.004', '0.009', '0.000', '0.009'],
    ['Intermediate', '0.037', '0.020', '0.021', '0.020'],
    ['CNN TL', '0.068', '0.045', '0.023', '0.045'],
    ['Intermediate TL', '0.076', '0.054', '0.030', '0.054'],
    ['Late Fusion', '0.004', '0.009', '0.000', '0.009'],
], 5.05, 3.34, 4.6, 2.2, cw=[1.7, 0.75, 0.75, 0.75, 0.75],
   hsz=8.5, rsz=7.5, highlight_rows=[0])

# Kesimpulan
kes = tb(s, 0.3, 5.6, 9.4, 1.2)
kes.text_frame.word_wrap = True
set_fill(kes, 'E6F4EA')
para(kes.text_frame, 'Kesimpulan Skema 2', sz=10, bold=True,
     color=COLOR_GREEN, first=True)
para(kes.text_frame,
     '• RAF-DB 4c paling robust untuk transfer (best 0.269 via Intermediate) — masih jauh dari Primer self (0.567)',
     sz=8.5, color=COLOR_TEXT)
para(kes.text_frame,
     '• KDEF → Primer catastrophic (<0.08 Macro F1 semua combo) — lab high-quality tidak transfer ke natural',
     sz=8.5, color=COLOR_TEXT)
para(kes.text_frame,
     '• FCNN (landmark) / Intermediate Fusion paling tahan domain shift; CNN murni paling rentan',
     sz=8.5, color=COLOR_TEXT)
para(kes.text_frame,
     '• Justifikasi kuat: Primer self-training + conf60 + fusion + augmentation lebih baik daripada TL dari public',
     sz=8.5, bold=True, color=COLOR_GREEN)


# ─────── Move all new slides to target position ───────
total_new = 7
new_start = initial  # first new slide index (0-based)
for i in range(total_new):
    # Move slide from new_start (always the next new slide to move) to after target
    # After each move, new_start stays the same (since new slides shift)
    move_slide_after(prs, new_start + i, INSERT_AFTER_IDX + i)

out = PPTX
prs.save(out)
print(f'Saved to {out}')
print(f'Total slides: {len(prs.slides)} (added {total_new})')
