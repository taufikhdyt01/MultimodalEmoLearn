"""
Rebuild slides 60-75 in PPT Bimbingan.pptx with proper tables and images.
Also adds image slides for charts from models/ folder.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy
import os

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'
PPTX_OUT = 'd:/MultimodalEmoLearn/docs/PPT_Bimbingan_updated.pptx'
MODELS_DIR = 'd:/MultimodalEmoLearn/models'

# Theme colors (Google Slides blue theme)
COLOR_HEADER_BG = RGBColor(0x1A, 0x73, 0xE8)   # Google Blue
COLOR_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)  # White
COLOR_ROW_ALT = RGBColor(0xE8, 0xF0, 0xFE)     # Light blue
COLOR_ROW_NORMAL = RGBColor(0xFF, 0xFF, 0xFF)  # White
COLOR_TEXT = RGBColor(0x20, 0x20, 0x20)         # Dark gray
COLOR_BEST = RGBColor(0x34, 0xA8, 0x53)         # Green for best results

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Title area (from existing slides)
TITLE_LEFT = Inches(0.33)
TITLE_TOP = Inches(0.47)
TITLE_WIDTH = Inches(9.1)
TITLE_HEIGHT = Inches(0.6)


def remove_slides(prs, start_idx, end_idx):
    """Remove slides from start_idx to end_idx (inclusive, 0-based)."""
    xml_slides = prs.slides._sldIdLst
    R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    # Remove in reverse to preserve indices
    for i in range(end_idx, start_idx - 1, -1):
        sldId_elem = xml_slides[i]
        rId = sldId_elem.get(f'{{{R_NS}}}id')
        xml_slides.remove(sldId_elem)
        # Properly drop the relationship AND its part
        prs.slides.part.drop_rel(rId)


def insert_slide_at(prs, index, layout_idx):
    """Add a new slide at the given index (0-based)."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst
    new_elem = xml_slides[-1]
    xml_slides.remove(new_elem)
    xml_slides.insert(index, new_elem)
    return slide


def set_title(slide, text):
    """Set the title shape text."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.text = text
            return shape
    # fallback: add text box
    txBox = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    txBox.text_frame.text = text
    return txBox


def add_title_box(slide, text):
    """Add a standalone title text box (for BLANK slides)."""
    txBox = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    tf = txBox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x20, 0x20, 0x20)
    return txBox


def add_text_box(slide, text_lines, left, top, width, height, font_size=12):
    """Add a text box with bullet lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text_lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = line
        para.font.size = Pt(font_size)
        para.font.color.rgb = COLOR_TEXT


def style_cell(cell, text, bold=False, bg_color=None, text_color=None, font_size=10, align=PP_ALIGN.LEFT):
    cell.text = text
    para = cell.text_frame.paragraphs[0]
    para.alignment = align
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if text_color:
        run.font.color.rgb = text_color
    if bg_color:
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = bg_color


def add_table(slide, headers, rows, left, top, width, height,
              col_widths=None, header_font_size=10, row_font_size=9,
              highlight_rows=None):
    """
    Add a formatted table to a slide.
    highlight_rows: list of row indices (0-based, excluding header) to highlight green
    """
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # Set column widths
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / total)

    # Header row
    for c, h in enumerate(headers):
        style_cell(tbl.cell(0, c), h,
                   bold=True,
                   bg_color=COLOR_HEADER_BG,
                   text_color=COLOR_HEADER_TEXT,
                   font_size=header_font_size,
                   align=PP_ALIGN.CENTER)

    # Data rows
    for r, row in enumerate(rows):
        is_highlighted = highlight_rows and r in highlight_rows
        bg = COLOR_BEST if is_highlighted else (COLOR_ROW_ALT if r % 2 == 0 else COLOR_ROW_NORMAL)
        text_col = COLOR_HEADER_TEXT if is_highlighted else COLOR_TEXT

        for c, val in enumerate(row):
            bold = is_highlighted
            style_cell(tbl.cell(r + 1, c), str(val),
                       bold=bold,
                       bg_color=bg,
                       text_color=text_col,
                       font_size=row_font_size,
                       align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)

    return tbl


def add_image_slide(prs, index, title_text, img_path, subtitle=None):
    """Add a slide with an image."""
    slide = insert_slide_at(prs, index, 10)  # BLANK layout
    add_title_box(slide, title_text)

    top = Inches(1.1)
    h = Inches(4.2)
    img = slide.shapes.add_picture(img_path, Inches(0.5), top, height=h)
    # Center the image
    img_w = img.width
    img.left = int((SLIDE_W - img_w) / 2)

    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(9), Inches(0.4))
        tf = txBox.text_frame
        tf.text = subtitle
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    return slide


# ───────────���─────────────────────────────────────────────
# LOAD
# ────────────────────────────────────────────────��────────
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Remove old slides 60-75 (indices 59-74)
remove_slides(prs, 59, 74)
print(f'After removal: {len(prs.slides)} slides')

# ─────────��───────────────────────────────────────────────
# INSERT INDEX tracker (current insertion point)
# ────────────────────────────────────────��────────────────
# We insert slides one by one in order, incrementing idx
idx = 59  # Start inserting after slide 59 (Agenda)

# ── SLIDE: Data yang Terkumpul ──────────────────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Data yang Terkumpul')
add_table(slide,
    headers=['Batch', 'Jumlah Mahasiswa', 'Sudut Kamera', 'Periode'],
    rows=[
        ['Batch 1 (lama)', '20 mahasiswa', 'Depan saja', 'April-Mei 2025'],
        ['Batch 2 (baru)', '17 mahasiswa', 'Depan + Samping', 'November 2025'],
        ['Total', '37 mahasiswa', '', ''],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(1.5),
    col_widths=[2, 2.5, 2.5, 2],
)
add_text_box(slide, [
    'Catatan Batch 2: Awalnya direkam 20 mahasiswa, namun 3 file rekaman tidak tersimpan',
    'di hardisk PC perekaman. Setelah ditelusuri, file gagal tersimpan saat proses recording.',
    'Target proposal: 38 mahasiswa  |  Tersedia: 37 mahasiswa'
], Inches(0.5), Inches(2.9), Inches(9), Inches(1.2), font_size=11)

# ── SLIDE: Pipeline Preprocessing ───────��──────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Pipeline Preprocessing')
add_table(slide,
    headers=['Tahap', 'Proses', 'Tool / Output'],
    rows=[
        ['1. Ekstraksi Frame', 'Berdasarkan timestamp emosi di database', 'ffmpeg → frame berlabel'],
        ['2. Face Detection', 'Deteksi & crop wajah 224x224 px', 'MediaPipe (bukan dlib)'],
        ['3. Landmark Extraction', '478 titik → map ke 68 titik standar', '68 titik → 136 fitur (x,y)'],
        ['4. Dataset Preparation', 'Matching label + split train/val/test', 'User-based split'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(2.5),
    col_widths=[2.5, 4, 2.5],
    header_font_size=10, row_font_size=10,
)
add_text_box(slide, [
    'Perubahan dari proposal: dlib diganti MediaPipe karena dlib tidak support side view (batch 2)'
], Inches(0.5), Inches(3.9), Inches(9), Inches(0.4), font_size=10)

# ── SLIDE: Hasil Preprocessing ─────────────���───────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Hasil Preprocessing')
add_table(slide,
    headers=['Tahap', 'Batch 1', 'Batch 2', 'Total'],
    rows=[
        ['Frame diekstrak', '3,849', '6,553', '10,402'],
        ['Face terdeteksi', '3,824', '6,070', '9,894'],
        ['Gagal deteksi', '25', '483', '508'],
        ['Detection rate', '99.4%', '92.6%', '95.1%'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(6), height=Inches(2.2),
    col_widths=[2.5, 1.2, 1.2, 1.1],
    header_font_size=10, row_font_size=10,
)
add_text_box(slide, [
    'Detail Batch 2:',
    '  Front view: 3,275 frame  |  Side view: 3,278 frame',
    'Kegagalan deteksi: wajah menunduk, keluar frame, tertutup tangan',
], Inches(0.5), Inches(3.6), Inches(9), Inches(1.0), font_size=10)

# ── SLIDE: Distribusi Emosi ─────────────────────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Distribusi Emosi Dataset')
add_table(slide,
    headers=['Emosi', 'Jumlah', 'Persentase', 'Keterangan'],
    rows=[
        ['Neutral',   '8,356', '84.5%', 'Kelas mayoritas'],
        ['Happy',       '783',  '7.9%', ''],
        ['Sad',         '576',  '5.8%', ''],
        ['Surprised',    '79',  '0.8%', 'Kelas langka'],
        ['Angry',        '63',  '0.6%', 'Kelas langka'],
        ['Disgusted',    '24',  '0.2%', 'Kelas sangat langka'],
        ['Fearful',      '13',  '0.1%', 'Kelas sangat langka'],
        ['Total',     '9,894', '100%',  ''],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(7), height=Inches(3.5),
    col_widths=[1.8, 1.2, 1.5, 2.5],
    header_font_size=10, row_font_size=10,
)
add_text_box(slide, [
    'Sangat imbalanced: neutral mendominasi 84.5%',
    'Konsisten dengan literatur — mahasiswa fokus coding cenderung netral (Coto et al., 2022)',
], Inches(7.6), Inches(1.5), Inches(2.2), Inches(2), font_size=9)

# ── SLIDE: Penanganan Class Imbalance ───────────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Penanganan Class Imbalance — 3 Skenario Perbandingan')
add_table(slide,
    headers=['Skenario', 'Strategi', 'Dataset Train', 'Metrik'],
    rows=[
        ['B1 — Baseline', 'Tanpa penanganan', '7,064 samples', 'Macro F1'],
        ['B2 — Class Weights', 'Weighted cross-entropy (Cui et al., 2019)', '7,064 samples', 'Macro F1'],
        ['B3 — Weights + Aug', 'Class weights + augmentasi minoritas', '7,519 samples (+455)', 'Macro F1'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(1.5),
    col_widths=[2, 4, 2, 1],
    header_font_size=10, row_font_size=10,
)
add_table(slide,
    headers=['Emosi', 'Original', 'Setelah Augmentasi', 'Ditambah', 'Weight (B2)'],
    rows=[
        ['Angry',     '48',  '150', '+102', '21.3x'],
        ['Fearful',    '8',  '150', '+142', '125.0x'],
        ['Disgusted', '19',  '150', '+131', '52.9x'],
        ['Surprised', '70',  '150',  '+80', '14.7x'],
    ],
    left=Inches(0.5), top=Inches(3.0),
    width=Inches(7), height=Inches(2.2),
    col_widths=[1.5, 1.2, 2, 1.2, 1.1],
    header_font_size=9, row_font_size=9,
)
add_text_box(slide, [
    'Augmentasi: flip horizontal,',
    'rotasi ±15°, brightness.',
    'Hanya train set — val/test',
    'tetap original (evaluasi fair)',
], Inches(7.7), Inches(3.0), Inches(2), Inches(2), font_size=9)

# ── SLIDE: Split Dataset ────────────────────────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Split Dataset — User-Based Split')
add_table(slide,
    headers=['Split', 'Samples', 'Users', 'Persentase'],
    rows=[
        ['Train',       '7,064', '29', '71.4%'],
        ['Validation',  '1,174',  '3', '11.9%'],
        ['Test',        '1,656',  '5', '16.7%'],
        ['Total',       '9,894', '37', '100%'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(5), height=Inches(2.0),
    col_widths=[1.5, 1.2, 1, 1.3],
    header_font_size=10, row_font_size=10,
)
add_table(slide,
    headers=['Emosi', 'Train', 'Val', 'Test'],
    rows=[
        ['Neutral',   '5,678', '1,090', '1,588'],
        ['Happy',       '575',    '22',    '10'],
        ['Sad',         '402',    '48',    '38'],
        ['Angry',        '43',     '2',    '13'],
        ['Fearful',       '7',     '4',     '1'],
        ['Disgusted',    '19',     '2',     '3'],
        ['Surprised',    '39',     '6',     '3'],
    ],
    left=Inches(5.7), top=Inches(1.2),
    width=Inches(4), height=Inches(3.5),
    col_widths=[1.5, 0.9, 0.8, 0.8],
    header_font_size=9, row_font_size=9,
)
add_text_box(slide, [
    'Split by user: mencegah data leaking (model tidak "menghafal" wajah)',
    'Rasio 71/12/17 normal untuk user-based split (standar: 70/10-15/15-20)',
    'Smart distribution: semua 7 emosi terwakili di setiap split',
], Inches(0.5), Inches(3.4), Inches(5), Inches(1.4), font_size=9)

# ── SLIDE: Validasi Ahli ─────────────��──────────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Validasi Ahli Psikologi')
add_table(slide,
    headers=['Opsi', 'Total Sample', 'Strategi', 'Estimasi Waktu'],
    rows=[
        ['A', '1,938', 'Semua non-neutral + 400 neutral', '~8 jam'],
        ['B', '1,067', 'Stratified 10% per kelas (min. 30)', '~4 jam'],
        ['C',   '583', 'Stratified 5% per kelas (min. 30)',  '~2 jam'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(1.6),
    col_widths=[0.5, 1.5, 5, 2],
    header_font_size=10, row_font_size=10,
)
add_table(slide,
    headers=['Fitur Web Tool (Streamlit)', 'Status'],
    rows=[
        ['Tampilkan gambar wajah + label + confidence score', 'Selesai'],
        ['Klik Setuju / pilih emosi yang benar', 'Selesai'],
        ['Multi-validator, progress terpisah per ahli', 'Selesai'],
        ['Auto-hitung Cohen\'s Kappa & inter-rater agreement', 'Selesai'],
        ['Download hasil CSV', 'Selesai'],
        ['Deploy online (Streamlit Cloud) — tanpa instalasi', 'Siap deploy'],
    ],
    left=Inches(0.5), top=Inches(3.0),
    width=Inches(9), height=Inches(2.3),
    col_widths=[7.5, 1.5],
    header_font_size=10, row_font_size=9,
)

# ── SLIDE: Perubahan dari Proposal ─────────────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Perubahan dari Proposal')
add_table(slide,
    headers=['Aspek', 'Di Proposal', 'Implementasi', 'Alasan'],
    rows=[
        ['Face detection', 'dlib', 'MediaPipe', 'dlib gagal deteksi side view (batch 2)'],
        ['Landmark', '68 titik (dlib)', '68 titik (mapped dari MediaPipe)', 'Setara — output tetap 136 fitur'],
        ['Jml mahasiswa', '38', '37', '1 data tidak tersimpan di hardisk'],
        ['Fusion strategy', 'Late + Intermediate', 'Late + Intermediate', 'Tetap sesuai proposal'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(2.5),
    col_widths=[1.8, 1.8, 2.5, 2.9],
    header_font_size=10, row_font_size=10,
)
add_text_box(slide, [
    '(KONSULTASI 3) Apakah perubahan dlib → MediaPipe perlu direvisi di dokumen proposal,',
    'atau cukup dijelaskan di BAB 4 (Implementasi)?',
], Inches(0.5), Inches(4.0), Inches(9), Inches(0.8), font_size=10)

# ── SLIDE: Rencana Selanjutnya ──────────────────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Rencana Selanjutnya')
add_table(slide,
    headers=['No', 'Tahap', 'Status'],
    rows=[
        ['1', 'Pengumpulan data', 'SELESAI'],
        ['2', 'Preprocessing (frame, crop, landmark)', 'SELESAI'],
        ['3', 'Prepare dataset (numpy arrays)', 'SELESAI'],
        ['4', 'Class weights + augmentasi data', 'SELESAI'],
        ['5', 'Training 7-class from scratch (12 eksperimen)', 'SELESAI'],
        ['6', 'Training 4-class from scratch (12 eksperimen)', 'SELESAI'],
        ['7', 'Transfer Learning ResNet18 (24 eksperimen)', 'SELESAI'],
        ['8', 'Tool validasi ahli psikologi (web app)', 'SELESAI — siap deploy'],
        ['9', 'Deploy tool + kirim ke ahli psikologi', 'MENUNGGU KEPUTUSAN'],
        ['10', 'Penulisan BAB 4 & BAB 5', 'SELANJUTNYA'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(4.0),
    col_widths=[0.5, 5.5, 3],
    header_font_size=10, row_font_size=9,
)

# ── SLIDE: Hasil Training 7-Class ──────────────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Hasil Training 7-Class — From Scratch (12 Eksperimen)')
add_table(slide,
    headers=['Rank', 'Model', 'Skenario', 'Accuracy', 'Macro F1', 'Weighted F1'],
    rows=[
        ['1',  'FCNN',               'B1 Baseline', '95.8%', '0.234', '0.952'],
        ['2',  'Late Fusion',         'B1 Baseline', '95.8%', '0.230', '0.951'],
        ['3',  'FCNN',               'B2 Weights',  '89.1%', '0.189', '0.912'],
        ['4',  'Late Fusion',         'B2 Weights',  '89.7%', '0.189', '0.909'],
        ['5',  'Late Fusion',         'B3 Aug',      '92.5%', '0.182', '0.929'],
        ['6',  'FCNN',               'B3 Aug',      '92.3%', '0.182', '0.927'],
        ['7',  'Intermediate Fusion', 'B2 Weights',  '84.5%', '0.140', '0.881'],
        ['9',  'CNN',                'B2 Weights',   '82.8%', '0.134', '0.872'],
        ['10', 'CNN',                'B1 Baseline',  '84.2%', '0.133', '0.881'],
        ['12', 'Intermediate Fusion', 'B1 Baseline', '63.3%', '0.111', '0.744'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(4.0),
    col_widths=[0.5, 2.2, 1.5, 1.2, 1.2, 1.4],
    header_font_size=10, row_font_size=9,
    highlight_rows=[0],
)

# ── SLIDE: Image — 7-class comparison charts ───────────
add_image_slide(prs, idx,
    'Visualisasi Hasil 7-Class — Perbandingan Skenario',
    os.path.join(MODELS_DIR, 'comparison_grouped_bar.png'),
    'Grouped bar chart: Macro F1 per model per skenario (7-class from scratch)'
)
idx += 1

add_image_slide(prs, idx,
    'Heatmap Hasil 7-Class',
    os.path.join(MODELS_DIR, 'comparison_heatmap.png'),
    'Heatmap Macro F1: model vs skenario (7-class from scratch)'
)
idx += 1

# ── SLIDE: Analisis Hasil ──────────────────────��────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Analisis Hasil — Temuan Utama (7-Class)')
add_table(slide,
    headers=['Temuan', 'Deskripsi', 'Implikasi'],
    rows=[
        ['1. FCNN > CNN',
         'FCNN: 0.234 vs CNN: 0.134',
         'Fitur geometrik lebih robust terhadap variasi cahaya & sudut'],
        ['2. Fusion tidak membantu',
         'Late Fusion optimal: 90% FCNN + 10% CNN\nIntermediate lebih buruk',
         'CNN noisy mengganggu fitur landmark yang sudah baik'],
        ['3. Class weights tidak membantu',
         'B1 Baseline terbaik\nWeight 125x pada fearful (8 sample)',
         'Weight ekstrem membuat training tidak stabil'],
        ['4. Kelas minoritas tidak terdeteksi',
         'Happy: 0% recall | Angry: 0% | Fearful: 0%\nHanya Neutral (99%) & Sad (32%)',
         'Test set: 10 happy, 1 fearful, 3 disgusted'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(4.0),
    col_widths=[2, 3.5, 3.5],
    header_font_size=10, row_font_size=9,
)

# ── SLIDE: Hasil Training 4-Class ────────���─────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Hasil Training 4-Class — From Scratch')
add_text_box(slide, [
    '4 kelas: neutral, happy, sad, negative (angry + fearful + disgusted + surprised digabung)',
], Inches(0.5), Inches(1.05), Inches(9), Inches(0.25), font_size=10)
add_table(slide,
    headers=['Rank', 'Model', 'Skenario', 'Accuracy', 'Macro F1', 'Weighted F1'],
    rows=[
        ['1',  'FCNN',               'B3 Aug',      '94.4%', '0.394', '0.943'],
        ['2',  'Late Fusion',         'B3 Aug',      '94.6%', '0.385', '0.943'],
        ['3',  'FCNN',               'B1 Baseline', '95.8%', '0.330', '0.943'],
        ['4',  'Late Fusion',         'B1 Baseline', '95.8%', '0.330', '0.943'],
        ['5',  'FCNN',               'B2 Weights',  '92.9%', '0.327', '0.929'],
        ['7',  'CNN',                'B2 Weights',  '92.8%', '0.296', '0.929'],
        ['9',  'Intermediate Fusion', 'B2 Weights',  '87.3%', '0.258', '0.895'],
        ['12', 'CNN',                'B3 Aug',      '79.2%', '0.238', '0.853'],
    ],
    left=Inches(0.5), top=Inches(1.35),
    width=Inches(9), height=Inches(3.4),
    col_widths=[0.5, 2.2, 1.5, 1.2, 1.2, 1.4],
    header_font_size=10, row_font_size=9,
    highlight_rows=[0],
)

# ── SLIDE: Image — 4-class charts ──────────────────────
add_image_slide(prs, idx,
    'Visualisasi Hasil 4-Class — Perbandingan Skenario',
    os.path.join(MODELS_DIR, '4class', 'comparison_4class_scenarios.png'),
    'Grouped bar chart: Macro F1 per model per skenario (4-class from scratch)'
)
idx += 1

add_image_slide(prs, idx,
    'Perbandingan 7-Class vs 4-Class',
    os.path.join(MODELS_DIR, '4class', 'comparison_4class_vs_7class.png'),
    'Perbandingan Macro F1 terbaik: 7-class vs 4-class per model'
)
idx += 1

# ── SLIDE: Perbandingan 7-Class vs 4-Class ─────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Perbandingan 7-Class vs 4-Class (From Scratch)')
add_table(slide,
    headers=['Model', '7-Class Best (Macro F1)', '4-Class Best (Macro F1)', 'Peningkatan'],
    rows=[
        ['FCNN',               '0.234 (B1)', '0.394 (B3)', '+68%'],
        ['Late Fusion',         '0.230 (B1)', '0.385 (B3)', '+67%'],
        ['CNN',                '0.134 (B2)', '0.296 (B2)', '+121%'],
        ['Intermediate Fusion', '0.140 (B2)', '0.258 (B2)', '+84%'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(1.9),
    col_widths=[2.5, 2.5, 2.5, 1.5],
    header_font_size=10, row_font_size=10,
    highlight_rows=[0],
)
add_table(slide,
    headers=['Temuan', 'Penjelasan'],
    rows=[
        ['5. Penggabungan kelas langka signifikan',
         'Kelas "negative" punya 145 sample vs maks 70 di 7-class'],
        ['6. Augmentasi efektif di 4-class, tidak di 7-class',
         '4-class: 145 sample → augmentasi bermakna. 7-class: fearful 8 sample → tidak bermakna'],
        ['7. FCNN konsisten terbaik di kedua konfigurasi',
         'Fitur geometrik selalu mengungguli fitur penampilan (citra wajah)'],
    ],
    left=Inches(0.5), top=Inches(3.2),
    width=Inches(9), height=Inches(2.0),
    col_widths=[3, 6],
    header_font_size=10, row_font_size=9,
)

# ── SLIDE: Transfer Learning ────────────────────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Hasil Transfer Learning — ResNet18 Pretrained ImageNet')
add_text_box(slide, [
    'Strategi: EmotionCNN (from scratch) → EmotionCNNTransfer (ResNet18)  |  Fine-tune LR: 0.00005'
], Inches(0.5), Inches(1.05), Inches(9), Inches(0.2), font_size=9)
add_table(slide,
    headers=['Model', '7-Class Scratch', '7-Class TL', 'Delta', '4-Class Scratch', '4-Class TL', 'Delta'],
    rows=[
        ['CNN',                 '0.134', '0.177', '+32%', '0.296', '0.407', '+37%'],
        ['FCNN',                '0.234', '0.234',  '—',   '0.394', '0.394',  '—'],
        ['Late Fusion',         '0.230', '0.234',  '+2%', '0.385', '0.442', '+15%'],
        ['Intermediate Fusion', '0.140', '0.232', '+66%', '0.258', '0.376', '+46%'],
    ],
    left=Inches(0.5), top=Inches(1.35),
    width=Inches(9), height=Inches(2.0),
    col_widths=[2, 1.1, 1.0, 0.7, 1.3, 1.0, 0.9],
    header_font_size=9, row_font_size=9,
    highlight_rows=[2],
)
add_text_box(slide, [
    'BEST OVERALL: Late Fusion TL 4-class B2 — Macro F1: 0.442',
    '',
    'Temuan 8: ResNet18 pretrained (ImageNet) meningkatkan CNN secara signifikan (+37% di 4-class)',
    'Temuan 9: FCNN tidak terpengaruh TL — landmark adalah data numerik, tidak butuh pretrained',
    'Temuan 10: Dengan TL, Late Fusion akhirnya menjadi model terbaik (CNN TL tidak lagi noisy)',
], Inches(0.5), Inches(3.5), Inches(9), Inches(1.8), font_size=10)

# ── SLIDE: Image — Final comparison charts ─────────────
add_image_slide(prs, idx,
    'Visualisasi Final — Perbandingan 48 Kombinasi',
    os.path.join(MODELS_DIR, 'final_comparison_bar.png'),
    'Bar chart: Macro F1 semua kombinasi (from scratch vs transfer learning, 7-class vs 4-class)'
)
idx += 1

add_image_slide(prs, idx,
    'Heatmap Final — Semua Eksperimen',
    os.path.join(MODELS_DIR, 'final_comparison_heatmap.png'),
    'Heatmap Macro F1: seluruh 48 kombinasi eksperimen'
)
idx += 1

# ── SLIDE: Perbandingan Lengkap ───────────��─────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Perbandingan Lengkap — 48 Kombinasi Eksperimen')
add_table(slide,
    headers=['Tahap', 'Konfigurasi', 'Model Terbaik', 'Macro F1', 'Catatan'],
    rows=[
        ['Tahap 1', '7-class, From Scratch', 'FCNN B1 Baseline',      '0.234', 'Baseline awal'],
        ['Tahap 2', '4-class, From Scratch', 'FCNN B3 Aug',           '0.394', '+68% dari Tahap 1'],
        ['Tahap 3', 'Transfer Learning',      'Late Fusion TL B2 4-class', '0.442', '+12% dari Tahap 2'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(1.5),
    col_widths=[1, 2, 2.5, 1.2, 2.3],
    header_font_size=10, row_font_size=10,
    highlight_rows=[2],
)
add_table(slide,
    headers=['Rank', 'Model', 'CNN Variant', 'Kelas', 'Skenario', 'Macro F1'],
    rows=[
        ['1', 'Late Fusion',         'Transfer Learning', '4-class', 'B2', '0.442'],
        ['2', 'CNN',                 'Transfer Learning', '4-class', 'B2', '0.407'],
        ['3', 'FCNN',                'Transfer Learning', '4-class', 'B3', '0.394'],
        ['3', 'FCNN',                'From Scratch',      '4-class', 'B3', '0.394'],
        ['5', 'Late Fusion',         'From Scratch',      '4-class', 'B3', '0.385'],
    ],
    left=Inches(0.5), top=Inches(2.9),
    width=Inches(9), height=Inches(2.3),
    col_widths=[0.5, 1.8, 2.2, 1.2, 1.2, 1.1],
    header_font_size=10, row_font_size=9,
    highlight_rows=[0],
)

# ── SLIDE: Diskusi & Konsultasi ────────────���────────────
slide = insert_slide_at(prs, idx, 10); idx += 1
add_title_box(slide, 'Diskusi & Konsultasi')
add_table(slide,
    headers=['RQ', 'Pertanyaan Penelitian', 'Jawaban / Hasil'],
    rows=[
        ['RQ1', 'Performa CNN',
         'From scratch: 0.134 (7c) / 0.296 (4c)\nTransfer Learning: 0.177 (7c) / 0.407 (4c, +37%)'],
        ['RQ2', 'Performa FCNN',
         '0.234 (7-class) / 0.394 (4-class)\nLandmark terbukti lebih robust'],
        ['RQ3', 'Perbandingan Fusion',
         'From scratch: FCNN > Late Fusion > Intermediate\nTL: Late Fusion TL terbaik (0.442)'],
    ],
    left=Inches(0.5), top=Inches(1.2),
    width=Inches(9), height=Inches(2.0),
    col_widths=[0.6, 2.5, 5.9],
    header_font_size=10, row_font_size=9,
)
add_table(slide,
    headers=['No', 'Topik Konsultasi', 'Pertanyaan'],
    rows=[
        ['1', 'Transfer Learning sebagai kontribusi',
         'TL meningkatkan F1: 0.394 → 0.442. Cukup signifikan sebagai kontribusi tesis?'],
        ['2', 'FCNN > Fusion (from scratch)',
         'Bertentangan hipotesis proposal — namun dengan TL, fusion akhirnya terbaik. Narasi tesis?'],
        ['3', 'Macro F1 0.442',
         'Acceptable untuk tesis S2? Dataset naturalistik, nyata, sangat imbalanced'],
        ['4', 'Validasi ahli psikologi',
         'Opsi A/B/C? 1 atau 2 ahli? Perlu honorarium?'],
    ],
    left=Inches(0.5), top=Inches(3.3),
    width=Inches(9), height=Inches(2.0),
    col_widths=[0.4, 2.5, 6.1],
    header_font_size=10, row_font_size=9,
)

# ─────────��───────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────
prs.save(PPTX_OUT)
print(f'Saved to: {PPTX_OUT}')
print(f'Total slides: {len(prs.slides)}')

# Verify
prs2 = Presentation(PPTX_OUT)
print(f'\nSlide list (from slide 58):')
for i in range(57, len(prs2.slides)):
    slide = prs2.slides[i]
    title = ''
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            title = shape.text_frame.text.strip()[:55].replace('\n', ' ')
            break
    print(f'  Slide {i+1:2d}: {title}')
