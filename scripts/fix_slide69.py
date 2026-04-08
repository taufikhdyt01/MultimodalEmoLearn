"""Replace slide 69 (Tahap 1) with bullet-point format instead of checkmark table."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'
COLOR_HEADER_BG  = RGBColor(0x1A, 0x73, 0xE8)
COLOR_HEADER_TEXT= RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ROW_ALT    = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_ROW_NORMAL = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT       = RGBColor(0x20, 0x20, 0x20)
COLOR_BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)


def add_title(slide, text, font_size=22):
    tb = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(9.3), Inches(0.7))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT


def add_colored_box(slide, lines, left, top, width, height, bg_color, font_size=10):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    spPr = tb._element.spPr
    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    r, g, b = bg_color[0], bg_color[1], bg_color[2]
    srgbClr.set('val', '%02X%02X%02X' % (r, g, b))
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(font_size)
        para.font.color.rgb = COLOR_TEXT
        if i == 0:
            para.font.bold = True


def style_cell(cell, text, bold=False, bg=None, fg=None, size=9, align=PP_ALIGN.LEFT):
    cell.text = text
    para = cell.text_frame.paragraphs[0]
    para.alignment = align
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if fg:
        run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def add_table(slide, headers, rows, left, top, width, height,
              col_widths=None, h_size=10, r_size=9):
    tbl = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, width, height).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / total)
    for c, h in enumerate(headers):
        style_cell(tbl.cell(0, c), h, bold=True,
                   bg=COLOR_HEADER_BG, fg=COLOR_HEADER_TEXT,
                   size=h_size, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_ROW_ALT if r % 2 == 0 else COLOR_ROW_NORMAL
        for c, val in enumerate(row):
            style_cell(tbl.cell(r + 1, c), str(val),
                       bg=bg, fg=COLOR_TEXT, size=r_size,
                       align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)


prs = Presentation(PPTX_PATH)

# Slide 69 = index 68
slide = prs.slides[68]
print('Clearing slide 69:', next(
    (s.text_frame.text.strip()[:50] for s in slide.shapes if s.has_text_frame
     and s.text_frame.text.strip()), ''))

# Clear all shapes
sp_tree = slide.shapes._spTree
for child in list(sp_tree):
    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
    if tag not in ('bg', 'bgPr'):
        sp_tree.remove(child)

# Rebuild with bullet format
add_title(slide, 'Rancangan Eksperimen — Tahap 1: 7-Class From Scratch')

# Left column: 4 Model
add_colored_box(slide, [
    '4 Model yang Dibandingkan:',
    'CNN                   -- fitur citra wajah 224x224 px',
    'FCNN                  -- fitur landmark geometrik (136 koordinat)',
    'Late Fusion           -- CNN + FCNN digabung di level keputusan',
    'Intermediate Fusion   -- CNN + FCNN digabung di level fitur',
], Inches(0.4), Inches(1.0), Inches(4.5), Inches(2.1),
   bg_color=COLOR_BLUE_LIGHT, font_size=10)

# Right column: 3 Skenario
add_colored_box(slide, [
    '3 Skenario Penanganan Class Imbalance:',
    'B1 = Baseline: training biasa, tanpa penanganan',
    'B2 = Class Weights (Cui et al., 2019):',
    '     penalty loss lebih besar untuk kelas langka',
    'B3 = Class Weights + Augmentasi:',
    '     flip, rotasi +-15 derajat, brightness',
    '     hanya untuk kelas < 150 sample',
], Inches(5.1), Inches(1.0), Inches(4.5), Inches(2.1),
   bg_color=COLOR_BLUE_LIGHT, font_size=10)

# Summary table
add_table(slide,
    headers=['Konfigurasi', 'Jumlah Eksperimen', 'Kelas', 'GPU', 'Metrik'],
    rows=[
        ['4 model x 3 skenario', '12 eksperimen',
         '7 kelas (neutral, happy, sad, angry, fearful, disgusted, surprised)',
         'NVIDIA T4', 'Macro F1'],
    ],
    left=Inches(0.4), top=Inches(3.25),
    width=Inches(9.2), height=Inches(0.85),
    col_widths=[2.0, 1.8, 3.6, 1.0, 0.8],
    h_size=10, r_size=10,
)

prs.save(PPTX_PATH)
print('Saved!')

# Verify
prs2 = Presentation(PPTX_PATH)
slide = prs2.slides[68]
print('Slide 69 shapes:', len(list(slide.shapes)))
for shape in slide.shapes:
    if shape.has_text_frame:
        print(' ', repr(shape.text_frame.text[:60].replace('\n', ' ')))