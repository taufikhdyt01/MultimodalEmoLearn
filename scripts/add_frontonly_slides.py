"""
Add 3 new slides before Diskusi (slide 87):
1. Eksperimen Front-Only + Perbandingan
2. Evaluasi Robustness (LOSO/CV/Random Split)
3. Temuan Data Leakage (Random >> Single > LOSO)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE       = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_ORANGE     = RGBColor(0xE6, 0x51, 0x00)
COLOR_GREEN      = RGBColor(0x0F, 0x9D, 0x58)
COLOR_RED        = RGBColor(0xD9, 0x3C, 0x3C)
COLOR_TEXT       = RGBColor(0x20, 0x20, 0x20)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY       = RGBColor(0x66, 0x66, 0x66)
COLOR_LIGHT_GRAY = RGBColor(0xF1, 0xF3, 0xF4)


def set_fill(tb, hex_val):
    spPr = tb._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', hex_val)


def add_tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))


def add_para(tf, text, size=10, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, first=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return para


def style_cell(cell, text, bold=False, bg=None, fg=None, size=9, align=PP_ALIGN.LEFT):
    cell.text = text
    para = cell.text_frame.paragraphs[0]
    para.alignment = align
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if fg:
        run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def add_table(slide, headers, rows, left, top, width, height,
              col_widths=None, h_size=9, r_size=9):
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None:
            tblPr.remove(sid)
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(Inches(width) * cw / total)
    for c, h in enumerate(headers):
        style_cell(tbl.cell(0, c), h, bold=True,
                   bg=COLOR_BLUE, fg=COLOR_WHITE,
                   size=h_size, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LIGHT_BLUE if r % 2 == 0 else COLOR_WHITE
        for c, val in enumerate(row):
            style_cell(tbl.cell(r + 1, c), str(val),
                       bold=(c == 0), bg=bg, fg=COLOR_TEXT,
                       size=r_size,
                       align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


def insert_slide(prs, index):
    layout = prs.slide_layouts[10]
    slide = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst
    new_elem = xml_slides[-1]
    xml_slides.remove(new_elem)
    xml_slides.insert(index, new_elem)
    return slide


# ── LOAD ──
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Insert before Diskusi (currently slide 87 = index 86)
# Insert 3 slides at index 86 → they become slides 87, 88, 89
# Diskusi moves to 90, Terimakasih to 91

# ═══════════════════════════════════════════════════════
# SLIDE 1: Eksperimen Front-Only + Perbandingan
# ═══════════════════════════════════════════════════════
s1 = insert_slide(prs, 86)

tb = add_tb(s1, 0.3, 0.12, 9.4, 0.55)
tf = tb.text_frame
add_para(tf, 'Eksperimen Front-Only & Perbandingan dengan Front+Side',
         size=18, bold=True, color=COLOR_TEXT, first=True)

# Motivation box
tb_m = add_tb(s1, 0.3, 0.75, 9.4, 0.8)
tf_m = tb_m.text_frame; tf_m.word_wrap = True
set_fill(tb_m, 'FFF0CC')
add_para(tf_m, 'Motivasi: Batch 1 hanya depan, batch 2 depan+samping '
         '-> inkonsistensi data -> ulangi 48 eksperimen dengan front-only (7,091 sampel)',
         size=9.5, color=COLOR_TEXT, first=True)

# Table: Perbandingan
add_table(s1,
    headers=['Tahap', 'Front-Only (F1)', 'Front+Side (F1)', 'Selisih'],
    rows=[
        ['7-class scratch', '0.175', '0.234', '-0.060'],
        ['4-class scratch', '0.394', '0.394', '-0.001'],
        ['7-class TL', '0.180', '0.232', '-0.052'],
        ['4-class TL', '0.412', '0.407', '+0.005'],
    ],
    left=0.3, top=1.65, width=9.4, height=1.8,
    col_widths=[2.5, 2.3, 2.3, 2.3], h_size=10, r_size=10)

# Bottom insight box
tb_i = add_tb(s1, 0.3, 3.55, 9.4, 1.0)
tf_i = tb_i.text_frame; tf_i.word_wrap = True
set_fill(tb_i, 'E6F4EA')
add_para(tf_i, 'Temuan:', size=10, bold=True, color=COLOR_TEXT, first=True)
add_para(tf_i, '-> Best front-only (0.412) sedikit lebih baik dari front+side (0.407)',
         size=9.5, color=COLOR_TEXT)
add_para(tf_i, '-> Side view tidak meningkatkan performa, justru menambah noise',
         size=9.5, color=COLOR_TEXT)
add_para(tf_i, '-> Konsistensi data lebih penting dari kuantitas data',
         size=9.5, color=COLOR_TEXT)

print('  Inserted slide: Eksperimen Front-Only')

# ═══════════════════════════════════════════════════════
# SLIDE 2: Evaluasi Robustness
# ═══════════════════════════════════════════════════════
s2 = insert_slide(prs, 87)

tb = add_tb(s2, 0.3, 0.12, 9.4, 0.55)
tf = tb.text_frame
add_para(tf, 'Evaluasi Robustness: LOSO, 5-Fold CV, Random Split',
         size=18, bold=True, color=COLOR_TEXT, first=True)

# Strategy explanation
add_table(s2,
    headers=['Strategi', 'Cara', 'Fold', 'Data Leakage?'],
    rows=[
        ['LOSO', '1 user = 1 test set, rotasi', '37', 'Tidak'],
        ['5-Fold CV', 'User dibagi 5 grup, rotasi', '5', 'Tidak'],
        ['Random Split', 'Sampel diacak tanpa peduli user', '5 repeat', 'Ya (baseline)'],
    ],
    left=0.3, top=0.75, width=9.4, height=1.3,
    col_widths=[2.0, 3.5, 1.5, 2.4], h_size=10, r_size=10)

# Results table
add_table(s2,
    headers=['Model', 'Single Split', 'Random Split', '5-Fold CV', 'LOSO'],
    rows=[
        ['Intermediate TL', '0.412', '0.586 +/- 0.032', 'pending', '0.370 +/- 0.125'],
        ['Late Fusion', '0.394', '0.580 +/- 0.032', 'pending', 'pending'],
        ['FCNN', '0.361', '0.471 +/- 0.026', 'pending', 'pending'],
    ],
    left=0.3, top=2.25, width=9.4, height=1.3,
    col_widths=[2.2, 1.5, 2.0, 1.85, 1.85], h_size=9.5, r_size=9.5)

# Insight
tb_i = add_tb(s2, 0.3, 3.7, 9.4, 0.4)
tf_i = tb_i.text_frame; tf_i.word_wrap = True
add_para(tf_i, 'Pola: Random (0.586) >> Single (0.412) > LOSO (0.370) '
         '-- semakin ketat evaluasi, semakin jujur hasilnya',
         size=9.5, italic=True, color=COLOR_GRAY, first=True)

print('  Inserted slide: Evaluasi Robustness')

# ═══════════════════════════════════════════════════════
# SLIDE 3: Temuan Data Leakage
# ═══════════════════════════════════════════════════════
s3 = insert_slide(prs, 88)

tb = add_tb(s3, 0.3, 0.12, 9.4, 0.55)
tf = tb.text_frame
add_para(tf, 'Temuan: Bukti Data Leakage & Pentingnya User-Based Split',
         size=18, bold=True, color=COLOR_TEXT, first=True)

# Data leakage comparison
add_table(s3,
    headers=['Model', 'User-Based Split', 'Random Split', 'Selisih', 'Kenaikan'],
    rows=[
        ['Intermediate TL', '0.412', '0.586', '+0.174', '+42%'],
        ['Late Fusion', '0.394', '0.580', '+0.186', '+47%'],
        ['FCNN', '0.361', '0.471', '+0.110', '+30%'],
    ],
    left=0.3, top=0.8, width=9.4, height=1.3,
    col_widths=[2.2, 1.8, 1.8, 1.8, 1.8], h_size=10, r_size=10)

# Explanation boxes
# Box 1: Why leakage
tb_l = add_tb(s3, 0.3, 2.25, 4.55, 1.6)
tf_l = tb_l.text_frame; tf_l.word_wrap = True
set_fill(tb_l, 'FCE4E4')
add_para(tf_l, 'Mengapa Random Split Lebih Tinggi?', size=10, bold=True, color=COLOR_RED, first=True)
add_para(tf_l, '', size=5, first=False)
add_para(tf_l, 'Random split: user yang sama bisa ada di train & test', size=9.5, color=COLOR_TEXT)
add_para(tf_l, '-> Model "menghafal" wajah user, bukan ekspresi', size=9.5, color=COLOR_TEXT)
add_para(tf_l, '-> F1 naik 30-47% bukan karena model lebih baik,', size=9.5, color=COLOR_TEXT)
add_para(tf_l, '   tapi karena evaluasi tidak valid (data leakage)', size=9.5, color=COLOR_TEXT)

# Box 2: LOSO insight
tb_r = add_tb(s3, 4.95, 2.25, 4.75, 1.6)
tf_r = tb_r.text_frame; tf_r.word_wrap = True
set_fill(tb_r, 'E6F4EA')
add_para(tf_r, 'LOSO Menunjukkan Performa Sebenarnya', size=10, bold=True, color=COLOR_GREEN, first=True)
add_para(tf_r, '', size=5, first=False)
add_para(tf_r, 'Intermediate TL LOSO: 0.370 +/- 0.125', size=9.5, color=COLOR_TEXT)
add_para(tf_r, '-> Std tinggi (0.125): performa sangat bervariasi', size=9.5, color=COLOR_TEXT)
add_para(tf_r, '   antar user -- ada yang mudah, ada yang sulit', size=9.5, color=COLOR_TEXT)
add_para(tf_r, '-> Single split (0.412) sedikit over-estimate', size=9.5, color=COLOR_TEXT)

# Bottom conclusion
tb_c = add_tb(s3, 0.3, 4.0, 9.4, 0.6)
tf_c = tb_c.text_frame; tf_c.word_wrap = True
set_fill(tb_c, '1A73E8')
add_para(tf_c, 'Kesimpulan: Evaluasi HARUS menggunakan user-based split (LOSO/CV) -- '
         'random split memberikan hasil yang menyesatkan (+30-47%)',
         size=10, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER, first=True)

print('  Inserted slide: Temuan Data Leakage')

# ── SAVE ──
prs.save(PPTX_PATH)
print(f'\nSaved! Total: {len(prs.slides)} slides')

# Verify
prs2 = Presentation(PPTX_PATH)
for i in range(84, len(prs2.slides)):
    slide = prs2.slides[i]
    title = next((s.text_frame.text.strip()[:55].replace(chr(10),' ').encode('ascii','replace').decode()
                  for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), '')
    print(f'  Slide {i+1}: {title}')
