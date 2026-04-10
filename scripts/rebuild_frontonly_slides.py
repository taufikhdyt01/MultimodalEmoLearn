"""
Rebuild slides 87-89 with complete front-only results:
- Slide 87: Rancangan Eksperimen Front-Only (motivasi)
- Slide 88: Hasil Training 7-Class Front-Only (12 eksperimen)
- Slide 89: Hasil Training 4-Class Front-Only (12 eksperimen)
- Slide 90: Hasil Transfer Learning Front-Only (tabel 7+4 class)
- Slide 91: Perbandingan Front-Only vs Front+Side (best per tahap)
- Slide 92: Evaluasi Robustness (LOSO/CV/Random Split)
- Slide 93: Temuan Data Leakage
Then Diskusi + Terimakasih follow after.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE       = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_GREEN      = RGBColor(0x0F, 0x9D, 0x58)
COLOR_RED        = RGBColor(0xD9, 0x3C, 0x3C)
COLOR_TEXT       = RGBColor(0x20, 0x20, 0x20)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY       = RGBColor(0x66, 0x66, 0x66)


def set_fill(tb, hex_val):
    spPr = tb._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', hex_val)


def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def para(tf, text, sz=10, bold=False, italic=False, color=None, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = color
    return p


def cell(c, text, bold=False, bg=None, fg=None, sz=9, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]; p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    if fg: r.font.color.rgb = fg
    if bg: c.fill.solid(); c.fill.fore_color.rgb = bg


def table(slide, headers, rows, l, t, w, h, cw=None, hsz=9, rsz=9, highlight_best=None):
    shape = slide.shapes.add_table(len(rows)+1, len(headers), Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None: tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw): tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WHITE, sz=hsz, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LIGHT_BLUE if r % 2 == 0 else COLOR_WHITE
        is_best = (highlight_best is not None and r == highlight_best)
        for c, val in enumerate(row):
            fg = COLOR_TEXT
            b = (c == 0)
            if is_best:
                b = True
                bg_r = RGBColor(0xC6, 0xEF, 0xCE)  # light green for best row
            else:
                bg_r = bg
            cell(tbl.cell(r+1, c), str(val), bold=b, bg=bg_r, fg=fg, sz=rsz,
                 align=PP_ALIGN.LEFT if c <= 1 else PP_ALIGN.CENTER)


def insert(prs, index):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    xml_slides = prs.slides._sldIdLst
    el = xml_slides[-1]; xml_slides.remove(el); xml_slides.insert(index, el)
    return slide


def clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            sp_tree.remove(child)


# ── LOAD ──
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Delete old slides 87-89 (indices 86-88) — clear and reuse
for i in [86, 87, 88]:
    clear_slide(prs.slides[i])
print('Cleared slides 87-89')

# We need 7 slides total (87-93), currently have 3 cleared + need 4 more
# Insert 4 new slides at index 89 (after cleared slides)
for _ in range(4):
    insert(prs, 89)
print('Inserted 4 new slides')

# Now slides 87-93 are available (indices 86-92)
# 87 = Rancangan Front-Only
# 88 = Hasil 7-Class Front-Only
# 89 = Hasil 4-Class Front-Only
# 90 = Hasil TL Front-Only
# 91 = Perbandingan Front-Only vs Front+Side
# 92 = Evaluasi Robustness
# 93 = Temuan Data Leakage
# 94 = Diskusi
# 95 = Terimakasih

# ═══ SLIDE 87: Rancangan Eksperimen Front-Only ═══
s = prs.slides[86]
t = tb(s, 0.3, 0.12, 9.4, 0.6)
para(t.text_frame, 'Rancangan Eksperimen \u2014 Tahap 4: Front-Only', sz=20, bold=True, color=COLOR_TEXT, first=True)

# Motivation box
m = tb(s, 0.3, 0.85, 9.4, 1.3)
m.text_frame.word_wrap = True; set_fill(m, 'FFF0CC')
para(m.text_frame, 'Motivasi:', sz=10, bold=True, color=COLOR_TEXT, first=True)
para(m.text_frame, '-> Batch 1 hanya memiliki sudut kamera depan (front)', sz=9.5, color=COLOR_TEXT)
para(m.text_frame, '-> Batch 2 memiliki depan + samping (side)', sz=9.5, color=COLOR_TEXT)
para(m.text_frame, '-> Inkonsistensi ini bisa bias hasil eksperimen', sz=9.5, color=COLOR_TEXT)
para(m.text_frame, '-> Solusi: ulangi seluruh eksperimen dengan front-only', sz=9.5, color=COLOR_TEXT)

# Dataset table
table(s, ['', 'Front+Side (sebelumnya)', 'Front-Only'], [
    ['Batch 1', '3,824 (front)', '3,824 (front)'],
    ['Batch 2', '~6,070 (front+side)', '~3,267 (front saja)'],
    ['Total', '9,894', '7,091'],
    ['Train / Val / Test', '7,064 / 1,174 / 1,656', '5,348 / 707 / 1,036'],
], 0.3, 2.3, 9.4, 1.7, cw=[2.5, 3.5, 3.4], hsz=10, rsz=10)

b = tb(s, 0.3, 4.1, 9.4, 0.3)
para(b.text_frame, '48 eksperimen front-only (4 model x 3 skenario x 2 kelas x 2 CNN variant) | User split identik',
     sz=10, color=COLOR_TEXT, first=True)
print('  87: Rancangan Front-Only')

# ═══ SLIDE 88: Hasil 7-Class Front-Only ═══
s = prs.slides[87]
t = tb(s, 0.3, 0.12, 9.4, 0.6)
para(t.text_frame, 'Hasil Training 7-Class Front-Only (12 Eksperimen)', sz=20, bold=True, color=COLOR_TEXT, first=True)

rows_7c = [
    ['1', 'Late Fusion', 'B3 Aug', '93.3%', '0.175', '0.919'],
    ['2', 'Late Fusion', 'B1 Baseline', '93.1%', '0.171', '0.917'],
    ['3', 'FCNN', 'B1 Baseline', '90.9%', '0.158', '0.906'],
    ['4', 'FCNN', 'B3 Aug', '88.3%', '0.157', '0.892'],
    ['5', 'Late Fusion', 'B2 Weights', '94.1%', '0.151', '0.921'],
    ['6', 'FCNN', 'B2 Weights', '91.7%', '0.148', '0.913'],
    ['7', 'CNN', 'B1 Baseline', '91.8%', '0.137', '0.911'],
    ['8', 'Intermediate', 'B1 Baseline', '91.4%', '0.137', '0.909'],
    ['9', 'CNN', 'B2 Weights', '90.5%', '0.137', '0.905'],
    ['10', 'CNN', 'B3 Aug', '89.3%', '0.136', '0.902'],
    ['11', 'Intermediate', 'B2 Weights', '84.9%', '0.136', '0.885'],
    ['12', 'Intermediate', 'B3 Aug', '85.7%', '0.135', '0.876'],
]
table(s, ['Rank', 'Model', 'Skenario', 'Accuracy', 'Macro F1', 'W-F1'],
      rows_7c, 0.3, 0.8, 9.4, 3.8, cw=[0.6, 2.0, 1.5, 1.3, 1.5, 1.5],
      hsz=10, rsz=9.5, highlight_best=0)

b = tb(s, 0.3, 4.7, 9.4, 0.3)
para(b.text_frame, 'Best: Late Fusion B3 (Macro F1: 0.175) | FCNN dan Late Fusion mendominasi',
     sz=10, bold=True, color=COLOR_TEXT, first=True)
print('  88: Hasil 7-Class')

# ═══ SLIDE 89: Hasil 4-Class Front-Only ═══
s = prs.slides[88]
t = tb(s, 0.3, 0.12, 9.4, 0.6)
para(t.text_frame, 'Hasil Training 4-Class Front-Only (12 Eksperimen)', sz=20, bold=True, color=COLOR_TEXT, first=True)

rows_4c = [
    ['1', 'Late Fusion', 'B3 Aug', '92.3%', '0.394', '0.924'],
    ['2', 'FCNN', 'B3 Aug', '89.4%', '0.361', '0.909'],
    ['3', 'FCNN', 'B1 Baseline', '89.7%', '0.307', '0.901'],
    ['4', 'Late Fusion', 'B1 Baseline', '89.7%', '0.307', '0.901'],
    ['5', 'CNN', 'B2 Weights', '90.1%', '0.240', '0.908'],
    ['6', 'Late Fusion', 'B2 Weights', '91.9%', '0.276', '0.914'],
    ['7', 'Intermediate', 'B2 Weights', '87.5%', '0.269', '0.892'],
    ['8', 'CNN', 'B3 Aug', '90.2%', '0.265', '0.906'],
    ['9', 'Intermediate', 'B1 Baseline', '94.3%', '0.243', '0.919'],
    ['10', 'CNN', 'B1 Baseline', '89.3%', '0.238', '0.901'],
    ['11', 'Intermediate', 'B3 Aug', '90.7%', '0.239', '0.903'],
    ['12', 'FCNN', 'B2 Weights', '84.2%', '0.245', '0.872'],
]
table(s, ['Rank', 'Model', 'Skenario', 'Accuracy', 'Macro F1', 'W-F1'],
      rows_4c, 0.3, 0.8, 9.4, 3.8, cw=[0.6, 2.0, 1.5, 1.3, 1.5, 1.5],
      hsz=10, rsz=9.5, highlight_best=0)

b = tb(s, 0.3, 4.7, 9.4, 0.3)
para(b.text_frame, 'Best: Late Fusion B3 (Macro F1: 0.394) | +125% dari 7-class (0.175)',
     sz=10, bold=True, color=COLOR_TEXT, first=True)
print('  89: Hasil 4-Class')

# ═══ SLIDE 90: Hasil TL Front-Only ═══
s = prs.slides[89]
t = tb(s, 0.3, 0.12, 9.4, 0.6)
para(t.text_frame, 'Hasil Transfer Learning Front-Only (ResNet18)', sz=20, bold=True, color=COLOR_TEXT, first=True)

# 7-class TL
t1 = tb(s, 0.3, 0.75, 9.4, 0.3)
para(t1.text_frame, '7-Class TL: Best per Model', sz=11, bold=True, color=COLOR_TEXT, first=True)

table(s, ['Model', 'From Scratch (Best)', 'Transfer Learning (Best)', 'Peningkatan'], [
    ['CNN', '0.137 (B1)', '0.154 (B1)', '+12%'],
    ['FCNN', '0.158 (B1)', '0.158 (B1)', '- (sama)'],
    ['Late Fusion', '0.175 (B3)', '0.167 (B1)', '-5%'],
    ['Intermediate', '0.137 (B1)', '0.180 (B3)', '+31%'],
], 0.3, 1.1, 9.4, 1.5, cw=[2.2, 2.4, 2.4, 2.4], hsz=9.5, rsz=9.5, highlight_best=3)

# 4-class TL
t2 = tb(s, 0.3, 2.75, 9.4, 0.3)
para(t2.text_frame, '4-Class TL: Best per Model', sz=11, bold=True, color=COLOR_TEXT, first=True)

table(s, ['Model', 'From Scratch (Best)', 'Transfer Learning (Best)', 'Peningkatan'], [
    ['CNN', '0.265 (B3)', '0.274 (B1)', '+3%'],
    ['FCNN', '0.361 (B3)', '0.361 (B3)', '- (sama)'],
    ['Late Fusion', '0.394 (B3)', '0.372 (B3)', '-6%'],
    ['Intermediate', '0.269 (B2)', '0.412 (B1)', '+53%'],
], 0.3, 3.1, 9.4, 1.5, cw=[2.2, 2.4, 2.4, 2.4], hsz=9.5, rsz=9.5, highlight_best=3)

b = tb(s, 0.3, 4.7, 9.4, 0.3)
para(b.text_frame, 'Best Overall: Intermediate Fusion TL 4-class B1 (Macro F1: 0.412)',
     sz=10, bold=True, color=COLOR_TEXT, first=True)
print('  90: Hasil TL')

# ═══ SLIDE 91: Perbandingan Front-Only vs Front+Side ═══
s = prs.slides[90]
t = tb(s, 0.3, 0.12, 9.4, 0.6)
para(t.text_frame, 'Perbandingan Front-Only vs Front+Side', sz=20, bold=True, color=COLOR_TEXT, first=True)

table(s, ['Tahap', 'Front-Only (F1)', 'Front+Side (F1)', 'Selisih'], [
    ['7-class scratch', '0.175 (Late Fusion B3)', '0.234 (FCNN B1)', '-0.060'],
    ['4-class scratch', '0.394 (Late Fusion B3)', '0.394 (FCNN B3)', '-0.001'],
    ['7-class TL', '0.180 (Intermediate TL B3)', '0.232 (Intermediate TL B1)', '-0.052'],
    ['4-class TL', '0.412 (Intermediate TL B1)', '0.407 (CNN TL B2)', '+0.005'],
], 0.3, 0.8, 9.4, 1.6, cw=[2.0, 2.8, 2.8, 1.8], hsz=10, rsz=10, highlight_best=3)

# Temuan boxes
m1 = tb(s, 0.3, 2.55, 4.5, 1.6)
m1.text_frame.word_wrap = True; set_fill(m1, 'E6F4EA')
para(m1.text_frame, 'Temuan Positif', sz=10, bold=True, color=COLOR_GREEN, first=True)
para(m1.text_frame, '', sz=4)
para(m1.text_frame, '-> Best overall front-only (0.412) sedikit', sz=9.5, color=COLOR_TEXT)
para(m1.text_frame, '   lebih baik dari front+side (0.407)', sz=9.5, color=COLOR_TEXT)
para(m1.text_frame, '-> Konsistensi data > kuantitas data', sz=9.5, color=COLOR_TEXT)
para(m1.text_frame, '-> CNN membaik tanpa variasi sudut', sz=9.5, color=COLOR_TEXT)

m2 = tb(s, 4.9, 2.55, 4.8, 1.6)
m2.text_frame.word_wrap = True; set_fill(m2, 'FCE4E4')
para(m2.text_frame, 'Trade-off', sz=10, bold=True, color=COLOR_RED, first=True)
para(m2.text_frame, '', sz=4)
para(m2.text_frame, '-> FCNN turun -0.076 di 7-class (kehilangan', sz=9.5, color=COLOR_TEXT)
para(m2.text_frame, '   variasi landmark dari side view)', sz=9.5, color=COLOR_TEXT)
para(m2.text_frame, '-> Dataset berkurang 28% (9,894 -> 7,091)', sz=9.5, color=COLOR_TEXT)
para(m2.text_frame, '-> Tapi best model tidak terdampak', sz=9.5, color=COLOR_TEXT)

print('  91: Perbandingan')

# ═══ SLIDE 92: Evaluasi Robustness ═══
s = prs.slides[91]
clear_slide(s)

t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Evaluasi Robustness: LOSO, 5-Fold CV, Random Split', sz=18, bold=True, color=COLOR_TEXT, first=True)

table(s, ['Strategi', 'Cara', 'Fold', 'Data Leakage?'], [
    ['LOSO', '1 user = 1 test set, rotasi', '37', 'Tidak'],
    ['5-Fold CV', 'User dibagi 5 grup, rotasi', '5', 'Tidak'],
    ['Random Split', 'Sampel diacak tanpa peduli user', '5x', 'Ya (baseline)'],
], 0.3, 0.75, 9.4, 1.25, cw=[2.0, 3.5, 1.2, 2.7], hsz=10, rsz=10)

table(s, ['Model', 'Single Split', 'Random Split', '5-Fold CV', 'LOSO'], [
    ['Intermediate TL', '0.412', '0.586 +/- 0.032', 'pending', '0.370 +/- 0.125'],
    ['Late Fusion', '0.394', '0.580 +/- 0.032', 'pending', 'pending'],
    ['FCNN', '0.361', '0.471 +/- 0.026', 'pending', 'pending'],
], 0.3, 2.2, 9.4, 1.25, cw=[2.2, 1.5, 2.0, 1.85, 1.85], hsz=10, rsz=10)

b = tb(s, 0.3, 3.6, 9.4, 0.35)
para(b.text_frame, 'Pola: Random (0.586) >> Single (0.412) > LOSO (0.370) -- semakin ketat, semakin jujur',
     sz=10, italic=True, color=COLOR_GRAY, align=PP_ALIGN.CENTER, first=True)
print('  92: Evaluasi Robustness')

# ═══ SLIDE 93: Temuan Data Leakage ═══
s = prs.slides[92]
clear_slide(s)

t = tb(s, 0.3, 0.12, 9.4, 0.55)
para(t.text_frame, 'Temuan: Bukti Data Leakage & Pentingnya User-Based Split', sz=18, bold=True, color=COLOR_TEXT, first=True)

table(s, ['Model', 'User-Based', 'Random Split', 'Selisih', 'Kenaikan'], [
    ['Intermediate TL', '0.412', '0.586', '+0.174', '+42%'],
    ['Late Fusion', '0.394', '0.580', '+0.186', '+47%'],
    ['FCNN', '0.361', '0.471', '+0.110', '+30%'],
], 0.3, 0.8, 9.4, 1.2, cw=[2.2, 1.8, 1.8, 1.8, 1.8], hsz=10, rsz=10)

# Boxes
m1 = tb(s, 0.3, 2.15, 4.5, 1.5)
m1.text_frame.word_wrap = True; set_fill(m1, 'FCE4E4')
para(m1.text_frame, 'Mengapa Random Split Lebih Tinggi?', sz=10, bold=True, color=COLOR_RED, first=True)
para(m1.text_frame, '', sz=4)
para(m1.text_frame, 'User yang sama di train & test', sz=9.5, color=COLOR_TEXT)
para(m1.text_frame, '-> Model "menghafal" wajah, bukan ekspresi', sz=9.5, color=COLOR_TEXT)
para(m1.text_frame, '-> F1 naik 30-47% bukan karena model', sz=9.5, color=COLOR_TEXT)
para(m1.text_frame, '   lebih baik, tapi evaluasi tidak valid', sz=9.5, color=COLOR_TEXT)

m2 = tb(s, 4.9, 2.15, 4.8, 1.5)
m2.text_frame.word_wrap = True; set_fill(m2, 'E6F4EA')
para(m2.text_frame, 'LOSO: Performa Sebenarnya', sz=10, bold=True, color=COLOR_GREEN, first=True)
para(m2.text_frame, '', sz=4)
para(m2.text_frame, 'Intermediate TL LOSO: 0.370 +/- 0.125', sz=9.5, color=COLOR_TEXT)
para(m2.text_frame, '-> Std tinggi: performa bervariasi antar user', sz=9.5, color=COLOR_TEXT)
para(m2.text_frame, '-> Single split (0.412) sedikit over-estimate', sz=9.5, color=COLOR_TEXT)
para(m2.text_frame, '-> LOSO = evaluasi paling jujur', sz=9.5, color=COLOR_TEXT)

# Conclusion bar
c = tb(s, 0.3, 3.8, 9.4, 0.5)
c.text_frame.word_wrap = True; set_fill(c, '1A73E8')
para(c.text_frame, 'Kesimpulan: Evaluasi HARUS user-based split (LOSO/CV) -- random split menyesatkan (+30-47%)',
     sz=10, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER, first=True)
print('  93: Temuan Data Leakage')

# ── SAVE ──
prs.save(PPTX_PATH)
print(f'\nSaved! Total: {len(prs.slides)} slides')

for i in range(85, len(prs.slides)):
    slide = prs.slides[i]
    title = next((s.text_frame.text.strip()[:55].replace(chr(10),' ').encode('ascii','replace').decode()
                  for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), '')
    print(f'  Slide {i+1}: {title}')
