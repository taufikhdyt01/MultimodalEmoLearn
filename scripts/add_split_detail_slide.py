"""
Add slide with batch-level split detail after the existing Split Dataset slide.
Shows: Batch 1 vs Batch 2 breakdown per train/val/test.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_BLUE       = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_TEXT       = RGBColor(0x20, 0x20, 0x20)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY       = RGBColor(0x66, 0x66, 0x66)


def set_fill(shape, hex_val):
    spPr = shape._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', hex_val)


def add_tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def add_para(tf, text, sz=10, bold=False, italic=False, color=None, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = color


def cell(c, text, bold=False, bg=None, fg=None, sz=9, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]; p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    if fg: r.font.color.rgb = fg
    if bg: c.fill.solid(); c.fill.fore_color.rgb = bg


def table(slide, headers, rows, l, t, w, h, cw=None, hsz=9.5, rsz=9.5):
    shape = slide.shapes.add_table(len(rows)+1, len(headers), Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None: tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw): tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WHITE, sz=hsz, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LIGHT_BLUE if r % 2 == 0 else COLOR_WHITE
        for c, val in enumerate(row):
            b = (c == 0) or (r == len(rows) - 1)  # bold first col and last row
            cell(tbl.cell(r+1, c), str(val), bold=b, bg=bg, fg=COLOR_TEXT, sz=rsz,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


def insert(prs, index):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    xml_slides = prs.slides._sldIdLst
    el = xml_slides[-1]; xml_slides.remove(el); xml_slides.insert(index, el)
    return slide


# ── LOAD ──
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Find Split Dataset slide
target_idx = None
for i in range(60, 70):
    slide = prs.slides[i]
    for s in slide.shapes:
        if s.has_text_frame and 'Split Dataset' in s.text_frame.text:
            target_idx = i
            break
    if target_idx is not None:
        break

print(f'Split Dataset slide: {target_idx + 1}')

# Insert after Split Dataset slide
s = insert(prs, target_idx + 1)

# Title
t = add_tb(s, 0.3, 0.12, 9.4, 0.55)
add_para(t.text_frame, 'Detail Split per Batch (Front-Only)', sz=20, bold=True, color=COLOR_TEXT, first=True)

# Main table: Batch x Split
table(s,
    headers=['', 'Train', 'Val', 'Test', 'Total'],
    rows=[
        ['Batch 1 (20 user)', '18 user (3,334)', '1 user (232)', '1 user (258)', '20 user (3,824)'],
        ['Batch 2 (17 user)', '11 user (2,014)', '2 user (475)', '4 user (778)', '17 user (3,267)'],
        ['Total', '29 user (5,348)', '3 user (707)', '5 user (1,036)', '37 user (7,091)'],
    ],
    l=0.3, t=0.8, w=9.4, h=1.3,
    cw=[2.2, 2.0, 1.8, 1.8, 1.6], hsz=10, rsz=10)

# User IDs detail
tb_detail = add_tb(s, 0.3, 2.25, 9.4, 2.2)
tf = tb_detail.text_frame; tf.word_wrap = True
set_fill(tb_detail, 'F1F3F4')

add_para(tf, 'Detail User ID per Split:', sz=10, bold=True, color=COLOR_TEXT, first=True)
add_para(tf, '', sz=4)
add_para(tf, 'Train (29 user):', sz=9.5, bold=True, color=COLOR_BLUE)
add_para(tf, '  Batch 1: 97, 99, 100, 101, 102, 103, 104, 106, 107, 108, 109, 110, 113, 114, 115, 116, 117, 118',
         sz=9, color=COLOR_TEXT)
add_para(tf, '  Batch 2: 200, 201, 203, 205, 208, 210, 211, 212, 213, 215, 216',
         sz=9, color=COLOR_TEXT)
add_para(tf, '', sz=4)
add_para(tf, 'Validation (3 user):', sz=9.5, bold=True, color=COLOR_BLUE)
add_para(tf, '  Batch 1: 112  |  Batch 2: 207, 214',
         sz=9, color=COLOR_TEXT)
add_para(tf, '', sz=4)
add_para(tf, 'Test (5 user):', sz=9.5, bold=True, color=COLOR_BLUE)
add_para(tf, '  Batch 1: 111  |  Batch 2: 197, 202, 206, 209',
         sz=9, color=COLOR_TEXT)

# Note
tb_note = add_tb(s, 0.3, 4.55, 9.4, 0.35)
add_para(tb_note.text_frame,
         'Split identik untuk eksperimen front+side dan front-only (user yang sama di setiap split)',
         sz=8.5, italic=True, color=COLOR_GRAY, align=PP_ALIGN.CENTER, first=True)

print(f'Inserted slide after {target_idx + 1}')

# ── SAVE ──
prs.save(PPTX_PATH)
print(f'Saved! Total: {len(prs.slides)} slides')
