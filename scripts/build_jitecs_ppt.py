"""Build JITeCS paper proposal PPT (separate from bimbingan PPT)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT = 'd:/MultimodalEmoLearn/docs/PPT JITeCS.pptx'

# ─────────── Style constants (match bimbingan PPT) ───────────
COLOR_BLUE = RGBColor(0x1A, 0x73, 0xE8)
COLOR_LB = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_WH = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x20, 0x20, 0x20)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_RED = RGBColor(0xD9, 0x3C, 0x3C)
COLOR_GREEN = RGBColor(0x0F, 0x9D, 0x58)
COLOR_WARN = RGBColor(0xFC, 0xE4, 0xE4)
COLOR_OK = RGBColor(0xE6, 0xF4, 0xEA)
COLOR_YELLOW = RGBColor(0xFF, 0xF0, 0xCC)


def set_fill(shape, hex_val):
    spPr = shape._element.spPr
    for nf in spPr.findall(qn('a:noFill')):
        spPr.remove(nf)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', hex_val)


def tb(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def para(tf, text, sz=10, bold=False, italic=False, color=None,
         align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    if color:
        r.font.color.rgb = color


def cell(c, text, bold=False, bg=None, fg=None, sz=9, align=PP_ALIGN.LEFT):
    c.text = text
    p = c.text_frame.paragraphs[0]
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    if fg:
        r.font.color.rgb = fg
    if bg:
        c.fill.solid()
        c.fill.fore_color.rgb = bg


def table(slide, headers, rows, l, t, w, h, cw=None,
          hsz=9, rsz=8.5, highlight_rows=None):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                    Inches(l), Inches(t),
                                    Inches(w), Inches(h))
    tbl = shape.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        sid = tblPr.find(qn('a:tableStyleId'))
        if sid is not None:
            tblPr.remove(sid)
    if cw:
        tot = sum(cw)
        for i, c in enumerate(cw):
            tbl.columns[i].width = int(Inches(w) * c / tot)
    for c, h in enumerate(headers):
        cell(tbl.cell(0, c), h, bold=True, bg=COLOR_BLUE, fg=COLOR_WH,
             sz=hsz, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_LB if r % 2 == 0 else COLOR_WH
        if highlight_rows and r in highlight_rows:
            bg = COLOR_OK
        for c, val in enumerate(row):
            cell(tbl.cell(r + 1, c), str(val), bold=(c == 0),
                 bg=bg, fg=COLOR_TEXT, sz=rsz,
                 align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)


def title(s, text):
    t = tb(s, 0.3, 0.12, 9.4, 0.55)
    para(t.text_frame, text, sz=18, bold=True, color=COLOR_TEXT, first=True)


def new_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


# ─────────── Build presentation ───────────
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# ─────────── SLIDE 1: Title ───────────
s = new_slide(prs)
set_fill(s.shapes.add_shape(1, Inches(0), Inches(0),
                             Inches(10), Inches(5.625)), 'F5F9FF')

t = tb(s, 0.5, 1.4, 9.0, 0.9)
para(t.text_frame,
     'Multimodal Fusion of Facial Image and Landmark Features',
     sz=22, bold=True, color=COLOR_BLUE, align=PP_ALIGN.CENTER, first=True)
para(t.text_frame,
     'with Transfer Learning for Emotion Recognition',
     sz=22, bold=True, color=COLOR_BLUE, align=PP_ALIGN.CENTER)
para(t.text_frame,
     'in Programming Learning Context',
     sz=22, bold=True, color=COLOR_BLUE, align=PP_ALIGN.CENTER)

t2 = tb(s, 0.5, 3.1, 9.0, 0.4)
para(t2.text_frame, 'Proposal Publikasi Paper — JITeCS',
     sz=14, italic=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER, first=True)

t3 = tb(s, 0.5, 3.8, 9.0, 0.8)
para(t3.text_frame, 'Taufik Hidayat', sz=13, bold=True,
     color=COLOR_TEXT, align=PP_ALIGN.CENTER, first=True)
para(t3.text_frame, 'Magister Ilmu Komputer',
     sz=11, color=COLOR_GRAY, align=PP_ALIGN.CENTER)
para(t3.text_frame, 'Subset dari Tahap 4 Tesis — April 2026',
     sz=10, italic=True, color=COLOR_GRAY, align=PP_ALIGN.CENTER)

# ─────────── Helper paragraph renderers ───────────
def sec_line(tf, num, text, sz=11.5, first=False):
    """Section heading: large bold blue."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    r1 = p.add_run(); r1.text = f'{num}  '
    r1.font.size = Pt(sz); r1.font.bold = True; r1.font.color.rgb = COLOR_BLUE
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(sz); r2.font.bold = True; r2.font.color.rgb = COLOR_BLUE


def sub_line(tf, num, text, sz=10):
    """Subsection heading: bold dark text, italic."""
    p = tf.add_paragraph()
    r1 = p.add_run(); r1.text = f'    {num}  '
    r1.font.size = Pt(sz); r1.font.bold = True; r1.font.color.rgb = COLOR_TEXT
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(sz); r2.font.bold = True
    r2.font.italic = True; r2.font.color.rgb = COLOR_TEXT


def point_line(tf, text, sz=8.5, indent='        '):
    """Content bullet: small gray italic, indented."""
    p = tf.add_paragraph()
    r = p.add_run(); r.text = f'{indent}•  {text}'
    r.font.size = Pt(sz); r.font.color.rgb = COLOR_GRAY


def spacer(tf, sz=4):
    sp = tf.add_paragraph()
    r = sp.add_run(); r.text = ''
    r.font.size = Pt(sz)


# ─────────── SLIDE 2: Daftar Isi Paper (Part 1) ───────────
s = new_slide(prs)
title(s, 'Daftar Isi Paper (1/2)')

# Legend box at top
lg = tb(s, 0.3, 0.75, 9.4, 0.32)
lg.text_frame.word_wrap = True
p = lg.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = 'BAB  '
r1.font.size = Pt(9); r1.font.bold = True; r1.font.color.rgb = COLOR_BLUE
r2 = p.add_run(); r2.text = '= judul bab   |   '
r2.font.size = Pt(9); r2.font.color.rgb = COLOR_GRAY
r3 = p.add_run(); r3.text = 'Sub bab  '
r3.font.size = Pt(9); r3.font.bold = True; r3.font.italic = True
r3.font.color.rgb = COLOR_TEXT
r4 = p.add_run(); r4.text = '= judul sub bab   |   '
r4.font.size = Pt(9); r4.font.color.rgb = COLOR_GRAY
r5 = p.add_run(); r5.text = '•  '
r5.font.size = Pt(9); r5.font.color.rgb = COLOR_GRAY
r6 = p.add_run(); r6.text = 'isi poin yang dibahas dalam sub bab'
r6.font.size = Pt(9); r6.font.italic = True; r6.font.color.rgb = COLOR_GRAY

# Left column
left = tb(s, 0.3, 1.15, 4.7, 4.3)
left.text_frame.word_wrap = True

sec_line(left.text_frame, 'Abstract', '', sz=11.5, first=True)
point_line(left.text_frame, 'Problem: FER di konteks programming natural')
point_line(left.text_frame, 'Method: 5 arsitektur (CNN/FCNN/Early/Intermediate/Late) × TL')
point_line(left.text_frame, 'Data: 7.091 samples dari 37 mahasiswa')
point_line(left.text_frame, 'Best: Intermediate Fusion TL 4c B1 = F1 0.412')
point_line(left.text_frame, 'Insight: Intermediate + TL unggul di data natural')
spacer(left.text_frame)

sec_line(left.text_frame, '1.', 'Introduction')
point_line(left.text_frame,
           'Motivasi: emosi mahasiswa → learning outcome programming')
point_line(left.text_frame,
           'Gap: FER existing fokus data lab (posed), bukan natural')
point_line(left.text_frame, 'Rumusan masalah & tujuan penelitian')
point_line(left.text_frame,
           'Kontribusi: dataset baru + studi sistematis fusion × TL')
point_line(left.text_frame, 'Struktur paper')

# Right column
right = tb(s, 5.1, 1.15, 4.6, 4.3)
right.text_frame.word_wrap = True

sec_line(right.text_frame, '2.', 'Related Work', first=True)
sub_line(right.text_frame, '2.1', 'Deep Learning for FER')
sub_line(right.text_frame, '2.2', 'Multimodal Fusion of Image and Landmark Features')
sub_line(right.text_frame, '2.3', 'Transfer Learning for FER')
sub_line(right.text_frame, '2.4', 'Affective Computing in Education')
spacer(right.text_frame)

sec_line(right.text_frame, '3.', 'Proposed Method')
sub_line(right.text_frame, '3.1', 'Dataset')
point_line(right.text_frame, 'Akuisisi 37 user, 2 batch, sesi programming')
point_line(right.text_frame, 'Anotasi Face API + confidence score')
point_line(right.text_frame, 'Statistik & distribusi (imbalanced)')
point_line(right.text_frame, 'Preprocessing (resize, MediaPipe landmark)')
point_line(right.text_frame, 'Split strategi user-wise 80/10/10')
sub_line(right.text_frame, '3.2', 'Multimodal Architecture')
point_line(right.text_frame, 'CNN branch (vanilla + ResNet18 TL)')
point_line(right.text_frame, 'FCNN branch (68 landmark, 136-dim)')
point_line(right.text_frame, 'Fusion: Intermediate, Late')

# ─────────── SLIDE 3: Daftar Isi Paper (Part 2) ───────────
s = new_slide(prs)
title(s, 'Daftar Isi Paper (2/2)')

# Left column — continues section 3 and starts 4
left = tb(s, 0.3, 0.85, 4.7, 4.5)
left.text_frame.word_wrap = True

sec_line(left.text_frame, '3.', 'Proposed Method (lanjutan)', first=True)
sub_line(left.text_frame, '3.3', 'Training Setup')
point_line(left.text_frame, 'Adam optimizer, ReduceLROnPlateau scheduler')
point_line(left.text_frame, 'Early stopping by macro F1 (patience=15)')
sub_line(left.text_frame, '3.4', 'Experimental Design')
point_line(left.text_frame,
           'Scenarios: B1 Baseline, B2 Class Weights, B3 Augmentation')
point_line(left.text_frame,
           'Class config: 7-class vs 4-class (minoritas → negative)')
point_line(left.text_frame,
           'Evaluation metrics: Macro F1 (utama), Micro F1, Weighted F1')
spacer(left.text_frame)

sec_line(left.text_frame, '4.', 'Experimental Results')
sub_line(left.text_frame, '4.1', 'Overall Performance')
sub_line(left.text_frame, '4.2', 'Effect of Transfer Learning')
sub_line(left.text_frame, '4.3', 'Effect of Fusion Strategy')
sub_line(left.text_frame, '4.4', 'Effect of Class Granularity')
sub_line(left.text_frame, '4.5', 'Per-Class Analysis')

# Right column — section 5, 6, references
right = tb(s, 5.1, 0.85, 4.6, 4.5)
right.text_frame.word_wrap = True

sec_line(right.text_frame, '5.', 'Discussion', first=True)
sub_line(right.text_frame, '5.1', 'Multimodal Fusion vs Single-Modality')
point_line(right.text_frame, 'Menjawab RQ1 (fusi image+landmark vs single)')
sub_line(right.text_frame, '5.2', 'Fusion Strategy Comparison: Intermediate vs Late')
point_line(right.text_frame, 'Menjawab RQ2 (strategi fusion terbaik)')
sub_line(right.text_frame, '5.3', 'Transfer Learning Effectiveness')
point_line(right.text_frame, 'Menjawab RQ3 (gain TL pada data natural)')
sub_line(right.text_frame, '5.4', 'Limitations')
point_line(right.text_frame, 'Minority class langka, 37 subjek, label noise')
sub_line(right.text_frame, '5.5', 'Implications for Learning Analytics')
spacer(right.text_frame)

sec_line(right.text_frame, '6.', 'Conclusion')
point_line(right.text_frame, 'Ringkasan 3 kontribusi & best result')
point_line(right.text_frame,
           'Future work: lebih banyak user, anotasi manual, real-time')
spacer(right.text_frame)

sec_line(right.text_frame, '', 'References')
point_line(right.text_frame,
           '20-25 sitasi (SOTA FER + affective computing + multimodal DL)')

# ─────────── SLIDE 4: Latar Belakang ───────────
s = new_slide(prs)
title(s, 'Latar Belakang')

# Problem box
b1 = tb(s, 0.3, 0.85, 9.4, 1.1)
b1.text_frame.word_wrap = True
set_fill(b1, 'E8F0FE')
para(b1.text_frame, 'Masalah:', sz=11, bold=True, color=COLOR_BLUE, first=True)
para(b1.text_frame,
     '• Emosi mahasiswa saat sesi pemrograman memengaruhi learning outcome '
     '(frustrasi, kebingungan, fokus).',
     sz=10, color=COLOR_TEXT)
para(b1.text_frame,
     '• Deteksi otomatis diperlukan untuk learning analytics, '
     'tetapi pengenalan emosi wajah (FER) berbasis satu modalitas rentan error '
     'pada kondisi natural.',
     sz=10, color=COLOR_TEXT)

# Gap box
b2 = tb(s, 0.3, 2.1, 9.4, 1.4)
b2.text_frame.word_wrap = True
set_fill(b2, 'FFF0CC')
para(b2.text_frame, 'Gap Penelitian:', sz=11, bold=True,
     color=COLOR_TEXT, first=True)
para(b2.text_frame,
     '• FER yang ada mayoritas dievaluasi pada dataset lab '
     '(CK+, JAFFE, KDEF) dengan ekspresi dibuat-buat (posed).',
     sz=10, color=COLOR_TEXT)
para(b2.text_frame,
     '• Belum banyak riset FER multimodal (citra wajah + facial landmark) '
     'pada konteks pembelajaran pemrograman (natural expression).',
     sz=10, color=COLOR_TEXT)
para(b2.text_frame,
     '• Perbandingan sistematis strategi fusion × transfer learning '
     'pada data natural masih terbatas.',
     sz=10, color=COLOR_TEXT)

# Contribution box
b3 = tb(s, 0.3, 3.7, 9.4, 1.55)
b3.text_frame.word_wrap = True
set_fill(b3, 'E6F4EA')
para(b3.text_frame, 'Kontribusi Paper:', sz=11, bold=True,
     color=COLOR_GREEN, first=True)
para(b3.text_frame,
     '1. Dataset multimodal baru: 7.091 sampel dari 37 mahasiswa '
     'selama sesi pemrograman natural.',
     sz=10, color=COLOR_TEXT)
para(b3.text_frame,
     '2. Studi komparatif sistematis 5 arsitektur (CNN, FCNN, Early Fusion, '
     'Intermediate Fusion, Late Fusion) dengan/tanpa transfer learning — '
     'mencakup seluruh spektrum fusion point (input → feature → decision).',
     sz=10, color=COLOR_TEXT)
para(b3.text_frame,
     '3. Evaluasi menyeluruh 60 konfigurasi dengan 3 metrik F1 '
     '(Macro, Micro, Weighted) untuk data imbalanced.',
     sz=10, color=COLOR_TEXT)

# ─────────── SLIDE 5: Rumusan Masalah & Tujuan ───────────
s = new_slide(prs)
title(s, 'Rumusan Masalah & Tujuan')

rm = tb(s, 0.3, 0.85, 9.4, 2.65)
rm.text_frame.word_wrap = True
set_fill(rm, 'FCE4E4')
para(rm.text_frame, 'Rumusan Masalah (Research Questions)', sz=12,
     bold=True, color=COLOR_RED, first=True)
para(rm.text_frame, '', sz=3)

# RQ1
p = rm.text_frame.add_paragraph()
r1 = p.add_run(); r1.text = 'RQ1: '
r1.font.size = Pt(10); r1.font.bold = True; r1.font.color.rgb = COLOR_RED
r2 = p.add_run(); r2.text = (
    'Apakah fusi multimodal antara citra wajah dan facial landmark '
    'memberikan kinerja yang lebih baik dibandingkan pendekatan '
    'single-modality untuk pengenalan emosi pada konteks pembelajaran '
    'pemrograman?')
r2.font.size = Pt(10); r2.font.color.rgb = COLOR_TEXT
para(rm.text_frame, '', sz=3)

# RQ2
p = rm.text_frame.add_paragraph()
r1 = p.add_run(); r1.text = 'RQ2: '
r1.font.size = Pt(10); r1.font.bold = True; r1.font.color.rgb = COLOR_RED
r2 = p.add_run(); r2.text = (
    'Strategi fusion manakah — Intermediate Fusion atau Late Fusion — '
    'yang lebih efektif dalam menangani data ekspresi wajah yang '
    'natural dan imbalanced?')
r2.font.size = Pt(10); r2.font.color.rgb = COLOR_TEXT
para(rm.text_frame, '', sz=3)

# RQ3
p = rm.text_frame.add_paragraph()
r1 = p.add_run(); r1.text = 'RQ3: '
r1.font.size = Pt(10); r1.font.bold = True; r1.font.color.rgb = COLOR_RED
r2 = p.add_run(); r2.text = (
    'Bagaimana pengaruh transfer learning (ResNet18 pretrained '
    'ImageNet) terhadap kinerja model pengenalan emosi wajah '
    'multimodal pada dataset kecil dengan ekspresi natural?')
r2.font.size = Pt(10); r2.font.color.rgb = COLOR_TEXT

# Tujuan (moved below RQ)
tuj = tb(s, 0.3, 3.6, 9.4, 1.65)
tuj.text_frame.word_wrap = True
set_fill(tuj, 'E6F4EA')
para(tuj.text_frame, 'Tujuan Penelitian', sz=12, bold=True,
     color=COLOR_GREEN, first=True)
para(tuj.text_frame, '', sz=2)
para(tuj.text_frame,
     '1. Mengembangkan pipeline FER multimodal (citra + landmark) '
     'untuk konteks pembelajaran pemrograman.',
     sz=10, color=COLOR_TEXT)
para(tuj.text_frame,
     '2. Menentukan arsitektur fusion dan strategi transfer learning '
     'yang optimal melalui studi 48 konfigurasi eksperimen.',
     sz=10, color=COLOR_TEXT)
para(tuj.text_frame,
     '3. Menyajikan analisis per-kelas emosi untuk memahami perilaku '
     'model pada data minoritas dalam konteks natural.',
     sz=10, color=COLOR_TEXT)

# Target venue
v = tb(s, 0.3, 3.95, 9.4, 1.25)
v.text_frame.word_wrap = True
set_fill(v, 'E8F0FE')
para(v.text_frame, 'Target Publikasi: JITeCS (Journal of Information Technology '
     'and Computer Science)', sz=11, bold=True, color=COLOR_BLUE, first=True)
para(v.text_frame,
     '• Terindeks SINTA 2 • Accredited oleh Kemenristekdikti • Peer-reviewed',
     sz=9.5, color=COLOR_TEXT)
para(v.text_frame,
     '• Scope: IT, data science, machine learning, applied computing '
     '— relevan dengan topik FER untuk education tech.',
     sz=9.5, color=COLOR_TEXT)
para(v.text_frame,
     '• Format: 8-12 halaman IEEE, 20-30 referensi.',
     sz=9.5, color=COLOR_TEXT)

# ─────────── SLIDE 6: Related Work (placeholder) ───────────
s = new_slide(prs)
title(s, 'Related Work')

# Categories to be reviewed
cat = tb(s, 0.3, 0.85, 9.4, 3.1)
cat.text_frame.word_wrap = True
set_fill(cat, 'E8F0FE')
para(cat.text_frame, 'Kategori Literatur yang akan Direview',
     sz=12, bold=True, color=COLOR_BLUE, first=True)
para(cat.text_frame, '', sz=3)
para(cat.text_frame, '1. Deep Learning untuk Facial Expression Recognition',
     sz=11, bold=True, color=COLOR_TEXT)
para(cat.text_frame,
     '   CNN, ResNet, EfficientNet, Vision Transformer pada FER klasik.',
     sz=10, color=COLOR_GRAY)
para(cat.text_frame, '', sz=2)
para(cat.text_frame, '2. Multimodal Fusion untuk FER',
     sz=11, bold=True, color=COLOR_TEXT)
para(cat.text_frame,
     '   Strategi early / intermediate / late fusion pada kombinasi '
     'citra wajah + landmark / audio / fisiologis.',
     sz=10, color=COLOR_GRAY)
para(cat.text_frame, '', sz=2)
para(cat.text_frame, '3. Transfer Learning pada FER',
     sz=11, bold=True, color=COLOR_TEXT)
para(cat.text_frame,
     '   ImageNet pretrained backbone, domain adaptation, fine-tuning '
     'strategy untuk dataset kecil.',
     sz=10, color=COLOR_GRAY)
para(cat.text_frame, '', sz=2)
para(cat.text_frame, '4. Affective Computing dalam Konteks Pendidikan',
     sz=11, bold=True, color=COLOR_TEXT)
para(cat.text_frame,
     '   Emotion recognition untuk learning analytics, MOOC engagement, '
     'adaptive tutoring system.',
     sz=10, color=COLOR_GRAY)

# TBD note
note = tb(s, 0.3, 4.15, 9.4, 1.1)
note.text_frame.word_wrap = True
set_fill(note, 'FFF0CC')
para(note.text_frame, 'Catatan', sz=10, bold=True,
     color=COLOR_TEXT, first=True)
para(note.text_frame,
     '• Pemilihan paper spesifik (20-25 sitasi) dilakukan saat drafting — '
     'fokus pada relevansi metodologis.',
     sz=9.5, color=COLOR_TEXT)
para(note.text_frame,
     '• Posisi paper ini akan dirumuskan setelah review selesai '
     '(konsultasi dengan pembimbing).',
     sz=9.5, color=COLOR_TEXT)

# ─────────── SLIDE 7: Proposed Method — Dataset ───────────
s = new_slide(prs)
title(s, 'Proposed Method — Dataset (Section 3.1)')

# Overview
ov = tb(s, 0.3, 0.85, 4.6, 2.0)
ov.text_frame.word_wrap = True
set_fill(ov, 'E8F0FE')
para(ov.text_frame, 'Karakteristik', sz=11, bold=True,
     color=COLOR_BLUE, first=True)
para(ov.text_frame, '', sz=3)
para(ov.text_frame, '• 37 mahasiswa (2 batch data collection)',
     sz=10, color=COLOR_TEXT)
para(ov.text_frame, '• 7.091 frame wajah (front-view)',
     sz=10, color=COLOR_TEXT)
para(ov.text_frame, '• Durasi: sesi pemrograman 30-60 menit',
     sz=10, color=COLOR_TEXT)
para(ov.text_frame, '• 7 kelas emosi (FER standar)',
     sz=10, color=COLOR_TEXT)
para(ov.text_frame, '• Anotasi: Face API (confidence score)',
     sz=10, color=COLOR_TEXT)
para(ov.text_frame, '• Split: user-wise 80/10/10 (tidak ada leakage)',
     sz=10, color=COLOR_TEXT)

# Distribution table
para(tb(s, 5.0, 0.85, 4.7, 0.25).text_frame,
     'Distribusi Kelas Emosi', sz=10, bold=True,
     color=COLOR_TEXT, first=True)
table(s, ['Emosi', 'Jumlah', '%'],
      [
          ['Neutral', '5.519', '77.8%'],
          ['Happy', '651', '9.2%'],
          ['Sad', '361', '5.1%'],
          ['Surprised', '39', '0.5%'],
          ['Angry', '32', '0.5%'],
          ['Disgusted', '16', '0.2%'],
          ['Fearful', '5', '0.1%'],
      ],
      5.0, 1.15, 4.7, 1.9, cw=[2.0, 1.5, 1.2], hsz=9, rsz=8,
      highlight_rows=None)

# Preprocessing + split
pp = tb(s, 0.3, 3.0, 9.4, 2.3)
pp.text_frame.word_wrap = True
set_fill(pp, 'FFF0CC')
para(pp.text_frame, 'Preprocessing Pipeline', sz=11, bold=True,
     color=COLOR_TEXT, first=True)
para(pp.text_frame, '', sz=3)
para(pp.text_frame,
     '1. Face detection → crop wajah → resize 224×224 RGB (untuk CNN branch).',
     sz=10, color=COLOR_TEXT)
para(pp.text_frame,
     '2. Facial landmark extraction via MediaPipe (478 titik → '
     'subset 68 klasik = 136 koordinat).',
     sz=10, color=COLOR_TEXT)
para(pp.text_frame,
     '3. Normalisasi: piksel [0,1], landmark normalized ke frame '
     '(agar invariant terhadap resolusi).',
     sz=10, color=COLOR_TEXT)
para(pp.text_frame,
     '4. Remap 4-class opsional: neutral / happy / sad / negative '
     '(gabung angry+fearful+disgusted+surprised).',
     sz=10, color=COLOR_TEXT)
para(pp.text_frame,
     '5. Split user-wise: 29 user train / 3 user val / 5 user test '
     '→ mencegah identity-leakage.',
     sz=10, color=COLOR_TEXT)

# ─────────── SLIDE 8: Proposed Method — Arsitektur Model ───────────
s = new_slide(prs)
title(s, 'Proposed Method — Arsitektur Model (Section 3.2)')

table(s, ['Model', 'Input', 'Fusion Point', 'Keterangan'],
      [
          ['CNN', 'Citra 224×224×3', '—',
           'Single modal: visual'],
          ['FCNN', 'Landmark 136-dim', '—',
           'Single modal: geometric'],
          ['Early Fusion', 'Citra + heatmap (224×224×4)',
           'Input level (0%)',
           'Landmark Gaussian heatmap sbg channel ke-4'],
          ['Intermediate Fusion', 'Citra + Landmark',
           'Feature level (50%)',
           'CNN+FCNN concat di hidden layer'],
          ['Late Fusion', 'Citra + Landmark',
           'Decision level (95%)',
           'Softmax averaging 2 model terpisah'],
      ],
      0.3, 0.85, 9.4, 2.0, cw=[1.6, 2.2, 1.8, 3.0], hsz=9, rsz=8)

# TL note
tl = tb(s, 0.3, 3.05, 9.4, 1.1)
tl.text_frame.word_wrap = True
set_fill(tl, 'E8F0FE')
para(tl.text_frame, 'Transfer Learning Backbone', sz=11, bold=True,
     color=COLOR_BLUE, first=True)
para(tl.text_frame, '', sz=2)
para(tl.text_frame,
     '• ResNet18 pretrained ImageNet (1000-class) → fine-tune untuk FER',
     sz=9.5, color=COLOR_TEXT)
para(tl.text_frame,
     '• Early Fusion TL: first Conv2d dimodifikasi dari 3→4 channel; '
     'weight RGB di-copy, weight heatmap diinisialisasi dari mean(RGB)',
     sz=9.5, color=COLOR_TEXT)

# Training setup
tr = tb(s, 0.3, 4.25, 9.4, 1.1)
tr.text_frame.word_wrap = True
set_fill(tr, 'FFF0CC')
para(tr.text_frame, 'Training Setup', sz=11, bold=True,
     color=COLOR_TEXT, first=True)
para(tr.text_frame, '', sz=2)
para(tr.text_frame,
     '• Adam optimizer • ReduceLROnPlateau scheduler (factor=0.5, patience=8)',
     sz=10, color=COLOR_TEXT)
para(tr.text_frame,
     '• Early stopping by Macro F1 (patience=15) • Max 50 epoch '
     '• Batch size 32',
     sz=10, color=COLOR_TEXT)

# ─────────── SLIDE 9: Proposed Method — Skenario Eksperimen ───────────
s = new_slide(prs)
title(s, 'Proposed Method — Matriks Eksperimen (Section 3.3-3.6)')

# Dimensions
dm = tb(s, 0.3, 0.85, 4.6, 2.5)
dm.text_frame.word_wrap = True
set_fill(dm, 'E8F0FE')
para(dm.text_frame, 'Dimensi Eksperimen', sz=11, bold=True,
     color=COLOR_BLUE, first=True)
para(dm.text_frame, '', sz=3)
para(dm.text_frame, '5 arsitektur × 2 backbone × 3 skenario × 2 kelas',
     sz=10, bold=True, color=COLOR_TEXT)
para(dm.text_frame, '= 60 total konfigurasi', sz=10, bold=True,
     color=COLOR_GREEN)
para(dm.text_frame, '', sz=4)
para(dm.text_frame, 'Arsitektur:', sz=10, bold=True, color=COLOR_TEXT)
para(dm.text_frame, '  CNN, FCNN, Early, Intermediate, Late',
     sz=9.5, color=COLOR_TEXT)
para(dm.text_frame, 'Backbone: from-scratch, ResNet18 TL',
     sz=9.5, color=COLOR_TEXT)
para(dm.text_frame, 'Skenario: B1, B2, B3', sz=9.5, color=COLOR_TEXT)
para(dm.text_frame, 'Kelas: 7-class, 4-class', sz=9.5, color=COLOR_TEXT)

# Scenarios
sc = tb(s, 5.1, 0.85, 4.6, 2.5)
sc.text_frame.word_wrap = True
set_fill(sc, 'FFF0CC')
para(sc.text_frame, 'Tiga Skenario untuk Imbalance', sz=11, bold=True,
     color=COLOR_TEXT, first=True)
para(sc.text_frame, '', sz=3)
para(sc.text_frame, 'B1 — Baseline:', sz=10, bold=True, color=COLOR_TEXT)
para(sc.text_frame,
     '   Standard cross-entropy, tanpa intervention.',
     sz=9.5, color=COLOR_TEXT)
para(sc.text_frame, '', sz=2)
para(sc.text_frame, 'B2 — Class Weights:', sz=10, bold=True, color=COLOR_TEXT)
para(sc.text_frame,
     '   Weighted CE loss: weight ∝ 1/freq class',
     sz=9.5, color=COLOR_TEXT)
para(sc.text_frame, '   (beri bobot lebih ke kelas minoritas).',
     sz=9.5, color=COLOR_TEXT)
para(sc.text_frame, '', sz=2)
para(sc.text_frame, 'B3 — Data Augmentation:', sz=10, bold=True,
     color=COLOR_TEXT)
para(sc.text_frame,
     '   Rotation, flip, brightness untuk kelas minoritas '
     'sampai seimbang.',
     sz=9.5, color=COLOR_TEXT)

# Metrics
me = tb(s, 0.3, 3.5, 9.4, 1.75)
me.text_frame.word_wrap = True
set_fill(me, 'E6F4EA')
para(me.text_frame, 'Metrik Evaluasi (3 Jenis F1)', sz=11, bold=True,
     color=COLOR_GREEN, first=True)
para(me.text_frame, '', sz=3)
para(me.text_frame,
     '• Macro F1 — rata-rata F1 per-kelas (treats all classes equal) '
     '→ metrik utama untuk imbalanced data.',
     sz=10, color=COLOR_TEXT)
para(me.text_frame,
     '• Micro F1 — = accuracy (global precision/recall).',
     sz=10, color=COLOR_TEXT)
para(me.text_frame,
     '• Weighted F1 — rata-rata F1 dibobot jumlah sampel '
     '→ didominasi kelas mayoritas.',
     sz=10, color=COLOR_TEXT)

# ─────────── SLIDE 10: Hasil Utama ───────────
s = new_slide(prs)
title(s, 'Hasil Utama — Top 10 Konfigurasi (4-Class)')

table(s,
      ['Rank', 'Model', 'Backbone', 'Skenario', 'Macro F1', 'Accuracy'],
      [
          ['1', 'Intermediate Fusion', 'TL (ResNet18)', 'B1', '0.412', '0.740'],
          ['2', 'Late Fusion', 'Scratch', 'B3', '0.394', '0.715'],
          ['3', 'Late Fusion', 'TL', 'B3', '0.372', '0.702'],
          ['4', 'FCNN', '—', 'B3', '0.361', '0.688'],
          ['5', 'Late Fusion', 'TL', 'B1', '0.309', '0.695'],
          ['6', 'CNN', 'TL', 'B2', '0.297', '0.680'],
          ['7', 'Intermediate Fusion', 'Scratch', 'B2', '0.286', '0.671'],
          ['8', 'CNN', 'TL', 'B1', '0.274', '0.664'],
          ['9', 'CNN', 'TL', 'B3', '0.268', '0.658'],
          ['10', 'Intermediate Fusion', 'Scratch', 'B1', '0.243', '0.629'],
      ],
      0.3, 0.85, 9.4, 3.3, cw=[0.6, 2.2, 1.6, 1.2, 1.4, 1.4],
      hsz=9, rsz=8.5, highlight_rows=[0])

# Key finding
kf = tb(s, 0.3, 4.3, 9.4, 0.9)
kf.text_frame.word_wrap = True
set_fill(kf, 'E6F4EA')
para(kf.text_frame, 'Temuan Utama', sz=11, bold=True,
     color=COLOR_GREEN, first=True)
para(kf.text_frame,
     'Intermediate Fusion + Transfer Learning (ResNet18) unggul — '
     'joint feature learning dari citra & landmark lebih efektif daripada late fusion '
     'pada natural expression data yang imbalanced.',
     sz=10, color=COLOR_TEXT)

# ─────────── SLIDE 11: Analisis per Kelas ───────────
s = new_slide(prs)
title(s, 'Analisis per-Kelas (Best Model: Intermediate TL 4c B1)')

table(s,
      ['Kelas', 'Support Test', 'Precision', 'Recall', 'F1',
       'Catatan'],
      [
          ['Neutral', '816', '0.90', '0.97', '0.94', 'Dominan — mudah'],
          ['Happy', '57', '0.68', '0.53', '0.59', 'Cukup reliable'],
          ['Sad', '40', '0.43', '0.25', '0.32', 'Menengah'],
          ['Negative', '16', '0.40', '0.13', '0.20',
           'Minoritas — sulit (combined class)'],
      ],
      0.3, 0.85, 9.4, 1.6, cw=[1.2, 1.2, 1.1, 1.1, 1.0, 3.0],
      hsz=9, rsz=8.5, highlight_rows=[0])

# Interpretation
it = tb(s, 0.3, 2.65, 9.4, 2.6)
it.text_frame.word_wrap = True
set_fill(it, 'FFF0CC')
para(it.text_frame, 'Interpretasi', sz=11, bold=True,
     color=COLOR_TEXT, first=True)
para(it.text_frame, '', sz=3)
para(it.text_frame,
     '• Kelas minoritas (negative) menjadi bottleneck: '
     'hanya 16 sampel test → 1 kesalahan = 6% drop F1.',
     sz=10, color=COLOR_TEXT)
para(it.text_frame,
     '• Happy (F1 0.59) cukup baik meski hanya 9% dari total → '
     'menunjukkan model bisa pelajari ekspresi jelas.',
     sz=10, color=COLOR_TEXT)
para(it.text_frame,
     '• Sad (F1 0.32) menengah — ekspresi sedih sering ambigu '
     'dengan neutral di konteks pemrograman natural.',
     sz=10, color=COLOR_TEXT)
para(it.text_frame,
     '• Macro F1 0.412 mencerminkan kesulitan kelas minoritas, '
     'bukan kegagalan model secara keseluruhan.',
     sz=10, color=COLOR_TEXT)
para(it.text_frame,
     '• Confusion matrix menunjukkan mayoritas kesalahan adalah '
     '"other class → neutral" (over-prediction mayoritas).',
     sz=10, color=COLOR_TEXT)

# ─────────── SLIDE 12: Diskusi (Menjawab Rumusan Masalah) ───────────
s = new_slide(prs)
title(s, 'Diskusi — Menjawab Rumusan Masalah')

# RQ1: Multimodal vs single
rq1 = tb(s, 0.3, 0.85, 9.4, 1.4)
rq1.text_frame.word_wrap = True
set_fill(rq1, 'E8F0FE')
para(rq1.text_frame,
     '5.1  RQ1: Seberapa efektif fusi multimodal vs single-modality?',
     sz=11, bold=True, color=COLOR_BLUE, first=True)
para(rq1.text_frame, '', sz=2)
para(rq1.text_frame,
     '• Jawaban: Multimodal (Intermediate TL) 0.412 > CNN saja '
     '(0.274) > FCNN saja (0.361).',
     sz=9.5, color=COLOR_TEXT)
para(rq1.text_frame,
     '• Fusi citra + landmark memberikan gain signifikan karena '
     'kedua modality saling melengkapi — texture dari citra, '
     'geometric dari landmark.',
     sz=9.5, color=COLOR_TEXT)

# RQ2: Fusion strategy
rq2 = tb(s, 0.3, 2.35, 9.4, 1.4)
rq2.text_frame.word_wrap = True
set_fill(rq2, 'E6F4EA')
para(rq2.text_frame,
     '5.2  RQ2: Strategi fusion mana paling baik untuk data natural '
     'imbalanced? (Intermediate vs Late)',
     sz=11, bold=True, color=COLOR_GREEN, first=True)
para(rq2.text_frame, '', sz=2)
para(rq2.text_frame,
     '• Jawaban: Intermediate Fusion (0.412) > Late Fusion '
     '(0.394 best scratch / 0.372 TL).',
     sz=9.5, color=COLOR_TEXT)
para(rq2.text_frame,
     '• Joint feature learning (Intermediate) lebih baik karena CNN & FCNN '
     'dioptimalkan bersama → interaksi cross-modal dipelajari end-to-end.',
     sz=9.5, color=COLOR_TEXT)

# RQ3: Transfer learning
rq3 = tb(s, 0.3, 3.85, 9.4, 1.4)
rq3.text_frame.word_wrap = True
set_fill(rq3, 'FFF0CC')
para(rq3.text_frame,
     '5.3  RQ3: Apakah transfer learning (ResNet18 ImageNet) '
     'memberi gain pada dataset natural programming?',
     sz=11, bold=True, color=COLOR_TEXT, first=True)
para(rq3.text_frame, '', sz=2)
para(rq3.text_frame,
     '• Jawaban: Ya — TL memberi gain konsisten (+0.17 Macro F1 '
     'Intermediate TL vs scratch di 4-class B1).',
     sz=9.5, color=COLOR_TEXT)
para(rq3.text_frame,
     '• Dataset kecil (7.091 sampel) rentan overfitting — TL menyediakan '
     'feature representation yang matang dari ImageNet sebagai starting point.',
     sz=9.5, color=COLOR_TEXT)

# ─────────── SLIDE 13: Kesimpulan & Future Work ───────────
s = new_slide(prs)
title(s, 'Kesimpulan & Future Work')

# Conclusion
con = tb(s, 0.3, 0.85, 9.4, 2.1)
con.text_frame.word_wrap = True
set_fill(con, 'E6F4EA')
para(con.text_frame, 'Kesimpulan', sz=12, bold=True,
     color=COLOR_GREEN, first=True)
para(con.text_frame, '', sz=3)
para(con.text_frame,
     '1. Dataset baru FER untuk konteks pembelajaran pemrograman '
     '(7.091 sampel, 37 user, anotasi otomatis) berhasil dibangun.',
     sz=10, color=COLOR_TEXT)
para(con.text_frame,
     '2. Dari 48 konfigurasi, Intermediate Fusion + Transfer Learning '
     '(ResNet18) 4-class B1 optimal: Macro F1 = 0.412.',
     sz=10, color=COLOR_TEXT)
para(con.text_frame,
     '3. Multimodal fusion (image + landmark) mengungguli single-modal; '
     'joint feature learning lebih efektif dari late averaging.',
     sz=10, color=COLOR_TEXT)
para(con.text_frame,
     '4. Transfer learning memberi gain ~5% F1 pada dataset natural '
     '— penting untuk data kecil.',
     sz=10, color=COLOR_TEXT)

# Future work
fw = tb(s, 0.3, 3.05, 9.4, 2.15)
fw.text_frame.word_wrap = True
set_fill(fw, 'FFF0CC')
para(fw.text_frame, 'Future Work', sz=12, bold=True,
     color=COLOR_TEXT, first=True)
para(fw.text_frame, '', sz=3)
para(fw.text_frame,
     '• Perluasan dataset: tambah user & sesi → minimal 100 subjek '
     'untuk generalisasi statistik.',
     sz=10, color=COLOR_TEXT)
para(fw.text_frame,
     '• Anotasi manual human-labeled untuk validasi Face API + '
     'koreksi label noise.',
     sz=10, color=COLOR_TEXT)
para(fw.text_frame,
     '• Eksplorasi arsitektur modern: Vision Transformer (ViT) + '
     'landmark graph neural network.',
     sz=10, color=COLOR_TEXT)
para(fw.text_frame,
     '• Deployment real-time untuk integrasi ke learning management '
     'system (LMS) — adaptive feedback berbasis emosi.',
     sz=10, color=COLOR_TEXT)

# ─────────── Save ───────────
prs.save(OUTPUT)
print(f'Saved: {OUTPUT}')
print(f'Total slides: {len(prs.slides)}')
