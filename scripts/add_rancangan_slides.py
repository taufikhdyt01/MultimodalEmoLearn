"""Add 3 'Rancangan Eksperimen' slides before Hasil Training slides."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'

COLOR_HEADER_BG   = RGBColor(0x1A, 0x73, 0xE8)
COLOR_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ROW_ALT     = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_ROW_NORMAL  = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT        = RGBColor(0x20, 0x20, 0x20)
COLOR_BLUE_LIGHT  = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_ORANGE_LIGHT = RGBColor(0xFF, 0xF0, 0xCC)
SLIDE_W = Inches(10)


def insert_slide_at(prs, index, layout_idx=10):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst
    new_elem = xml_slides[-1]
    xml_slides.remove(new_elem)
    xml_slides.insert(index, new_elem)
    return slide


def add_title(slide, text, font_size=22):
    tb = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(9.3), Inches(0.7))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT


def add_text(slide, lines, left, top, width, height, font_size=10):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(font_size)
        para.font.color.rgb = COLOR_TEXT


def add_colored_box(slide, lines, left, top, width, height, bg_color, font_size=10):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    # Set background color via XML
    spPr = tb._element.spPr
    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    r, g, b = bg_color[0], bg_color[1], bg_color[2]
    srgbClr.set('val', '%02X%02X%02X' % (r, g, b))
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(font_size)
        para.font.color.rgb = COLOR_TEXT
        if i == 0:
            para.font.bold = True


def style_cell(cell, text, bold=False, bg=None, fg=None, size=9, align=PP_ALIGN.LEFT):
    cell.text = text
    para = cell.text_frame.paragraphs[0]
    para.alignment = align
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if fg:
        run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def add_table(slide, headers, rows, left, top, width, height,
              col_widths=None, h_size=10, r_size=9):
    tbl = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, width, height).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / total)
    for c, h in enumerate(headers):
        style_cell(tbl.cell(0, c), h, bold=True,
                   bg=COLOR_HEADER_BG, fg=COLOR_HEADER_TEXT,
                   size=h_size, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        bg = COLOR_ROW_ALT if r % 2 == 0 else COLOR_ROW_NORMAL
        for c, val in enumerate(row):
            style_cell(tbl.cell(r + 1, c), str(val),
                       bg=bg, fg=COLOR_TEXT, size=r_size,
                       align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)


# ── LOAD ────────────────────────────────────────────────
prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# Insert before slide 69 (index 68 = Hasil Training 7-Class)
INSERT_IDX = 68

# ── SLIDE A: Tahap 1 — 7-Class From Scratch ─────────────
s1 = insert_slide_at(prs, INSERT_IDX)
add_title(s1, 'Rancangan Eksperimen — Tahap 1: 7-Class From Scratch')
add_table(s1,
    headers=['Model', 'Arsitektur', 'Input', 'B1 Baseline', 'B2 Class Weights', 'B3 Aug'],
    rows=[
        ['CNN',                'Custom CNN',         'Citra wajah 224x224', 'V', 'V', 'V'],
        ['FCNN',               'Fully-connected NN', 'Landmark 136 fitur',  'V', 'V', 'V'],
        ['Late Fusion',        'CNN + FCNN (late)',  'Citra + Landmark',    'V', 'V', 'V'],
        ['Intermediate Fusion','CNN + FCNN (mid)',   'Citra + Landmark',    'V', 'V', 'V'],
    ],
    left=Inches(0.4), top=Inches(1.0),
    width=Inches(9.2), height=Inches(2.1),
    col_widths=[2.2, 2.3, 2.2, 0.9, 1.7, 0.9],
    h_size=10, r_size=10,
)
add_colored_box(s1, [
    'Keterangan Skenario:',
    'B1 = Baseline: training biasa tanpa penanganan imbalance',
    'B2 = Class Weights (Cui et al., 2019): penalty loss lebih besar untuk kelas langka',
    'B3 = Class Weights + Augmentasi data: flip, rotasi +-15 derajat, brightness untuk kelas < 150 sample',
], Inches(0.4), Inches(3.25), Inches(9.2), Inches(1.3),
   bg_color=COLOR_BLUE_LIGHT, font_size=10)
add_text(s1, [
    '7 kelas: neutral, happy, sad, angry, fearful, disgusted, surprised',
    '4 model x 3 skenario = 12 eksperimen  |  GPU: NVIDIA T4  |  Metrik: Macro F1-Score',
], Inches(0.4), Inches(4.65), Inches(9.2), Inches(0.6), font_size=10)


# ── SLIDE B: Tahap 2 — 4-Class ──────────────────────────
s2 = insert_slide_at(prs, INSERT_IDX + 1)
add_title(s2, 'Rancangan Eksperimen — Tahap 2: 4-Class')
add_colored_box(s2, [
    'Motivasi dari hasil Tahap 1:',
    '-> Macro F1 terbaik hanya 0.234 -- emosi langka (angry, fearful, disgusted, surprised) tidak terdeteksi sama sekali',
    '-> Penyebab: kelas terkecil hanya 8 sample (fearful) -- terlalu sedikit untuk dipelajari model',
    '-> Solusi yang dicoba: gabungkan 4 emosi langka menjadi satu kelas "negative"',
], Inches(0.4), Inches(1.0), Inches(9.2), Inches(1.5),
   bg_color=COLOR_ORANGE_LIGHT, font_size=10)
add_table(s2,
    headers=['Kelas Baru', 'Kelas Asal yang Digabung', 'Total Sample', 'Alasan'],
    rows=[
        ['neutral',  'neutral',                                   '8,356', 'Tetap — kelas dominan'],
        ['happy',    'happy',                                       '783', 'Tetap — cukup data'],
        ['sad',      'sad',                                         '576', 'Tetap — cukup data'],
        ['negative', 'angry (63) + fearful (13) + disgusted (24) + surprised (79)', '179',
         'Digabung — masing-masing < 80 sample'],
    ],
    left=Inches(0.4), top=Inches(2.65),
    width=Inches(9.2), height=Inches(2.0),
    col_widths=[1.4, 4.0, 1.4, 2.4],
    h_size=10, r_size=10,
)
add_text(s2, [
    '4 model x 3 skenario = 12 eksperimen tambahan  |  Total kumulatif: 24 eksperimen',
], Inches(0.4), Inches(4.75), Inches(9.2), Inches(0.3), font_size=10)


# ── SLIDE C: Tahap 3 — Transfer Learning ────────────────
s3 = insert_slide_at(prs, INSERT_IDX + 2)
add_title(s3, 'Rancangan Eksperimen — Tahap 3: Transfer Learning')
add_colored_box(s3, [
    'Motivasi dari hasil Tahap 2:',
    '-> FCNN (landmark) sudah membaik: Macro F1 0.394  |  Namun CNN (citra) masih rendah: Macro F1 0.296',
    '-> Penyebab CNN lemah: dataset terlalu kecil (~9K sampel) untuk melatih CNN from scratch',
    '-> Solusi: gunakan ResNet18 pretrained ImageNet -- sudah memahami fitur visual dari 1.2 juta gambar',
], Inches(0.4), Inches(1.0), Inches(9.2), Inches(1.5),
   bg_color=COLOR_ORANGE_LIGHT, font_size=10)
add_table(s3,
    headers=['Komponen', 'Tahap 1 & 2 (From Scratch)', 'Tahap 3 (Transfer Learning)', 'Catatan'],
    rows=[
        ['CNN backbone',
         'EmotionCNN (dari nol)',
         'ResNet18 pretrained ImageNet',
         'Fine-tune seluruh layer'],
        ['FCNN',
         'Fully-connected (dari nol)',
         'Sama -- tidak berubah',
         'Landmark = numerik, tidak perlu pretrained'],
        ['Late Fusion',
         'CNN scratch + FCNN',
         'CNN TL + FCNN',
         ''],
        ['Intermediate Fusion',
         'CNN scratch + FCNN (feature-level)',
         'ResNet18 + FCNN (feature-level)',
         ''],
        ['Learning Rate',
         '0.0001',
         '0.00005',
         'Lebih kecil untuk fine-tuning'],
    ],
    left=Inches(0.4), top=Inches(2.65),
    width=Inches(9.2), height=Inches(2.35),
    col_widths=[2.0, 2.6, 2.6, 2.0],
    h_size=10, r_size=9,
)
add_text(s3, [
    '4 model x 3 skenario x 2 konfigurasi kelas (7-class + 4-class) = 24 eksperimen tambahan',
    'Total keseluruhan: 12 + 12 + 24 = 48 eksperimen',
], Inches(0.4), Inches(5.1), Inches(9.2), Inches(0.4), font_size=10)


# ── SAVE ────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f'Saved! Total slides: {len(prs.slides)}')

# Verify
prs2 = Presentation(PPTX_PATH)
print('New slides 68-75:')
for i in range(67, 75):
    slide = prs2.slides[i]
    title = next((s.text_frame.text.strip()[:62].replace(chr(10), ' ')
                  for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), '')
    print(f'  Slide {i+1}: {title}')