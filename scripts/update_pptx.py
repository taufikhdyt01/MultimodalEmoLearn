"""Update PPT Bimbingan.pptx with transfer learning and all training results."""
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation('d:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx')

LAYOUT_TITLE_AND_BODY = 2


def add_slide_after(prs, after_index, layout_idx):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst
    new_slide_elem = xml_slides[-1]
    xml_slides.remove(new_slide_elem)
    xml_slides.insert(after_index + 1, new_slide_elem)
    return slide


def set_slide_content(slide, title_text, body_lines):
    shapes = slide.shapes
    title_shape = None
    body_shape = None
    for shape in shapes:
        if shape.has_text_frame:
            if title_shape is None:
                title_shape = shape
            else:
                body_shape = shape
                break
    if title_shape:
        title_shape.text_frame.text = title_text
    if body_shape and body_lines:
        tf = body_shape.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = line


new_slides_data = [
    ("Data yang Terkumpul", [
        "Total 37 Mahasiswa",
        "",
        "Batch 1 (lama): 20 mahasiswa | Sudut depan saja | April-Mei 2025",
        "Batch 2 (baru): 17 mahasiswa | Depan + Samping  | November 2025",
        "Total: 37 mahasiswa (target 38 di proposal)",
        "",
        "Catatan Batch 2: 3 dari 20 rekaman tidak tersimpan di hardisk PC perekaman",
        "-> File rekaman gagal tersimpan saat proses recording",
    ]),
    ("Pipeline Preprocessing", [
        "1. Ekstraksi Frame (berdasarkan timestamp emosi di database)",
        "   -> Bukan setiap 5 detik secara blind -- hanya frame berlabel",
        "",
        "2. Face Detection & Cropping",
        "   -> MediaPipe (bukan dlib) -> support sudut samping",
        "   -> Crop 224x224 px",
        "",
        "3. Landmark Extraction",
        "   -> 68 titik -> 136 fitur koordinat (x,y)",
        "   -> Di-map dari 478 titik MediaPipe ke 68 standar",
        "",
        "4. Dataset Preparation",
        "   -> Matching label + split train/val/test by user",
    ]),
    ("Hasil Preprocessing", [
        "Statistik Ekstraksi Frame",
        "",
        "               Batch 1    Batch 2    Total",
        "Frame diekstrak  3,849      6,553    10,402",
        "Face terdeteksi  3,824      6,070     9,894",
        "Gagal deteksi       25        483       508",
        "Deteksi rate      99.4%      92.6%     95.1%",
        "",
        "Detail Batch 2: Front 3,275 frame | Side 3,278 frame",
        "Kegagalan: wajah menunduk, keluar frame, tertutup tangan",
    ]),
    ("Distribusi Emosi Dataset", [
        "Total: 9,894 sampel -- Sangat Imbalanced",
        "",
        "Neutral   : 8,356 (84.5%)",
        "Happy     :   783  (7.9%)",
        "Sad       :   576  (5.8%)",
        "Surprised :    79  (0.8%)",
        "Angry     :    63  (0.6%)",
        "Disgusted :    24  (0.2%)",
        "Fearful   :    13  (0.1%)",
        "",
        "Konsisten dengan literatur: dominasi neutral saat fokus coding",
        "(Coto et al., 2022)",
    ]),
    ("Penanganan Class Imbalance -- 3 Skenario", [
        "B1 -- Baseline (tanpa penanganan)",
        "     Training biasa | dataset: data/dataset/ (7,064 train)",
        "",
        "B2 -- Class Weights (Cui et al., 2019 - CVPR)",
        "     Weighted cross-entropy loss",
        "     Fearful: 125x | Disgusted: 53x | Angry: 21x | Surprised: 15x",
        "",
        "B3 -- Class Weights + Augmentasi Data",
        "     Augmentasi kelas < 150 sample -> target min. 150 per kelas",
        "     Flip, rotasi +/-15 derajat, brightness adjustment",
        "     Train: 7,519 (+455 aug) | Val/Test: original (evaluasi fair)",
        "",
        "Metrik evaluasi: Macro F1-Score (bukan accuracy)",
    ]),
    ("Split Dataset", [
        "Strategi: Split by User (mencegah data leaking antar split)",
        "",
        "Split        Samples   Users   Persen",
        "Train          7,064      29    71.4%",
        "Validation     1,174       3    11.9%",
        "Test           1,656       5    16.7%",
        "",
        "Mengapa bukan tepat 80/10/10?",
        "-> Split per-user tidak bisa tepat 80/10/10",
        "-> 71/12/17 masih dalam rentang standar user-based split",
        "",
        "Semua 7 emosi terwakili di train, validation, dan test",
        "(smart distribution untuk kelas langka)",
    ]),
    ("Validasi Ahli Psikologi", [
        "3 Opsi Set Validasi (sudah disiapkan):",
        "",
        "Opsi A -- 1,938 sample | Semua non-neutral + 400 neutral | ~8 jam",
        "Opsi B -- 1,067 sample | Stratified 10% per kelas         | ~4 jam",
        "Opsi C --   583 sample | Stratified 5% per kelas          | ~2 jam",
        "",
        "Web Tool Validasi (Streamlit) sudah siap deploy:",
        "-> Tampilkan gambar wajah + label otomatis + confidence score",
        "-> Ahli klik Setuju / pilih emosi yang benar",
        "-> Multi-validator, auto-save, progress tracker",
        "-> Otomatis hitung Cohen's Kappa & inter-rater agreement",
        "-> Download hasil ke CSV",
        "",
        "(KONSULTASI 2) Opsi mana yang tepat? 1 atau 2 ahli?",
    ]),
    ("Perubahan dari Proposal", [
        "Aspek            Di Proposal         Implementasi",
        "Face detection   dlib                MediaPipe",
        "Landmark         68 titik (dlib)     68 titik (mapped MediaPipe)",
        "Jml mahasiswa    38                  37",
        "Fusion strategy  Late + Intermediate Tetap sesuai proposal",
        "",
        "Alasan perubahan dlib -> MediaPipe:",
        "-> dlib gagal deteksi wajah sudut samping (side view batch 2)",
        "-> MediaPipe support multi-angle detection",
        "-> Output tetap 136 fitur -- substansi tidak berubah",
        "",
        "(KONSULTASI 3) Perlu direvisi di proposal, atau cukup BAB 4?",
    ]),
    ("Rencana Selanjutnya", [
        "No  Tahap                              Status",
        "1   Pengumpulan data                   SELESAI",
        "2   Preprocessing (frame, crop, lm)    SELESAI",
        "3   Prepare dataset (numpy arrays)     SELESAI",
        "4   Class weights + augmentasi         SELESAI",
        "5   Training 7-class (from scratch)    SELESAI",
        "6   Training 4-class (from scratch)    SELESAI",
        "7   Transfer Learning (ResNet18)       SELESAI",
        "8   Tool validasi ahli (web app)       SELESAI - siap deploy",
        "9   Deploy tool + kirim ke ahli        MENUNGGU KEPUTUSAN",
        "10  Penulisan BAB 4 & 5               SELANJUTNYA",
    ]),
    ("Hasil Training 7-Class -- From Scratch", [
        "4 Model x 3 Skenario = 12 Eksperimen | GPU: NVIDIA T4",
        "",
        "Rank  Model              Skenario     Accuracy  Macro F1",
        "  1   FCNN               B1 Baseline    95.8%    0.234",
        "  2   Late Fusion        B1 Baseline    95.8%    0.230",
        "  3   FCNN               B2 Weights     89.1%    0.189",
        "  4   Late Fusion        B2 Weights     89.7%    0.189",
        "  7   Intermediate       B2 Weights     84.5%    0.140",
        "  9   CNN                B2 Weights     82.8%    0.134",
        " 12   Intermediate       B1 Baseline    63.3%    0.111",
        "",
        "TERBAIK: FCNN + B1 Baseline (Macro F1: 0.234)",
    ]),
    ("Analisis Hasil -- Temuan Utama (7-Class)", [
        "Temuan 1: FCNN (Landmark) > CNN (Image)",
        "  FCNN: 0.234 | CNN: 0.134",
        "  Fitur geometrik lebih robust terhadap variasi cahaya & sudut wajah",
        "",
        "Temuan 2: Fusion tidak lebih baik dari FCNN saja",
        "  Late Fusion optimal: 90% FCNN + 10% CNN",
        "  Intermediate Fusion lebih buruk -- CNN noisy ganggu landmark",
        "",
        "Temuan 3: Class weights & augmentasi TIDAK membantu di 7-class",
        "  B1 Baseline terbaik",
        "  Weight 125x pada fearful (8 sample) terlalu ekstrem",
        "",
        "Temuan 4: Kelas minoritas tidak terdeteksi di test",
        "  Happy: 0% recall | Angry: 0% | Fearful: 0%",
        "  Hanya Neutral (99%) dan Sad (32%) yang terdeteksi",
    ]),
    ("Hasil Training 4-Class -- From Scratch", [
        "4 kelas: neutral, happy, sad, negative (angry+fearful+disgusted+surprised)",
        "",
        "Rank  Model              Skenario   Accuracy  Macro F1",
        "  1   FCNN               B3 Aug       94.4%    0.394  (+68%)",
        "  2   Late Fusion        B3 Aug       94.6%    0.385",
        "  3   FCNN               B1 Baseline  95.8%    0.330",
        "  7   CNN                B2 Weights   92.8%    0.296",
        "  9   Intermediate       B2 Weights   87.3%    0.258",
        "",
        "Temuan 5: Penggabungan kelas langka -> Macro F1 +68%",
        "Temuan 6: B3 (augmentasi) terbaik di 4-class vs B1 di 7-class",
        "  -> negative punya 145 sample -> augmentasi bermakna",
        "Temuan 7: FCNN tetap konsisten terbaik di kedua konfigurasi",
    ]),
    ("Perbandingan 7-Class vs 4-Class", [
        "Model          7-Class   4-Class   Peningkatan",
        "FCNN            0.234     0.394       +68%",
        "Late Fusion     0.230     0.385       +67%",
        "CNN             0.134     0.296      +121%",
        "Intermediate    0.140     0.258       +84%",
        "",
        "Mengapa 4-class jauh lebih baik?",
        "-> Kelas 'negative' punya 145 sampel vs maks 70 di 7-class",
        "-> Masalah bukan arsitektur -- tapi jumlah data minoritas",
        "",
        "Mengapa augmentasi berhasil di 4-class tapi gagal di 7-class?",
        "-> 4-class: 145 sample -> augmentasi bermakna",
        "-> 7-class: fearful 8 sample -> augmentasi tidak bermakna",
    ]),
    ("Hasil Transfer Learning -- ResNet18 Pretrained ImageNet", [
        "CNN diganti: EmotionCNN -> EmotionCNNTransfer (ResNet18) | LR: 0.00005",
        "FCNN tidak berubah (landmark = numerik, tidak butuh pretrained)",
        "",
        "7-Class:              From Scratch  Transfer Learning  Delta",
        "CNN                      0.134           0.177         +32%",
        "FCNN                     0.234           0.234          --",
        "Late Fusion              0.230           0.234          +2%",
        "Intermediate Fusion      0.140           0.232         +66%",
        "",
        "4-Class:              From Scratch  Transfer Learning  Delta",
        "CNN                      0.296           0.407         +37%",
        "FCNN                     0.394           0.394          --",
        "Late Fusion              0.385           0.442         +15%  BEST",
        "Intermediate Fusion      0.258           0.376         +46%",
        "",
        "BEST OVERALL: Late Fusion TL 4-class B2 -- Macro F1: 0.442",
    ]),
    ("Perbandingan Lengkap -- 48 Kombinasi Eksperimen", [
        "Ringkasan 3 Tahap Eksperimen:",
        "",
        "Tahap 1: 7-class from scratch",
        "  Best: FCNN B1 | Macro F1: 0.234",
        "",
        "Tahap 2: 4-class from scratch",
        "  Best: FCNN B3 | Macro F1: 0.394  (+68% dari tahap 1)",
        "",
        "Tahap 3: Transfer Learning (ResNet18)",
        "  Best: Late Fusion TL B2 4-class | Macro F1: 0.442  (+12%)",
        "",
        "Top 5 dari 48 kombinasi (4 model x 2 CNN variant x 3 skenario x 2 kelas):",
        "  1. Late Fusion TL   | 4-class | B2 | Macro F1: 0.442",
        "  2. CNN TL           | 4-class | B2 | Macro F1: 0.407",
        "  3. FCNN             | 4-class | B3 | Macro F1: 0.394",
        "  4. FCNN TL          | 4-class | B3 | Macro F1: 0.394",
        "  5. Late Fusion      | 4-class | B3 | Macro F1: 0.385",
    ]),
    ("Diskusi & Konsultasi", [
        "Jawaban Rumusan Masalah:",
        "RQ1 CNN: From scratch 0.134 -> Transfer Learning 0.407 (4-class, +37%)",
        "RQ2 FCNN: 0.234 (7-class) | 0.394 (4-class) -- landmark robust",
        "RQ3 Fusion: From scratch: FCNN > Late Fusion > Intermediate",
        "            Transfer Learning: Late Fusion TL terbaik (0.442)",
        "",
        "(KONSULTASI 4) Pertanyaan untuk Pembimbing:",
        "",
        "1. Transfer Learning sebagai kontribusi tesis?",
        "   TL meningkatkan F1: 0.394 -> 0.442. Cukup signifikan?",
        "",
        "2. FCNN > Fusion (from scratch) vs hipotesis proposal?",
        "   Dengan TL, fusion akhirnya menjadi terbaik.",
        "   Bagaimana narasi di tesis?",
        "",
        "3. Macro F1 0.442 -- acceptable untuk tesis S2?",
        "   Dataset nyata, naturalistik, sangat imbalanced.",
    ]),
]

insert_after = 58

for i, (title, body) in enumerate(new_slides_data):
    slide = add_slide_after(prs, insert_after + i, LAYOUT_TITLE_AND_BODY)
    set_slide_content(slide, title, body)
    print(f"Added slide {insert_after + i + 2}: {title}")

prs.save('d:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx')
print(f"\nDone! Total slides: {len(prs.slides)}")
