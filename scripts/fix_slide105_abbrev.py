"""Fix abbreviations (tr/te/va) in slide 105 table to full words."""
from pptx import Presentation

PPTX = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'
prs = Presentation(PPTX)

# Slide 105 (0-based index 104)
s = prs.slides[104]

replacements = {
    # Table cells (already fixed once, now update to clearer full format)
    '14,449 imgs (11,565 train / 2,884 test)':
        '15,339 imgs → 14,449 usable (11,565 train / 2,884 test)',
    '7,091 imgs (5,287 train / 579 val / 929 test)':
        'conf60: 6,795 imgs (5,287 train / 579 val / 929 test)',
    # Note
    'Catatan: RAF-DB pakai basic-emotion public release (15,339 imgs official) — BUKAN subset. '
    'KDEF ~32% di-drop karena MediaPipe gagal deteksi wajah di sudut profil penuh.':
        'Catatan: RAF-DB ~5.8% di-drop (15,339 → 14,449 usable). '
        'KDEF ~32% di-drop karena MediaPipe gagal deteksi wajah di sudut profil. '
        'Primer pakai versi conf60 (confidence ≥60%) yang sama dengan eksperimen utama.',
}

fixed = 0


def apply_to_runs(runs, label):
    global fixed
    for run in runs:
        for old, new in replacements.items():
            if old in run.text:
                run.text = run.text.replace(old, new)
                fixed += 1
                print(f'Fixed {label} #{fixed}: length {len(old)} -> {len(new)}')


for shape in s.shapes:
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    apply_to_runs(para.runs, 'table')
    elif shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            apply_to_runs(para.runs, 'textbox')

prs.save(PPTX)
print(f'Total cells fixed: {fixed}')
