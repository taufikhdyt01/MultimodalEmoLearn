"""
Reorder slides so narrative flows:
  Tahap1 → Hasil 7-class → Analisis → Tahap2 → Hasil 4-class → Perbandingan → Tahap3 → TL → ...
Also rebuild Analisis Hasil as bullet points (not table).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

PPTX_PATH = 'd:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx'
COLOR_TEXT       = RGBColor(0x20, 0x20, 0x20)
COLOR_HEADER_BG  = RGBColor(0x1A, 0x73, 0xE8)
COLOR_HEADER_TEXT= RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_GREEN_LIGHT= RGBColor(0xE6, 0xF4, 0xEA)

prs = Presentation(PPTX_PATH)
print(f'Loaded: {len(prs.slides)} slides')

# ── STEP 1: Reorder slides ───────────────────────────────
# Current layout (indices):
#   68 = Rancangan Tahap 1
#   69 = Rancangan Tahap 2   ← move to after idx 74 (Analisis)
#   70 = Rancangan Tahap 3   ← move to after idx 78 (Perbandingan 7vs4 table)
#   71 = Hasil Training 7-Class
#   72 = Visualisasi bar 7-class
#   73 = Heatmap 7-class
#   74 = Analisis Hasil
#   75 = Hasil Training 4-Class
#   76 = Visualisasi 4-class
#   77 = Perbandingan 7vs4 image
#   78 = Perbandingan 7vs4 table
#   79 = Hasil Transfer Learning
#   80 = Visualisasi final bar
#   81 = Heatmap final
#   82 = Perbandingan Lengkap
#   83 = Diskusi
#   84 = Terimakasih
#
# Target order for indices 68-84:
#   68, 71, 72, 73, 74, 69, 75, 76, 77, 78, 70, 79, 80, 81, 82, 83, 84

xml_slides = prs.slides._sldIdLst
all_elems = list(xml_slides)

# Extract the sldId elements we want to reorder (indices 68-84)
segment = all_elems[68:85]  # slides at positions 68..84

# New order within this segment (0-based within segment):
# segment[0]=Tahap1, segment[1]=Tahap2, segment[2]=Tahap3,
# segment[3]=7class, segment[4]=bar, segment[5]=heatmap, segment[6]=analisis,
# segment[7]=4class, segment[8]=4classbar, segment[9]=4classvs7img,
# segment[10]=4classvs7tbl, segment[11]=TL, segment[12]=finalbar,
# segment[13]=finalheatmap, segment[14]=perbandingan, segment[15]=diskusi, segment[16]=terimakasih
new_order = [0, 3, 4, 5, 6, 1, 7, 8, 9, 10, 2, 11, 12, 13, 14, 15, 16]
reordered = [segment[i] for i in new_order]

# Remove old segment and replace with reordered
for elem in segment:
    xml_slides.remove(elem)
for i, elem in enumerate(reordered):
    xml_slides.insert(68 + i, elem)

print('Slides reordered.')

# Verify new order
for i in range(67, len(prs.slides)):
    slide = prs.slides[i]
    title = next((s.text_frame.text.strip()[:55].replace(chr(10),' ')
                  for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), '')
    print(f'  Slide {i+1}: {title}')


# ── STEP 2: Rebuild Analisis Hasil as bullet points ──────
# After reorder, Analisis Hasil is now at index 72 (slide 73)
analisis_idx = None
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t.startswith('Analisis Hasil'):
                analisis_idx = i
                break
    if analisis_idx is not None:
        break

print(f'\nAnalisis Hasil at index {analisis_idx} (slide {analisis_idx+1})')
slide = prs.slides[analisis_idx]

# Clear all shapes
sp_tree = slide.shapes._spTree
for child in list(sp_tree):
    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
    if tag not in ('bg', 'bgPr'):
        sp_tree.remove(child)

# Add title
tb = slide.shapes.add_textbox(Inches(0.33), Inches(0.2), Inches(9.3), Inches(0.7))
tf = tb.text_frame
tf.text = 'Analisis Hasil — Temuan Utama (7-Class)'
p = tf.paragraphs[0]
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLOR_TEXT


def style_cell(cell, text, bold=False, bg=None, fg=None, size=9, align=PP_ALIGN.LEFT):
    cell.text = text
    para = cell.text_frame.paragraphs[0]
    para.alignment = align
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if fg:
        run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def add_label(slide, text, left, top, width, font_size=10):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT


def add_mini_table(slide, headers, rows, left, top, width, height,
                   col_widths=None, h_size=9, r_size=9, highlight_rows=None):
    tbl = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, width, height).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / total)
    for c, h in enumerate(headers):
        style_cell(tbl.cell(0, c), h, bold=True,
                   bg=COLOR_HEADER_BG, fg=COLOR_HEADER_TEXT,
                   size=h_size, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        is_hi = highlight_rows and r in highlight_rows
        bg = RGBColor(0x34, 0xA8, 0x53) if is_hi else \
             (COLOR_BLUE_LIGHT if r % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF))
        fg = RGBColor(0xFF, 0xFF, 0xFF) if is_hi else COLOR_TEXT
        for c, val in enumerate(row):
            style_cell(tbl.cell(r + 1, c), str(val),
                       bold=is_hi, bg=bg, fg=fg, size=r_size,
                       align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)


# --- Temuan 1: FCNN > CNN ---
add_label(slide, 'Temuan 1: FCNN (Landmark) > CNN (Citra)', Inches(0.4), Inches(1.0), Inches(4.5))
add_mini_table(slide,
    headers=['Model', 'Best Macro F1'],
    rows=[['FCNN', '0.234'], ['CNN', '0.134']],
    left=Inches(0.4), top=Inches(1.35), width=Inches(4.5), height=Inches(0.75),
    col_widths=[3, 1.5], highlight_rows=[0],
)

# --- Temuan 2: Fusion < FCNN ---
add_label(slide, 'Temuan 2: Fusion tidak lebih baik dari FCNN saja', Inches(5.1), Inches(1.0), Inches(4.5))
add_mini_table(slide,
    headers=['Model', 'Best Macro F1'],
    rows=[
        ['FCNN', '0.234'],
        ['Late Fusion (90% FCNN + 10% CNN)', '0.230'],
        ['Intermediate Fusion', '0.140'],
    ],
    left=Inches(5.1), top=Inches(1.35), width=Inches(4.5), height=Inches(1.0),
    col_widths=[3.5, 1], highlight_rows=[0],
)

# --- Temuan 3: Skenario ---
add_label(slide, 'Temuan 3: Class Weights & Augmentasi tidak membantu', Inches(0.4), Inches(2.5), Inches(4.5))
add_mini_table(slide,
    headers=['Skenario', 'FCNN F1', 'CNN F1'],
    rows=[
        ['B1 Baseline', '0.234', '0.133'],
        ['B2 Class Weights', '0.189', '0.134'],
        ['B3 Weights + Aug', '0.182', '0.119'],
    ],
    left=Inches(0.4), top=Inches(2.85), width=Inches(4.5), height=Inches(1.0),
    col_widths=[2.5, 1, 1], highlight_rows=[0],
)

# --- Temuan 4: Per-class recall ---
add_label(slide, 'Temuan 4: Kelas minoritas tidak terdeteksi (FCNN B1)', Inches(5.1), Inches(2.5), Inches(4.5))
add_mini_table(slide,
    headers=['Emosi', 'Support', 'Recall', 'F1'],
    rows=[
        ['Neutral',   '1,588', '99%', '0.98'],
        ['Sad',         '38', '32%', '0.41'],
        ['Happy',       '10',  '0%', '0.00'],
        ['Angry',       '13',  '0%', '0.00'],
        ['Fearful',      '1',  '0%', '0.00'],
    ],
    left=Inches(5.1), top=Inches(2.85), width=Inches(4.5), height=Inches(1.3),
    col_widths=[1.6, 1, 0.95, 0.95], highlight_rows=[0, 1],
)

# Kesimpulan
tb2 = slide.shapes.add_textbox(Inches(0.4), Inches(4.3), Inches(9.2), Inches(0.4))
tf2 = tb2.text_frame
tf2.text = ('Kesimpulan: Masalah bukan pada arsitektur — '
            'kelas minoritas terlalu sedikit untuk dipelajari model')
p2 = tf2.paragraphs[0]
p2.font.size = Pt(10)
p2.font.bold = True
p2.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)

# ── SAVE ────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f'\nSaved! Total slides: {len(prs.slides)}')